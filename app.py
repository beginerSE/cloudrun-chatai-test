# -*- coding: utf-8 -*-
import os
import re
import json
import ast
import asyncio
import tempfile
from datetime import datetime
from typing import Any, Dict, List, Optional
from contextlib import AsyncExitStack
import queue
import threading

import pandas as pd
import psycopg2
from psycopg2.extras import RealDictCursor
import bcrypt
from flask import Flask, request, jsonify, render_template, send_from_directory, Response, redirect, url_for, flash, session
from flask_login import LoginManager, UserMixin, login_user, logout_user, login_required, current_user
from urllib.parse import urlparse, urljoin
from dotenv import load_dotenv
from openai import OpenAI
import google.generativeai as genai
from google import genai as genai_new
from google.genai import types as genai_types

from mcp import ClientSession, StdioServerParameters
from mcp.client.stdio import stdio_client

from google.cloud import bigquery
from google.oauth2 import service_account as gcp_service_account
from google.api_core import exceptions as gcp_exceptions

load_dotenv()

# Custom JSON encoder for datetime objects
class DateTimeEncoder(json.JSONEncoder):
    def default(self, obj):
        if isinstance(obj, datetime):
            return obj.isoformat()
        return super().default(obj)

PROJECT_ID = os.getenv("PROJECT_ID", "your-project-id")
LOCATION = os.getenv("LOCATION", "US")
DEFAULT_DATASET = os.getenv("DEFAULT_DATASET", "your_dataset")
OPENAI_API_KEY = os.getenv("OPENAI_API_KEY", "")
OPENAI_MODEL = os.getenv("OPENAI_MODEL", "gpt-4o-mini")
GCP_SA_JSON = os.getenv("GCP_SA_JSON", "")
DEFAULT_MAX_SCAN_GB = float(os.getenv("DEFAULT_MAX_SCAN_GB", "10"))
DATABASE_URL = os.getenv("DATABASE_URL", "")

app = Flask(__name__)
app.config['SECRET_KEY'] = os.getenv('SESSION_SECRET', 'dev-secret-key-change-in-production')

@app.after_request
def add_header(response):
    """Disable caching for HTML pages during development"""
    if 'text/html' in response.content_type:
        response.headers['Cache-Control'] = 'no-cache, no-store, must-revalidate'
        response.headers['Pragma'] = 'no-cache'
        response.headers['Expires'] = '0'
    return response

# Configure Flask to use DateTimeEncoder for all JSON responses
from flask.json.provider import DefaultJSONProvider
class DateTimeJSONProvider(DefaultJSONProvider):
    def default(self, obj):
        if isinstance(obj, datetime):
            # Add 'Z' suffix to indicate UTC time
            return obj.isoformat() + 'Z'
        return super().default(obj)

app.json = DateTimeJSONProvider(app)

# Flask-Login setup
login_manager = LoginManager()
login_manager.init_app(app)
login_manager.login_view = 'login'
login_manager.login_message = 'このページにアクセスするにはログインが必要です。'

@login_manager.unauthorized_handler
def unauthorized_callback():
    """Custom unauthorized handler - returns JSON for API requests, redirect for pages"""
    if request.path.startswith('/api/'):
        # API request - return JSON error instead of redirect
        return jsonify({
            "error": "認証が必要です。再度ログインしてください。",
            "auth_required": True,
            "redirect": url_for('login')
        }), 401
    # Regular page request - redirect to login
    return redirect(url_for('login', next=request.url))

# Progress tracking for chat tasks - now using PostgreSQL for Cloud Run compatibility
# Legacy in-memory dict removed - all task state is stored in database

def create_chat_task(task_id: str, user_id: int, project_id: int) -> bool:
    """Create a new chat task in the database"""
    try:
        conn = get_db_connection()
        cur = conn.cursor()
        cur.execute('''
            INSERT INTO chat_tasks (task_id, user_id, project_id, status, steps, reasoning)
            VALUES (%s, %s, %s, 'running', '[]'::jsonb, '')
        ''', (task_id, user_id, project_id))
        conn.commit()
        cur.close()
        conn.close()
        return True
    except Exception as e:
        print(f"Error creating chat task: {e}")
        return False

def get_chat_task(task_id: str) -> dict:
    """Get chat task from database"""
    try:
        conn = get_db_connection()
        cur = conn.cursor(cursor_factory=RealDictCursor)
        cur.execute('''
            SELECT task_id, user_id, project_id, status, steps, result, error, 
                   reasoning, cancelled, session_id, traceback, created_at, updated_at
            FROM chat_tasks WHERE task_id = %s
        ''', (task_id,))
        task = cur.fetchone()
        cur.close()
        conn.close()
        return dict(task) if task else None
    except Exception as e:
        print(f"Error getting chat task: {e}")
        return None

def update_chat_task(task_id: str, **kwargs) -> bool:
    """Update chat task fields in database"""
    if not kwargs:
        return True
    try:
        conn = get_db_connection()
        cur = conn.cursor()
        
        # Build dynamic UPDATE query
        set_clauses = []
        params = []
        for key, value in kwargs.items():
            if key in ['steps', 'result']:
                set_clauses.append(f"{key} = %s::jsonb")
                params.append(json.dumps(value) if not isinstance(value, str) else value)
            else:
                set_clauses.append(f"{key} = %s")
                params.append(value)
        
        set_clauses.append("updated_at = CURRENT_TIMESTAMP")
        params.append(task_id)
        
        query = f"UPDATE chat_tasks SET {', '.join(set_clauses)} WHERE task_id = %s"
        cur.execute(query, params)
        conn.commit()
        cur.close()
        conn.close()
        return True
    except Exception as e:
        print(f"Error updating chat task: {e}")
        import traceback
        traceback.print_exc()
        return False

def add_chat_task_step(task_id: str, step: str) -> bool:
    """Add a step to chat task"""
    try:
        conn = get_db_connection()
        cur = conn.cursor()
        cur.execute('''
            UPDATE chat_tasks 
            SET steps = steps || %s::jsonb, updated_at = CURRENT_TIMESTAMP
            WHERE task_id = %s
        ''', (json.dumps([step]), task_id))
        conn.commit()
        cur.close()
        conn.close()
        return True
    except Exception as e:
        print(f"Error adding chat task step: {e}")
        return False

def append_chat_task_reasoning(task_id: str, text: str) -> bool:
    """Append text to chat task reasoning"""
    try:
        conn = get_db_connection()
        cur = conn.cursor()
        cur.execute('''
            UPDATE chat_tasks 
            SET reasoning = COALESCE(reasoning, '') || %s, updated_at = CURRENT_TIMESTAMP
            WHERE task_id = %s
        ''', (text, task_id))
        conn.commit()
        cur.close()
        conn.close()
        return True
    except Exception as e:
        print(f"Error appending chat task reasoning: {e}")
        return False

def is_chat_task_cancelled(task_id: str) -> bool:
    """Check if chat task is cancelled"""
    try:
        conn = get_db_connection()
        cur = conn.cursor()
        cur.execute('SELECT cancelled FROM chat_tasks WHERE task_id = %s', (task_id,))
        result = cur.fetchone()
        cur.close()
        conn.close()
        return result[0] if result else False
    except Exception as e:
        print(f"Error checking chat task cancelled: {e}")
        return False

def cleanup_old_chat_tasks(hours: int = 24) -> int:
    """Clean up old chat tasks (older than specified hours)"""
    try:
        conn = get_db_connection()
        cur = conn.cursor()
        cur.execute('''
            DELETE FROM chat_tasks 
            WHERE created_at < NOW() - INTERVAL '%s hours'
        ''', (hours,))
        deleted = cur.rowcount
        conn.commit()
        cur.close()
        conn.close()
        return deleted
    except Exception as e:
        print(f"Error cleaning up old chat tasks: {e}")
        return 0

# Database connection
def get_db_connection():
    """
    Cloud Run compatible database connection.
    Supports both Cloud SQL Unix socket and standard PostgreSQL URL.
    """
    # Check if running in Cloud Run with Cloud SQL
    cloud_sql_instance = os.getenv('CLOUD_SQL_CONNECTION_NAME')
    
    if cloud_sql_instance:
        # Cloud Run with Cloud SQL via Unix socket
        db_user = os.getenv('DB_USER', 'postgres')
        db_pass = os.getenv('DB_PASSWORD', '')
        db_name = os.getenv('DB_NAME', 'agent_db')
        
        # Unix socket connection for Cloud Run
        unix_socket = f'/cloudsql/{cloud_sql_instance}'
        connection_string = f'postgresql://{db_user}:{db_pass}@/{db_name}?host={unix_socket}'
        conn = psycopg2.connect(connection_string)
    else:
        # Local development or standard PostgreSQL URL
        conn = psycopg2.connect(DATABASE_URL)
    
    return conn

def run_startup_migrations():
    """Apply lightweight, idempotent schema migrations at startup"""
    try:
        conn = get_db_connection()
        cur = conn.cursor()
        cur.execute("ALTER TABLE projects ADD COLUMN IF NOT EXISTS max_scan_gb NUMERIC DEFAULT 10")
        cur.execute("ALTER TABLE project_memories ADD COLUMN IF NOT EXISTS category VARCHAR(50) DEFAULT 'general'")
        cur.execute("ALTER TABLE project_memories ADD COLUMN IF NOT EXISTS table_name VARCHAR(512)")
        cur.execute("ALTER TABLE projects ADD COLUMN IF NOT EXISTS enable_review BOOLEAN DEFAULT false")
        cur.execute("ALTER TABLE projects ADD COLUMN IF NOT EXISTS gemini_model VARCHAR(100) DEFAULT 'gemini-2.5-pro'")
        cur.execute('''
            CREATE TABLE IF NOT EXISTS project_documents (
                id SERIAL PRIMARY KEY,
                project_id INTEGER NOT NULL REFERENCES projects(id) ON DELETE CASCADE,
                user_id INTEGER NOT NULL REFERENCES users(id) ON DELETE CASCADE,
                filename VARCHAR(512) NOT NULL,
                file_type VARCHAR(20) NOT NULL,
                extracted_text TEXT,
                created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
            )
        ''')
        cur.execute("CREATE INDEX IF NOT EXISTS idx_project_documents_project_id ON project_documents(project_id)")
        cur.execute("ALTER TABLE project_memories ADD COLUMN IF NOT EXISTS document_id INTEGER REFERENCES project_documents(id) ON DELETE CASCADE")
        cur.execute('''
            CREATE TABLE IF NOT EXISTS table_baseline_stats (
                id SERIAL PRIMARY KEY,
                project_id INTEGER NOT NULL REFERENCES projects(id) ON DELETE CASCADE,
                user_id INTEGER NOT NULL REFERENCES users(id) ON DELETE CASCADE,
                dataset_id VARCHAR(512) NOT NULL,
                table_name VARCHAR(512) NOT NULL,
                row_count BIGINT,
                numeric_stats JSONB DEFAULT '[]'::jsonb,
                date_ranges JSONB DEFAULT '[]'::jsonb,
                error TEXT,
                computed_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP,
                UNIQUE (project_id, dataset_id, table_name)
            )
        ''')
        cur.execute("CREATE INDEX IF NOT EXISTS idx_table_baseline_stats_project ON table_baseline_stats(project_id)")
        cur.execute('''
            CREATE TABLE IF NOT EXISTS analysis_metrics (
                id SERIAL PRIMARY KEY,
                project_id INTEGER NOT NULL REFERENCES projects(id) ON DELETE CASCADE,
                user_id INTEGER NOT NULL REFERENCES users(id) ON DELETE CASCADE,
                session_id INTEGER,
                question TEXT,
                metric_name VARCHAR(255) NOT NULL,
                metric_value DOUBLE PRECISION NOT NULL,
                metric_unit VARCHAR(50),
                created_at TIMESTAMP DEFAULT CURRENT_TIMESTAMP
            )
        ''')
        cur.execute("CREATE INDEX IF NOT EXISTS idx_analysis_metrics_project ON analysis_metrics(project_id, created_at DESC)")
        conn.commit()
        cur.close()
        conn.close()
        print("DEBUG: Startup migrations applied (projects.max_scan_gb, projects.enable_review, project_memories.category/table_name)")
    except Exception as e:
        print(f"WARNING: Startup migration failed: {e}")

run_startup_migrations()

# Safe redirect helper
def is_safe_url(target):
    """Check if the target URL is safe to redirect to"""
    ref_url = urlparse(request.host_url)
    test_url = urlparse(urljoin(request.host_url, target))
    return test_url.scheme in ('http', 'https') and ref_url.netloc == test_url.netloc

def get_redirect_target():
    """Get a safe redirect target from request args or referrer"""
    for target in request.args.get('next'), request.referrer:
        if not target:
            continue
        if is_safe_url(target):
            return target
    return None

# User class
class User(UserMixin):
    def __init__(self, id, username, email):
        self.id = id
        self.username = username
        self.email = email

@login_manager.user_loader
def load_user(user_id):
    conn = get_db_connection()
    cur = conn.cursor(cursor_factory=RealDictCursor)
    cur.execute('SELECT id, username, email FROM users WHERE id = %s', (user_id,))
    user_data = cur.fetchone()
    cur.close()
    conn.close()
    
    if user_data:
        return User(user_data['id'], user_data['username'], user_data['email'])
    return None

def extract_project_id_from_sa_json(service_account_json_path: str) -> str:
    """Extract project_id from service account JSON file"""
    if not service_account_json_path:
        return None
    try:
        # Handle both file path and JSON string
        if os.path.exists(service_account_json_path):
            with open(service_account_json_path, 'r') as f:
                sa_data = json.load(f)
        else:
            # Try parsing as JSON string
            sa_data = json.loads(service_account_json_path)
        return sa_data.get('project_id')
    except Exception as e:
        print(f"Error extracting project_id from service account JSON: {e}")
        return None

def get_adc_project_id() -> str:
    """Get project ID from Application Default Credentials (ADC)"""
    try:
        import google.auth
        credentials, project = google.auth.default()
        if project:
            print(f"DEBUG: ADC project_id detected: {project}")
            return project
    except Exception as e:
        print(f"Error getting ADC project_id: {e}")
    return None

def get_active_project_config(user_id):
    """Get configuration from active project or fall back to env variables"""
    conn = get_db_connection()
    cur = conn.cursor(cursor_factory=RealDictCursor)
    cur.execute('''
        SELECT openai_api_key, gemini_api_key, ai_provider, bigquery_project_id, bigquery_dataset_id, service_account_json, use_env_json, max_scan_gb, enable_review, gemini_model
        FROM projects
        WHERE user_id = %s AND is_active = true
        LIMIT 1
    ''', (user_id,))
    project = cur.fetchone()
    cur.close()
    conn.close()
    
    if project:
        provider = project['ai_provider'] or 'openai'
        use_adc = project.get('use_env_json', False)
        service_account_json = project['service_account_json'] if not use_adc else None
        
        # Get project_id: prefer explicit setting, then extract from service account JSON, then ADC
        project_id = project['bigquery_project_id']
        if not project_id and service_account_json and not use_adc:
            project_id = extract_project_id_from_sa_json(service_account_json)
        if not project_id and use_adc:
            # Try to get project ID from ADC (Cloud Run automatically provides this)
            project_id = get_adc_project_id()
        if not project_id:
            project_id = PROJECT_ID
        
        # dataset_id: empty string or placeholder means not set
        dataset_id = project['bigquery_dataset_id']
        if not dataset_id or dataset_id == 'your_dataset':
            dataset_id = None
        
        return {
            'api_key': project['openai_api_key'] or OPENAI_API_KEY,
            'gemini_api_key': project['gemini_api_key'],
            'provider': provider,
            'project_id': project_id,
            'dataset_id': dataset_id,
            'service_account_json': service_account_json,
            'use_adc': use_adc,
            'max_scan_gb': float(project['max_scan_gb']) if project.get('max_scan_gb') else DEFAULT_MAX_SCAN_GB,
            'enable_review': bool(project.get('enable_review')),
            'gemini_model': project.get('gemini_model') or 'gemini-2.5-pro'
        }
    else:
        # Fall back to environment variables
        env_dataset = DEFAULT_DATASET if DEFAULT_DATASET != 'your_dataset' else None
        return {
            'api_key': OPENAI_API_KEY,
            'gemini_api_key': None,
            'provider': 'openai',
            'project_id': PROJECT_ID,
            'dataset_id': env_dataset,
            'service_account_json': GCP_SA_JSON,
            'use_adc': False,
            'max_scan_gb': DEFAULT_MAX_SCAN_GB,
            'enable_review': False,
            'gemini_model': 'gemini-2.5-pro'
        }

def build_memory_section(project_memories) -> str:
    """Build the PROJECT MEMORY section of the system prompt, grouped by category.

    Categories:
    - table_info: supplementary notes about specific tables (column meanings, caveats, join keys)
    - business_rule: business rules/definitions the AI must follow when building queries
    - general: everything else
    """
    if not project_memories:
        return ""
    table_notes = [m for m in project_memories if m.get('category') == 'table_info']
    business_rules = [m for m in project_memories if m.get('category') == 'business_rule']
    doc_context = [m for m in project_memories if m.get('category') == 'document_context']
    general = [m for m in project_memories if m.get('category') not in ('table_info', 'business_rule', 'document_context')]

    section = "\n\n## PROJECT MEMORY\nThe following information has been saved across all chat sessions for this project. Use this context to provide more relevant and personalized analysis:\n\n"

    if doc_context:
        section += "### PROJECT OBJECTIVES & ANALYSIS POLICY (extracted from uploaded proposal/minutes documents)\n"
        section += "Always align your analysis, interpretations, and insights with these objectives, policies, and KPI definitions. When presenting insights, explicitly connect them to the stated objectives where relevant:\n\n"
        for mem in doc_context:
            section += f"#### {mem['memory_key']}\n{mem['memory_value']}\n\n"

    if table_notes:
        section += "### TABLE NOTES (user-provided supplementary information about tables)\n"
        section += "Always take these notes into account when selecting tables, interpreting columns, and writing SQL:\n\n"
        for mem in table_notes:
            tbl = mem.get('table_name') or mem['memory_key']
            section += f"#### Table: {tbl}\n{mem['memory_value']}\n\n"

    if business_rules:
        section += "### BUSINESS RULES (must be followed when building queries and interpreting results)\n\n"
        for mem in business_rules:
            section += f"#### {mem['memory_key']}\n{mem['memory_value']}\n\n"

    if general:
        section += "### OTHER NOTES\n\n"
        for mem in general:
            section += f"#### {mem['memory_key']}\n{mem['memory_value']}\n\n"

    return section

MAX_REVIEW_ROUNDS = 2

def build_doc_context_section(project_memories) -> str:
    """Build a compact Japanese section of document-derived objectives/policies for the reviewer AI."""
    doc_context = [m for m in (project_memories or []) if m.get('category') == 'document_context']
    if not doc_context:
        return ""
    lines = []
    for mem in doc_context:
        lines.append(f"### {mem['memory_key']}\n{mem['memory_value']}")
    return "\n\n".join(lines)[:6000]


def load_baseline_stats(project_db_id, user_id):
    """Load saved table baseline stats for a project. Returns a list (possibly empty)."""
    if not project_db_id or not user_id:
        return []
    try:
        conn = get_db_connection()
        cur = conn.cursor(cursor_factory=RealDictCursor)
        cur.execute('''
            SELECT dataset_id, table_name, row_count, numeric_stats, date_ranges, computed_at
            FROM table_baseline_stats
            WHERE project_id = %s AND user_id = %s AND error IS NULL
            ORDER BY dataset_id, table_name
            LIMIT 30
        ''', (project_db_id, user_id))
        rows = cur.fetchall()
        cur.close()
        conn.close()
        return rows
    except Exception as e:
        print(f"Warning: Failed to load baseline stats: {e}")
        return []


def load_past_metrics(project_db_id, user_id, limit=40):
    """Load recently saved analysis metrics for a project. Returns a list (possibly empty)."""
    if not project_db_id or not user_id:
        return []
    try:
        conn = get_db_connection()
        cur = conn.cursor(cursor_factory=RealDictCursor)
        cur.execute('''
            SELECT question, metric_name, metric_value, metric_unit, created_at
            FROM analysis_metrics
            WHERE project_id = %s AND user_id = %s
            ORDER BY created_at DESC
            LIMIT %s
        ''', (project_db_id, user_id, limit))
        rows = cur.fetchall()
        cur.close()
        conn.close()
        return rows
    except Exception as e:
        print(f"Warning: Failed to load past analysis metrics: {e}")
        return []


def build_baseline_section(baseline_stats) -> str:
    """Build a compact Japanese section of table baseline stats for the reviewer AI."""
    if not baseline_stats:
        return ""
    lines = []
    for row in baseline_stats:
        parts = [f"### {row['dataset_id']}.{row['table_name']}（{row.get('computed_at')}時点）"]
        if row.get('row_count') is not None:
            parts.append(f"- 行数: {row['row_count']:,}")
        numeric_stats = row.get('numeric_stats') or []
        if isinstance(numeric_stats, str):
            try:
                numeric_stats = json.loads(numeric_stats)
            except Exception:
                numeric_stats = []
        for col in numeric_stats[:8]:
            parts.append(f"- {col.get('column')}: 合計={col.get('sum')}, 最小={col.get('min')}, 最大={col.get('max')}")
        date_ranges = row.get('date_ranges') or []
        if isinstance(date_ranges, str):
            try:
                date_ranges = json.loads(date_ranges)
            except Exception:
                date_ranges = []
        for col in date_ranges[:3]:
            parts.append(f"- {col.get('column')}: 期間 {col.get('min')} 〜 {col.get('max')}")
        lines.append("\n".join(parts))
    return "\n\n".join(lines)[:6000]


def build_past_metrics_section(past_metrics) -> str:
    """Build a compact Japanese section of past analysis metrics for the reviewer AI."""
    if not past_metrics:
        return ""
    lines = []
    for m in past_metrics:
        q = (m.get('question') or '')[:80]
        unit = m.get('metric_unit') or ''
        lines.append(f"- {m['metric_name']}: {m['metric_value']}{unit}（質問:「{q}」/ {m.get('created_at')}）")
    return "\n".join(lines)[:4000]


def run_review_ai(provider: str, api_key: str, gemini_api_key: str,
                  user_question: str, executed_queries: list,
                  result_data, answer: str, doc_context_section: str = "",
                  baseline_section: str = "", past_metrics_section: str = "") -> dict:
    """Run a reviewer AI that validates the final analysis answer.

    Returns a dict: {"verdict": "pass"|"needs_revision", "issues": [...], "reason": str}
    Fails open (returns pass with note) if the reviewer call itself fails.
    """
    # Build a compact result summary to keep tokens bounded
    result_summary = "（クエリ結果なし）"
    if result_data:
        try:
            summary_obj = {
                "row_count": len(result_data),
                "columns": list(result_data[0].keys()) if result_data else [],
                "sample_rows": result_data[:5]
            }
            result_summary = json.dumps(summary_obj, ensure_ascii=False, cls=DateTimeEncoder)[:4000]
        except Exception:
            result_summary = str(result_data)[:2000]

    sql_section = "\n\n".join(f"-- Query {i+1}\n{q}" for i, q in enumerate(executed_queries)) if executed_queries else "（SQLは実行されませんでした）"

    doc_block = ""
    extra_viewpoints = []
    if doc_context_section:
        doc_block += f"""

## プロジェクトの目的・分析方針（提案書・議事録から抽出）
{doc_context_section}
"""
        extra_viewpoints.append("回答・示唆が上記の「プロジェクトの目的・分析方針」に沿っているか（目的と無関係な示唆や、方針・KPI定義と矛盾する解釈がないか）")
    if baseline_section:
        doc_block += f"""

## テーブル基礎集計（数値整合性チェックの基準値）
{baseline_section}
"""
        extra_viewpoints.append("回答中の数値が上記の「テーブル基礎集計」と整合しているか（例: 件数が全体行数を超えていないか、合計値が基準の合計を超えていないか、期間がデータの存在範囲外でないか）。矛盾があれば指摘の先頭に【数値齟齬】を付け、どの数値がどの基準値とどう矛盾するかを具体的に示すこと")
    if past_metrics_section:
        doc_block += f"""

## 過去の分析で得られた主要指標
{past_metrics_section}
"""
        extra_viewpoints.append("同種の質問に対する過去の指標と比べて桁違いの乖離（10倍以上など）がないか。乖離があり説明がない場合は指摘の先頭に【数値齟齬】を付け、過去の値と今回の値を併記して警告すること（期間や条件が異なる正当な理由がSQLから読み取れる場合は指摘不要）")
    doc_viewpoint = "".join(f"\n{i + 5}. {v}" for i, v in enumerate(extra_viewpoints))

    review_prompt = f"""あなたはBigQueryデータ分析のレビュアーです。別のAIアナリストが生成した分析結果を厳密に検証してください。

## ユーザーの質問
{user_question}

## 実行されたSQL
{sql_section}

## クエリ結果サマリー
{result_summary}

## アナリストの最終回答
{answer}{doc_block}

## 検証観点
1. SQLがユーザーの質問の意図（指標・期間・フィルタ・粒度）に合致しているか
2. JOIN条件・期間解釈・集計方法などに意味的な誤りがないか
3. 最終回答がクエリ結果と整合しており、解釈が妥当か
4. 重大な見落とし（NULL処理、重複、単位の誤りなど）がないか{doc_viewpoint}

軽微な文言・スタイルの問題は指摘しないでください。分析の正しさに影響する問題のみ指摘してください。

以下のJSON形式のみで回答してください:
{{"verdict": "pass" または "needs_revision", "issues": ["具体的な指摘事項（日本語）", ...], "reason": "判定理由の要約（日本語、1〜2文）"}}"""

    try:
        if provider == 'gemini':
            client = genai_new.Client(api_key=gemini_api_key)
            response = client.models.generate_content(
                model='gemini-2.5-flash',
                contents=review_prompt,
                config=genai_types.GenerateContentConfig(response_mime_type="application/json")
            )
            text = response.text
        else:
            client = OpenAI(api_key=api_key)
            response = client.chat.completions.create(
                model=OPENAI_MODEL,
                messages=[
                    {"role": "system", "content": "You are a strict data analysis reviewer. Respond only with valid JSON."},
                    {"role": "user", "content": review_prompt}
                ],
                response_format={"type": "json_object"}
            )
            text = response.choices[0].message.content

        parsed = json.loads(text)
        verdict = parsed.get("verdict")
        if verdict not in ("pass", "needs_revision"):
            verdict = "pass"
        issues = parsed.get("issues") or []
        if not isinstance(issues, list):
            issues = [str(issues)]
        return {
            "verdict": verdict,
            "issues": [str(i) for i in issues],
            "reason": str(parsed.get("reason", ""))
        }
    except Exception as e:
        print(f"WARNING: Review AI call failed: {e}")
        return {
            "verdict": "pass",
            "issues": [],
            "reason": f"レビューAIの呼び出しに失敗したためスキップしました: {e}",
            "review_failed": True
        }

async def extract_payload_text(res) -> str:
    """Extract text from MCP response"""
    for p in res.content:
        if getattr(p, "type", None) == "application/json" and getattr(p, "text", None):
            return p.text
    for p in res.content:
        if hasattr(p, "text") and p.text:
            return p.text
    return ""

async def call_tool_text(session: ClientSession, tool: str, args: Dict[str, Any]) -> str:
    """Call MCP tool and return text response"""
    res = await session.call_tool(tool, arguments=args)
    text = await extract_payload_text(res)
    if not text:
        raise RuntimeError(f"{tool} returned empty payload.")
    return text

async def list_tables(session: ClientSession, project_id: str) -> List[str]:
    """List all tables in BigQuery project"""
    text = await call_tool_text(session, "list-tables", {"project": project_id})
    try:
        js = json.loads(text)
        if isinstance(js, dict):
            arr = js.get("tables") or js.get("data") or []
            if isinstance(arr, list):
                return [str(t) if isinstance(t, str) else t.get("table", "") for t in arr]
    except:
        pass
    lines = [ln.strip() for ln in text.splitlines() if ln.strip()]
    return lines

def parse_ddl_schema(ddl: str) -> List[Dict[str, str]]:
    """Parse BigQuery DDL to extract schema information"""
    import re
    schema = []
    
    # DDL から カラム定義部分を抽出（括弧内）
    match = re.search(r'\((.*)\)', ddl, re.DOTALL)
    if not match:
        return []
    
    columns_section = match.group(1)
    
    # 各行を処理
    for line in columns_section.split('\n'):
        line = line.strip()
        if not line or line.startswith('--'):
            continue
        
        # カラム名と型を抽出: "column_name TYPE" または "column_name TYPE,"
        # オプション句も考慮
        match = re.match(r'^([a-zA-Z0-9_]+)\s+([A-Z0-9_<>]+)', line)
        if match:
            col_name = match.group(1)
            col_type = match.group(2)
            schema.append({"name": col_name, "type": col_type})
    
    return schema

async def describe_table(session: ClientSession, project_id: str, dataset: str, table: str) -> List[Dict[str, str]]:
    """Get table schema"""
    variants = [
        {"project": project_id, "dataset": dataset, "table_name": table},
        {"project": project_id, "table_name": f"{dataset}.{table}"},
        {"dataset": dataset, "table_name": table},
        {"table_name": f"{dataset}.{table}"},
        {"table_name": f"{project_id}.{dataset}.{table}"},  # フルパス形式も試す
    ]
    
    for args in variants:
        try:
            text = await call_tool_text(session, "describe-table", args)
            print(f"DEBUG: describe-table response (first 200 chars): {text[:200]}")
            
            # まずJSONとして解析を試みる
            try:
                js = json.loads(text)
            except json.JSONDecodeError:
                # Pythonリテラル形式として評価
                print("DEBUG: describe-table JSON parse failed, trying Python literal eval...")
                try:
                    import datetime as dt_module
                    safe_dict = {
                        "datetime": dt_module,
                        "None": None,
                        "True": True,
                        "False": False,
                        "__builtins__": {}
                    }
                    exec(f"result = {text}", safe_dict)
                    js = safe_dict["result"]
                    print(f"DEBUG: Successfully evaluated describe-table Python literal")
                except Exception as eval_error:
                    print(f"DEBUG: describe-table Python literal eval failed: {eval_error}")
                    continue
            
            if isinstance(js, dict):
                # DDL形式の場合（'ddl'キーがある）
                ddl = js.get("ddl")
                if ddl and isinstance(ddl, str):
                    print(f"DEBUG: Parsing DDL: {ddl[:200]}...")
                    schema = parse_ddl_schema(ddl)
                    if schema:
                        return schema
                
                # 通常のスキーマ形式
                cols = js.get("columns") or js.get("schema") or js.get("fields")
                if isinstance(cols, list) and cols:
                    return [{"name": str(c.get("name", "")), "type": str(c.get("type", "UNKNOWN"))} for c in cols]
            elif isinstance(js, list):
                # リストが直接返される場合
                if js and isinstance(js[0], dict):
                    # DDL形式がリスト内にある場合
                    ddl = js[0].get("ddl")
                    if ddl and isinstance(ddl, str):
                        print(f"DEBUG: Parsing DDL from list: {ddl[:200]}...")
                        schema = parse_ddl_schema(ddl)
                        if schema:
                            return schema
                    # 通常のスキーマリスト
                    return [{"name": str(c.get("name", "")), "type": str(c.get("type", "UNKNOWN"))} for c in js if isinstance(c, dict)]
        except Exception as e:
            print(f"DEBUG: describe-table failed with args {args}: {e}")
            continue
    return []

def _make_bq_client(project_id: str, credentials_path: str = None) -> "bigquery.Client":
    """Create a BigQuery client with either a service account JSON key file or ADC"""
    if credentials_path and os.path.exists(credentials_path):
        creds = gcp_service_account.Credentials.from_service_account_file(credentials_path)
        return bigquery.Client(project=project_id, credentials=creds)
    # Fall back to Application Default Credentials (Cloud Run service account, etc.)
    return bigquery.Client(project=project_id)

def _format_bytes(num_bytes: int) -> str:
    """Format byte count into human readable string"""
    if num_bytes is None:
        return "不明"
    gb = num_bytes / (1024 ** 3)
    if gb >= 1:
        return f"{gb:.2f} GB"
    mb = num_bytes / (1024 ** 2)
    if mb >= 1:
        return f"{mb:.1f} MB"
    kb = num_bytes / 1024
    if kb >= 1:
        return f"{kb:.1f} KB"
    return f"{num_bytes} B"

def _serialize_bq_value(v):
    """Convert BigQuery cell values to JSON-friendly Python types"""
    import decimal
    import datetime as dt_module
    if isinstance(v, (dt_module.datetime, dt_module.date, dt_module.time)):
        return v.isoformat()
    if isinstance(v, decimal.Decimal):
        return float(v)
    if isinstance(v, bytes):
        return v.decode("utf-8", errors="replace")
    if isinstance(v, list):
        return [_serialize_bq_value(x) for x in v]
    if isinstance(v, dict):
        return {k: _serialize_bq_value(x) for k, x in v.items()}
    return v

# Transient BigQuery errors that are worth retrying with exponential backoff
_BQ_TRANSIENT_ERRORS = (
    gcp_exceptions.TooManyRequests,      # 429 rate limit
    gcp_exceptions.InternalServerError,  # 500
    gcp_exceptions.BadGateway,           # 502
    gcp_exceptions.ServiceUnavailable,   # 503
    gcp_exceptions.GatewayTimeout,       # 504
)

def _run_bq_query_sync(project_id: str, location: str, sql: str, timeout: int,
                       credentials_path: str = None, max_scan_gb: float = None,
                       add_step=None) -> List[Dict[str, Any]]:
    """Execute a BigQuery query synchronously with dry-run cost check, scan cap, and retries"""
    import time as time_module

    if max_scan_gb is None or max_scan_gb <= 0:
        max_scan_gb = DEFAULT_MAX_SCAN_GB
    max_scan_bytes = int(max_scan_gb * (1024 ** 3))

    try:
        client = _make_bq_client(project_id, credentials_path)
    except Exception as e:
        raise RuntimeError(
            f"BigQuery authentication failed: {e}. "
            f"Check that a service account JSON is uploaded in settings, or ADC is available (Cloud Run)."
        )
    labels = {"app": "bq-ai-agent", "component": "chat-agent"}

    # --- 1) Dry-run: estimate scan size before executing ---
    dry_config = bigquery.QueryJobConfig(dry_run=True, use_query_cache=False, labels=labels)
    try:
        dry_job = client.query(sql, job_config=dry_config, location=location or None)
    except gcp_exceptions.BadRequest as e:
        raise RuntimeError(f"SQL validation failed (dry-run): {e.message if hasattr(e, 'message') else e}")
    except gcp_exceptions.Forbidden as e:
        raise RuntimeError(f"BigQuery permission error (403): {e}")
    except gcp_exceptions.NotFound as e:
        raise RuntimeError(f"BigQuery resource not found (404): {e}")

    # Enforce read-only at the job level (string checks can be bypassed)
    stmt_type = getattr(dry_job, "statement_type", None)
    if stmt_type and stmt_type != "SELECT":
        raise RuntimeError(f"Non read-only statement type '{stmt_type}' is blocked. Only SELECT queries are allowed.")

    estimated_bytes = dry_job.total_bytes_processed or 0
    estimated_str = _format_bytes(estimated_bytes)
    print(f"DEBUG: Dry-run estimate: {estimated_str} (limit: {max_scan_gb} GB)")
    if add_step:
        add_step(f"💰 スキャン見積もり: {estimated_str}（上限 {max_scan_gb:g} GB）")

    if estimated_bytes > max_scan_bytes:
        raise RuntimeError(
            f"Query blocked before execution: estimated scan size {estimated_str} exceeds the project scan limit "
            f"of {max_scan_gb:g} GB. Rewrite the query to scan less data: "
            f"1) Filter on partitioned columns (e.g. _PARTITIONDATE / date column), "
            f"2) SELECT only needed columns instead of SELECT *, "
            f"3) Aggregate or sample the data. Note: LIMIT does NOT reduce scanned bytes."
        )

    # --- 2) Real execution with maximum_bytes_billed cap and retries ---
    job_config = bigquery.QueryJobConfig(
        maximum_bytes_billed=max_scan_bytes,
        labels=labels,
    )

    max_attempts = 3
    last_error = None
    for attempt in range(1, max_attempts + 1):
        try:
            query_job = client.query(sql, job_config=job_config, location=location or None)
            rows_iter = query_job.result(timeout=timeout)
            rows = [
                {k: _serialize_bq_value(v) for k, v in dict(row).items()}
                for row in rows_iter
            ]
            actual_bytes = query_job.total_bytes_processed
            if actual_bytes is not None:
                print(f"DEBUG: Query completed. Actual bytes processed: {_format_bytes(actual_bytes)}")
            return rows
        except _BQ_TRANSIENT_ERRORS as e:
            last_error = e
            if attempt < max_attempts:
                wait = 2 ** (attempt - 1)  # 1s, 2s
                print(f"WARNING: Transient BigQuery error (attempt {attempt}/{max_attempts}): {e}. Retrying in {wait}s...")
                if add_step:
                    add_step(f"⏳ 一時的なエラーのためリトライします（{attempt}/{max_attempts - 1}回目、{wait}秒待機）")
                time_module.sleep(wait)
                continue
            raise RuntimeError(f"BigQuery transient error persisted after {max_attempts} attempts: {e}")
        except gcp_exceptions.BadRequest as e:
            msg = str(e)
            if "bytesBilledLimitExceeded" in msg or "maximum_bytes_billed" in msg or "bytes billed" in msg.lower():
                raise RuntimeError(
                    f"Query aborted: it attempted to bill more than the scan limit ({max_scan_gb:g} GB). "
                    f"Rewrite the query to scan less data (filter partitions, select fewer columns, aggregate)."
                )
            raise RuntimeError(f"BigQuery query error: {msg[:500]}")
        except gcp_exceptions.Forbidden as e:
            raise RuntimeError(f"BigQuery permission error (403). Service account needs 'BigQuery Job User' role: {e}")
        except gcp_exceptions.NotFound as e:
            raise RuntimeError(f"BigQuery resource not found (404). Check project/dataset/table names: {e}")
        except TimeoutError:
            raise RuntimeError(
                f"Query execution timed out after {timeout} seconds. Try: 1) Filtering data with WHERE, "
                f"2) Using aggregation, 3) Reducing joined data."
            )
    raise RuntimeError(f"Query failed. Last error: {last_error}")

async def execute_query(session: ClientSession, project_id: str, location: str, sql: str, timeout: int = 300,
                        credentials_path: str = None, max_scan_gb: float = None,
                        add_step=None) -> List[Dict[str, Any]]:
    """Execute BigQuery SQL directly via google-cloud-bigquery with cost controls.

    The MCP session argument is kept for call-site compatibility but is no longer used;
    query execution now goes through the official BigQuery client library so that
    dry-run estimates, maximum_bytes_billed, labels, and retries can be applied.
    """
    lower = sql.lower()
    forbidden = ("insert ", "update ", "delete ", "merge ", "create ", "drop ", "alter ", "truncate ")
    if any(tok in lower for tok in forbidden):
        raise RuntimeError("Non read-only statements are blocked. Only SELECT queries are allowed.")

    if credentials_path is None and GCP_SA_JSON:
        credentials_path = GCP_SA_JSON

    try:
        return await asyncio.wait_for(
            asyncio.to_thread(
                _run_bq_query_sync, project_id, location, sql, timeout,
                credentials_path, max_scan_gb, add_step
            ),
            timeout=timeout + 30  # outer guard slightly above job timeout
        )
    except asyncio.TimeoutError:
        error_msg = f"Query execution timed out after {timeout} seconds. The query may be too complex or processing too much data. Try: 1) Filtering data with WHERE, 2) Using aggregation."
        print(f"TIMEOUT: {error_msg}")
        raise RuntimeError(error_msg)

def execute_python_code(code: str, dataframe_dict: Dict[str, pd.DataFrame] = None, timeout: int = 30) -> Dict[str, Any]:
    """
    Execute Python code in a restricted sandbox environment with access to BigQuery results as DataFrames
    
    Args:
        code: Python code to execute
        dataframe_dict: Dictionary of DataFrames to make available (e.g., {'df': query_result_dataframe})
        timeout: Execution timeout in seconds
        
    Returns:
        Dictionary with 'output' (stdout), 'result' (last expression value), 'plots' (base64 encoded images)
    """
    import io
    import sys
    import base64
    import matplotlib
    matplotlib.use('Agg')  # Non-interactive backend
    import matplotlib.pyplot as plt
    import numpy as np
    import seaborn as sns
    import sklearn
    from sklearn import (
        linear_model, tree, ensemble, svm, naive_bayes, 
        neighbors, cluster, decomposition, preprocessing,
        model_selection, metrics
    )
    import scipy
    from scipy import stats
    import statsmodels.api as sm
    
    # Capture stdout
    captured_output = io.StringIO()
    old_stdout = sys.stdout
    sys.stdout = captured_output
    
    # Safe import function that only allows specific modules
    def safe_import(name, globals=None, locals=None, fromlist=(), level=0):
        allowed_modules = {
            'pandas', 'numpy', 'matplotlib', 'seaborn', 'sklearn', 'scipy', 'statsmodels',
            'matplotlib.pyplot', 'sklearn.ensemble', 'sklearn.model_selection', 
            'sklearn.metrics', 'sklearn.linear_model', 'sklearn.tree', 'sklearn.svm',
            'sklearn.naive_bayes', 'sklearn.neighbors', 'sklearn.cluster',
            'sklearn.decomposition', 'sklearn.preprocessing', 'scipy.stats'
        }
        # Allow submodule imports
        if name in allowed_modules or any(name.startswith(mod + '.') for mod in allowed_modules):
            return __import__(name, globals, locals, fromlist, level)
        raise ImportError(f"Import of '{name}' is not allowed in sandbox")
    
    # Prepare restricted globals
    safe_globals = {
        '__builtins__': {
            'print': print,
            'len': len,
            'range': range,
            'enumerate': enumerate,
            'zip': zip,
            'map': map,
            'filter': filter,
            'sorted': sorted,
            'sum': sum,
            'min': min,
            'max': max,
            'abs': abs,
            'round': round,
            'int': int,
            'float': float,
            'str': str,
            'bool': bool,
            'list': list,
            'dict': dict,
            'tuple': tuple,
            'set': set,
            'type': type,
            'isinstance': isinstance,
            'hasattr': hasattr,
            'getattr': getattr,
            'True': True,
            'False': False,
            'None': None,
            '__import__': safe_import,
        },
        'pd': pd,
        'np': np,
        'plt': plt,
        'sns': sns,
        'DataFrame': pd.DataFrame,
        'Series': pd.Series,
        # Scikit-learn modules
        'sklearn': sklearn,
        'linear_model': linear_model,
        'tree': tree,
        'ensemble': ensemble,
        'svm': svm,
        'naive_bayes': naive_bayes,
        'neighbors': neighbors,
        'cluster': cluster,
        'decomposition': decomposition,
        'preprocessing': preprocessing,
        'model_selection': model_selection,
        'metrics': metrics,
        # Scipy
        'scipy': scipy,
        'stats': stats,
        # Statsmodels
        'sm': sm,
        'statsmodels': sm,
    }
    
    # Add DataFrames to globals
    if dataframe_dict:
        safe_globals.update(dataframe_dict)
    
    result = None
    plots = []
    
    try:
        # Clear any existing plots
        plt.close('all')
        
        # Compile and execute code
        compiled_code = compile(code, '<string>', 'exec')
        exec(compiled_code, safe_globals)
        
        # Check if there's a result variable
        if 'result' in safe_globals:
            result = safe_globals['result']
        
        # Capture any matplotlib plots
        figs = [plt.figure(n) for n in plt.get_fignums()]
        for fig in figs:
            buf = io.BytesIO()
            fig.savefig(buf, format='png', bbox_inches='tight', dpi=100)
            buf.seek(0)
            img_base64 = base64.b64encode(buf.read()).decode('utf-8')
            plots.append(img_base64)
            buf.close()
        
        plt.close('all')
        
    except Exception as e:
        sys.stdout = old_stdout
        raise RuntimeError(f"Python execution error: {str(e)}")
    finally:
        sys.stdout = old_stdout
    
    output_text = captured_output.getvalue()
    
    return {
        'output': output_text,
        'result': str(result) if result is not None else None,
        'plots': plots
    }

def build_openai_tools_schema() -> List[Dict[str, Any]]:
    """Build OpenAI function calling schema"""
    return [
        {
            "type": "function",
            "function": {
                "name": "list_tables",
                "description": "List all tables in BigQuery project",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "project": {"type": "string", "description": "GCP project ID"}
                    },
                    "required": ["project"]
                }
            }
        },
        {
            "type": "function",
            "function": {
                "name": "describe_table",
                "description": "Get the complete schema (column names, data types, and descriptions) of a BigQuery table. ALWAYS call this before writing SQL to ensure you use correct column names and understand data types. Pay special attention to: 1) Exact column name spelling (case-sensitive), 2) Data types (STRING, INTEGER, FLOAT64, TIMESTAMP, DATE, etc.), 3) Which columns are suitable for aggregation vs grouping.",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "project": {"type": "string", "description": "GCP project ID"},
                        "dataset": {"type": "string", "description": "Dataset name"},
                        "table": {"type": "string", "description": "Table name (exact spelling from list_tables)"}
                    },
                    "required": ["project", "dataset", "table"]
                }
            }
        },
        {
            "type": "function",
            "function": {
                "name": "execute_query",
                "description": "Execute a read-only BigQuery SQL SELECT query. IMPORTANT: 1) Always use backticks for table names: `project.dataset.table` or `dataset.table`, 2) Use exact column names from describe_table (case-sensitive), 3) Use SAFE_CAST() for type conversions to avoid errors, 4) Include ORDER BY for visualization data, 5) Use LIMIT to prevent large result sets, 6) Handle NULLs with IFNULL() or WHERE column IS NOT NULL. Only SELECT queries are allowed - no INSERT, UPDATE, DELETE, CREATE, DROP.",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "project": {"type": "string", "description": "GCP project ID"},
                        "location": {"type": "string", "description": "BigQuery location (e.g., 'US', 'asia-northeast1')"},
                        "query": {"type": "string", "description": "Complete SQL SELECT query with proper BigQuery syntax. Must use backticks for table names, exact column names, and include ORDER BY for sorted results."}
                    },
                    "required": ["project", "location", "query"]
                }
            }
        },
        {
            "type": "function",
            "function": {
                "name": "suggest_chart",
                "description": "Suggest the best chart visualization based on query results. Analyze the data structure and column types to determine the most appropriate chart type (bar, line, pie, doughnut, scatter) and axis configuration.",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "chart_type": {
                            "type": "string", 
                            "enum": ["bar", "line", "pie", "doughnut", "scatter", "none"],
                            "description": "Chart type: 'bar' for categorical comparisons, 'line' for time series/trends, 'pie' or 'doughnut' for part-to-whole relationships (max 10 categories), 'scatter' for correlations, 'none' if data is not suitable for visualization"
                        },
                        "x_axis": {"type": "string", "description": "Column name for X-axis (labels). Use date/time columns for line charts, categories for bar charts"},
                        "y_axis": {"type": "string", "description": "Column name for Y-axis (values). Must be numeric"},
                        "title": {"type": "string", "description": "Chart title in Japanese describing what the chart shows"}
                    },
                    "required": ["chart_type", "x_axis", "y_axis", "title"]
                }
            }
        },
        {
            "type": "function",
            "function": {
                "name": "execute_python",
                "description": "Execute Python code in a secure sandbox environment with access to BigQuery query results as Pandas DataFrames. Use this for advanced data analysis, machine learning, custom visualizations, statistical calculations, or data transformations that cannot be done with SQL alone. IMPORTANT: You CAN import and use external libraries in this sandbox. Available and importable libraries include: pandas, numpy, matplotlib, seaborn, sklearn (scikit-learn with all submodules like ensemble.RandomForestClassifier), scipy, and statsmodels. You can use import statements like 'from sklearn.ensemble import RandomForestClassifier' or 'import pandas as pd'. Pre-imported aliases available: pd (pandas), np (numpy), plt (matplotlib.pyplot), sns (seaborn), sklearn modules (linear_model, tree, ensemble, svm, naive_bayes, neighbors, cluster, decomposition, preprocessing, model_selection, metrics), scipy.stats (stats), statsmodels.api (sm). Query results are available as 'df' DataFrame. Use 'result' variable to return values. Use plt for creating custom charts.",
                "parameters": {
                    "type": "object",
                    "properties": {
                        "code": {
                            "type": "string",
                            "description": "Python code to execute. The most recent BigQuery query result is available as 'df' (a pandas DataFrame). Use print() for output, set 'result' variable for return values, and use plt.figure()/plt.plot()/etc for visualizations. Example: 'print(df.describe())' or 'result = df.groupby(\"category\")[\"value\"].mean()' or 'plt.figure(); plt.hist(df[\"column\"]); plt.title(\"Distribution\")'"
                        }
                    },
                    "required": ["code"]
                }
            }
        }
    ]

def build_gemini_tools_schema() -> List[Dict[str, Any]]:
    """Build Gemini function calling schema from OpenAI schema"""
    
    def convert_type_to_gemini(obj):
        """Recursively convert 'type' fields from lowercase to uppercase for Gemini"""
        if isinstance(obj, dict):
            new_obj = {}
            for key, value in obj.items():
                if key == "type" and isinstance(value, str):
                    # Convert type to uppercase (object -> OBJECT, string -> STRING, etc.)
                    new_obj[key] = value.upper()
                else:
                    new_obj[key] = convert_type_to_gemini(value)
            return new_obj
        elif isinstance(obj, list):
            return [convert_type_to_gemini(item) for item in obj]
        else:
            return obj
    
    openai_tools = build_openai_tools_schema()
    gemini_tools = []
    
    for tool in openai_tools:
        func = tool["function"]
        gemini_tool = {
            "name": func["name"],
            "description": func["description"],
            "parameters": convert_type_to_gemini(func["parameters"])
        }
        gemini_tools.append(gemini_tool)
    
    return gemini_tools

async def run_agent(user_question: str, conversation_history: List[Dict[str, str]], 
                    api_key: str = None, project_id: str = None, dataset_id: str = None, 
                    service_account_json: str = None, project_db_id: int = None, user_id: int = None,
                    provider: str = 'openai', gemini_api_key: str = None, use_adc: bool = False) -> Dict[str, Any]:
    """Main agent logic with MCP and OpenAI/Gemini"""
    # Use provided parameters or fall back to global variables
    api_key = api_key or OPENAI_API_KEY
    project_id = project_id or PROJECT_ID
    # dataset_id can be None - AI will discover tables dynamically
    
    # Only use service_account_json if not using ADC
    if not use_adc:
        service_account_json = service_account_json or GCP_SA_JSON
    else:
        service_account_json = None
    
    # Load project memories if project_db_id is provided
    project_memories = []
    if project_db_id and user_id:
        try:
            conn = get_db_connection()
            cur = conn.cursor(cursor_factory=RealDictCursor)
            cur.execute('''
                SELECT memory_key, memory_value, category, table_name, updated_at
                FROM project_memories
                WHERE project_id = %s AND user_id = %s
                ORDER BY updated_at DESC
            ''', (project_db_id, user_id))
            project_memories = cur.fetchall()
            cur.close()
            conn.close()
        except Exception as e:
            print(f"Warning: Failed to load project memories: {e}")
    
    env = os.environ.copy()
    # Only set GOOGLE_APPLICATION_CREDENTIALS if not using ADC
    # For ADC (Cloud Run), credentials are automatically available
    if service_account_json and not use_adc:
        env["GOOGLE_APPLICATION_CREDENTIALS"] = service_account_json
    
    server_params = StdioServerParameters(
        command="mcp-server-bigquery",
        args=["--project", project_id, "--location", LOCATION],
        env=env
    )
    
    steps = []
    result_data = None
    result_charts = []  # 複数のグラフ設定を保持
    
    async with AsyncExitStack() as stack:
        stdio_transport = await stack.enter_async_context(stdio_client(server_params))
        stdio, write = stdio_transport
        session = await stack.enter_async_context(ClientSession(stdio, write))
        
        await session.initialize()
        
        client = OpenAI(api_key=api_key)
        
        # Build memory section if memories exist
        memory_section = build_memory_section(project_memories)
        
        # Build dataset info section
        if dataset_id:
            dataset_info = f"- Default Dataset: {dataset_id}"
        else:
            dataset_info = "- Default Dataset: (Not set - use list_tables to discover available datasets and tables)"
        
        messages = [
            {"role": "system", "content": f"""You are an expert BigQuery data analyst assistant with deep knowledge of SQL optimization and data analysis.

## ENVIRONMENT
- BigQuery Project: {project_id}
{dataset_info}
{memory_section}
## STEP-BY-STEP REASONING FRAMEWORK
Follow this systematic approach for every user query:

### 1. UNDERSTAND THE QUESTION
- Identify the key metrics, dimensions, and filters requested
- Determine the time period if applicable
- Clarify any ambiguous terms by analyzing context

### 2. DISCOVER RELEVANT TABLES
- Use list_tables to see all available tables
- Identify tables that likely contain the requested data based on naming patterns
- Consider multiple tables if joins might be needed

### 3. EXAMINE TABLE SCHEMAS
- Use describe_table for each relevant table
- Carefully note:
  * Column names and their exact spelling
  * Data types (STRING, INTEGER, FLOAT, TIMESTAMP, DATE, BOOL, etc.)
  * Which columns can be used for aggregation vs grouping
  * Potential join keys between tables

### 4. CONSTRUCT SQL QUERY WITH BEST PRACTICES

**BigQuery-Specific Syntax:**
- Always use backticks for table references: `project.dataset.table` or `dataset.table`
- Use DATETIME/TIMESTAMP functions: CURRENT_DATETIME(), DATETIME_DIFF(), FORMAT_DATETIME()
- For dates: CURRENT_DATE(), DATE_DIFF(), FORMAT_DATE()
- Use SAFE_CAST() instead of CAST() to avoid errors
- Leverage IFNULL() or COALESCE() for NULL handling

**SQL Quality Guidelines:**
- Write explicit column names (avoid SELECT *)
- Use meaningful aliases with AS keyword
- Apply appropriate WHERE clauses for filtering
- Use proper GROUP BY with all non-aggregated columns
- Add ORDER BY for sorted results (essential for visualizations)
- Limit results appropriately (use LIMIT for large datasets)
- Handle NULLs explicitly to avoid incorrect aggregations

**Common Patterns:**
```sql
-- Aggregation with grouping
SELECT 
  column1,
  SUM(SAFE_CAST(column2 AS FLOAT64)) AS total_value,
  COUNT(*) AS count
FROM `{PROJECT_ID}.{DEFAULT_DATASET}.table_name`
WHERE column3 IS NOT NULL
GROUP BY column1
ORDER BY total_value DESC
LIMIT 100

-- Time-based analysis
SELECT 
  FORMAT_DATE('%Y-%m', date_column) AS month,
  AVG(numeric_column) AS avg_value
FROM `{PROJECT_ID}.{DEFAULT_DATASET}.table_name`
WHERE date_column >= DATE_SUB(CURRENT_DATE(), INTERVAL 6 MONTH)
GROUP BY month
ORDER BY month ASC

-- String matching
SELECT * 
FROM `{PROJECT_ID}.{DEFAULT_DATASET}.table_name`
WHERE LOWER(text_column) LIKE LOWER('%search_term%')
```

### 5. EXECUTE AND VALIDATE
- Run execute_query with the constructed SQL
- Check if results make sense (row count, value ranges)
- If error occurs, analyze the error message and fix:
  * Column name typos → recheck schema
  * Type mismatches → use SAFE_CAST()
  * NULL issues → add NULL handling
  * Syntax errors → verify BigQuery syntax

### 6. ADVANCED ANALYSIS WITH PYTHON (Optional)
For complex analysis beyond SQL capabilities, use execute_python with query results:

**When to Use Python:**
- Statistical analysis (correlation, regression, hypothesis testing)
- Custom data transformations or complex calculations
- Advanced visualizations (histograms, heatmaps, custom plots)
- Machine learning or predictive analytics
- Data cleaning or preprocessing

**Python Environment:**
- Query results available as 'df' (pandas DataFrame)
- Available libraries: pandas (pd), numpy (np), matplotlib.pyplot (plt), seaborn (sns)
- Use print() for text output, set 'result' variable for return values
- Use plt for custom visualizations (automatically captured)

**Example Usage:**
```python
# Statistical analysis
print(df.describe())
print(df.corr())

# Custom visualization
plt.figure(figsize=(10, 6))
plt.hist(df['column_name'], bins=30)
plt.title('Distribution of Column')
plt.xlabel('Value')
plt.ylabel('Frequency')

# Advanced calculation
result = df.groupby('category')['value'].agg(['mean', 'std', 'count'])
```

### 7. VISUALIZE INTELLIGENTLY
After getting query results, choose between suggest_chart (simple) or execute_python (advanced):

**suggest_chart - Use for standard visualizations:**
- **Line chart**: Time series data (x=date/time, y=metric), trends over time
- **Bar chart**: Categorical comparisons (x=category, y=value), rankings
- **Pie/Doughnut**: Part-to-whole with ≤10 categories, percentages, distributions
- **Scatter**: Correlation between two numeric variables (x=variable1, y=variable2)
- **None**: Single values, text-heavy data, or >50 data points without clear pattern

**execute_python - Use for custom visualizations:**
- Histograms, box plots, heatmaps, pair plots
- Multiple subplots or complex chart combinations
- Statistical plots (QQ plots, distribution fits)
- Custom styling or annotations

**Axis Selection (for suggest_chart):**
- x_axis: Use the dimension/category column (dates for trends, categories for comparisons)
- y_axis: Use the numeric metric column (always numeric values)
- Ensure columns match exactly with query result column names

### 8. PROVIDE INSIGHTS
- Explain what the data shows
- Highlight key findings (trends, outliers, patterns)
- Answer the user's original question clearly

## CRITICAL RULES
1. **Always check schemas before writing SQL** - Never assume column names
2. **Use exact column names from describe_table** - BigQuery is case-sensitive
3. **Test one table at a time** - Don't assume relationships exist
4. **Handle NULLs explicitly** - They affect aggregations
5. **Use appropriate data types** - SAFE_CAST when unsure
6. **Fully qualify table names** - Use `project.dataset.table` or `dataset.table`
7. **Always provide Japanese explanations** - Users expect Japanese responses

Think step-by-step and show your reasoning process."""}
        ]
        
        messages.extend(conversation_history)
        messages.append({"role": "user", "content": user_question})
        
        # Provider-based implementation
        if provider == 'gemini':
            # =====  GEMINI API IMPLEMENTATION =====
            genai.configure(api_key=gemini_api_key)
            gemini_tools = build_gemini_tools_schema()
            model = genai.GenerativeModel(
                model_name='gemini-2.5-pro',
                tools=gemini_tools,
                system_instruction=messages[0]["content"]  # Use system prompt
            )
            
            # Convert conversation history to Gemini format
            gemini_contents = []
            for msg in conversation_history:
                if msg["role"] == "user":
                    gemini_contents.append({"role": "user", "parts": [{"text": msg["content"]}]})
                elif msg["role"] == "assistant":
                    gemini_contents.append({"role": "model", "parts": [{"text": msg["content"]}]})
            
            # Add current question
            gemini_contents.append({"role": "user", "parts": [{"text": user_question}]})
            
            # Main iteration loop for Gemini
            for iteration in range(10):
                try:
                    # Generate response
                    response = model.generate_content(gemini_contents)
                    
                    # Check for valid response
                    if not response.candidates or not response.candidates[0].content.parts:
                        return {
                            "answer": "Geminiからレスポンスがありませんでした",
                            "steps": steps,
                            "data": result_data,
                            "charts": result_charts
                        }
                    
                    response_parts = response.candidates[0].content.parts
                    
                    # Check for function calls
                    has_function_call = False
                    for part in response_parts:
                        if hasattr(part, 'function_call') and part.function_call:
                            has_function_call = True
                            function_call = part.function_call
                            func_name = function_call.name
                            func_args = dict(function_call.args)
                            
                            steps.append(f"🔧 {func_name}({json.dumps(func_args, ensure_ascii=False)})")
                            
                            # Execute tool
                            try:
                                if func_name == "list_tables":
                                    result = await list_tables(session, func_args.get("project", project_id))
                                    func_result = {"tables": result}
                                    
                                elif func_name == "describe_table":
                                    result = await describe_table(
                                        session,
                                        func_args.get("project", project_id),
                                        func_args.get("dataset", dataset_id),
                                        func_args.get("table")
                                    )
                                    func_result = {"columns": result}
                                    
                                elif func_name == "execute_query":
                                    result = await execute_query(
                                        session,
                                        func_args.get("project", project_id),
                                        func_args.get("location", LOCATION),
                                        func_args.get("query")
                                    )
                                    result_data = result
                                    steps.append(f"📊 {len(result)}行のデータを取得しました")
                                    
                                    if len(result) > 100:
                                        func_result = {
                                            "row_count": len(result),
                                            "columns": list(result[0].keys()) if result else [],
                                            "sample_rows": result[:5],
                                            "message": f"Large dataset with {len(result)} rows."
                                        }
                                    else:
                                        func_result = {"rows": result}
                                        
                                elif func_name == "suggest_chart":
                                    chart_config = {
                                        "chart_type": func_args.get("chart_type", "bar"),
                                        "x_axis": func_args.get("x_axis"),
                                        "y_axis": func_args.get("y_axis"),
                                        "title": func_args.get("title", "")
                                    }
                                    result_charts.append(chart_config)
                                    func_result = {"chart": chart_config}
                                    steps.append(f"📈 グラフを提案しました: {chart_config['chart_type']}")
                                    
                                elif func_name == "execute_python":
                                    code = func_args.get("code")
                                    df_dict = {}
                                    if result_data:
                                        df_dict['df'] = pd.DataFrame(result_data)
                                    python_result = execute_python_code(code, df_dict, timeout=30)
                                    func_result = {
                                        "output": python_result.get('output', ''),
                                        "result": python_result.get('result'),
                                        "plots_count": len(python_result.get('plots', []))
                                    }
                                    steps.append(f"🐍 Python実行完了")
                                    
                                else:
                                    func_result = {"error": "Unknown function"}
                                    
                            except Exception as e:
                                error_msg = f"Error executing {func_name}: {str(e)}"
                                steps.append(f"❌ {error_msg}")
                                func_result = {"error": error_msg}
                            
                            # Add function call and response to conversation
                            gemini_contents.append({
                                "role": "model",
                                "parts": [{"function_call": {"name": func_name, "args": func_args}}]
                            })
                            gemini_contents.append({
                                "role": "user",
                                "parts": [{"function_response": {"name": func_name, "response": func_result}}]
                            })
                    
                    if not has_function_call:
                        # Final answer
                        answer_text = ""
                        for part in response_parts:
                            if hasattr(part, 'text'):
                                answer_text += part.text
                        return {
                            "answer": answer_text,
                            "steps": steps,
                            "data": result_data,
                            "charts": result_charts
                        }
                        
                except Exception as e:
                    error_msg = f"Gemini API error: {str(e)}"
                    print(f"ERROR: {error_msg}")
                    return {
                        "answer": f"エラーが発生しました: {error_msg}",
                        "steps": steps,
                        "data": result_data,
                        "charts": result_charts
                    }
            
            return {
                "answer": "最大反復回数に達しました",
                "steps": steps,
                "data": result_data,
                "charts": result_charts
            }
        
        else:  # provider == 'openai'
            # ===== OPENAI API IMPLEMENTATION =====
            client = OpenAI(api_key=api_key)
            tools = build_openai_tools_schema()
            
            for iteration in range(10):
                # gpt-5の場合は専用パラメーターを使用
                api_params = {
                    "model": OPENAI_MODEL,
                    "messages": messages,
                    "tools": tools,
                    "tool_choice": "auto"
                }
                
                # gpt-5の場合は reasoning_effort を追加
                if OPENAI_MODEL.startswith("gpt-5"):
                    api_params["reasoning_effort"] = "high"  # 深い推論を実行
                
                response = client.chat.completions.create(**api_params)
                
                assistant_message = response.choices[0].message
                messages.append(assistant_message.model_dump())
                
                if not assistant_message.tool_calls:
                    return {
                        "answer": assistant_message.content,
                        "steps": steps,
                        "data": result_data,
                        "charts": result_charts  # 配列で返す
                    }
                
                for tool_call in assistant_message.tool_calls:
                    func_name = tool_call.function.name
                    func_args = json.loads(tool_call.function.arguments)
                    
                    steps.append(f"🔧 {func_name}({json.dumps(func_args, ensure_ascii=False)})")
                    
                    try:
                        if func_name == "list_tables":
                            result = await list_tables(session, func_args.get("project", PROJECT_ID))
                            func_result = json.dumps({"tables": result}, ensure_ascii=False)
                            
                        elif func_name == "describe_table":
                            result = await describe_table(
                                session,
                                func_args.get("project", PROJECT_ID),
                                func_args.get("dataset", DEFAULT_DATASET),
                                func_args["table"]
                            )
                            func_result = json.dumps({"columns": result}, ensure_ascii=False)
                            
                        elif func_name == "execute_query":
                            result = await execute_query(
                                session,
                                func_args.get("project", PROJECT_ID),
                                func_args.get("location", LOCATION),
                                func_args["query"]
                            )
                            result_data = result
                            steps.append(f"📊 Query returned {len(result)} rows")
                            
                            # If data is large, return summary to avoid token limit
                            if len(result) > 100:
                                summary = {
                                    "row_count": len(result),
                                    "columns": list(result[0].keys()) if result else [],
                                    "sample_rows": result[:5],
                                    "message": f"Large dataset with {len(result)} rows. Data stored for analysis."
                                }
                                func_result = json.dumps(summary, ensure_ascii=False, cls=DateTimeEncoder)
                            else:
                                func_result = json.dumps({"rows": result}, ensure_ascii=False, cls=DateTimeEncoder)
                        else:
                            func_result = json.dumps({"error": "Unknown function"})
                            
                    except Exception as e:
                        error_msg = f"Error executing {func_name}: {str(e)}"
                        steps.append(f"❌ {error_msg}")
                        func_result = json.dumps({"error": error_msg}, ensure_ascii=False)
                    
                    messages.append({
                        "role": "tool",
                        "tool_call_id": tool_call.id,
                        "content": func_result
                    })
            
            return {
                "answer": "Maximum iterations reached.",
                "steps": steps,
                "data": result_data,
                "charts": result_charts  # 配列で返す
            }

async def run_agent_streaming(user_question: str, conversation_history: List[Dict[str, str]], msg_queue: queue.Queue, 
                              api_key: str = None, project_id: str = None, dataset_id: str = None, 
                              service_account_json: str = None, project_db_id: int = None, user_id: int = None,
                              provider: str = 'openai', gemini_api_key: str = None, use_adc: bool = False) -> Dict[str, Any]:
    """Main agent logic with MCP and OpenAI/Gemini - with streaming progress"""
    # Use provided parameters or fall back to global variables
    api_key = api_key or OPENAI_API_KEY
    project_id = project_id or PROJECT_ID
    dataset_id = dataset_id or DEFAULT_DATASET
    
    # Only use service_account_json if not using ADC
    if not use_adc:
        service_account_json = service_account_json or GCP_SA_JSON
    else:
        service_account_json = None
    
    # Load project memories if project_db_id is provided
    project_memories = []
    if project_db_id and user_id:
        try:
            conn = get_db_connection()
            cur = conn.cursor(cursor_factory=RealDictCursor)
            cur.execute('''
                SELECT memory_key, memory_value, category, table_name, updated_at
                FROM project_memories
                WHERE project_id = %s AND user_id = %s
                ORDER BY updated_at DESC
            ''', (project_db_id, user_id))
            project_memories = cur.fetchall()
            cur.close()
            conn.close()
        except Exception as e:
            print(f"Warning: Failed to load project memories: {e}")
    
    env = os.environ.copy()
    # Only set GOOGLE_APPLICATION_CREDENTIALS if not using ADC
    if service_account_json and not use_adc:
        env["GOOGLE_APPLICATION_CREDENTIALS"] = service_account_json
    
    server_params = StdioServerParameters(
        command="mcp-server-bigquery",
        args=["--project", project_id, "--location", LOCATION],
        env=env
    )
    
    steps = []
    result_data = None
    result_charts = []  # 複数のグラフ設定を保持
    
    async with AsyncExitStack() as stack:
        stdio_transport = await stack.enter_async_context(stdio_client(server_params))
        stdio, write = stdio_transport
        session = await stack.enter_async_context(ClientSession(stdio, write))
        
        await session.initialize()
        
        client = OpenAI(api_key=api_key)
        
        # Build memory section if memories exist
        memory_section = build_memory_section(project_memories)
        
        # Build dataset info section
        if dataset_id:
            dataset_info = f"- Default Dataset: {dataset_id}"
        else:
            dataset_info = "- Default Dataset: (Not set - use list_tables to discover available datasets and tables)"
        
        messages = [
            {"role": "system", "content": f"""You are an expert BigQuery data analyst assistant with deep knowledge of SQL optimization and data analysis.

## ENVIRONMENT
- BigQuery Project: {project_id}
{dataset_info}
{memory_section}
## STEP-BY-STEP REASONING FRAMEWORK
Follow this systematic approach for every user query:

### 1. UNDERSTAND THE QUESTION
- Identify the key metrics, dimensions, and filters requested
- Determine the time period if applicable
- Clarify any ambiguous terms by analyzing context

### 2. DISCOVER RELEVANT TABLES
- Use list_tables to see all available tables
- Identify tables that likely contain the requested data based on naming patterns
- Consider multiple tables if joins might be needed

### 3. EXAMINE TABLE SCHEMAS
- Use describe_table for each relevant table
- Carefully note:
  * Column names and their exact spelling
  * Data types (STRING, INTEGER, FLOAT, TIMESTAMP, DATE, BOOL, etc.)
  * Which columns can be used for aggregation vs grouping
  * Potential join keys between tables

### 4. CONSTRUCT SQL QUERY WITH BEST PRACTICES

**BigQuery-Specific Syntax:**
- Always use backticks for table references: `project.dataset.table` or `dataset.table`
- Use DATETIME/TIMESTAMP functions: CURRENT_DATETIME(), DATETIME_DIFF(), FORMAT_DATETIME()
- For dates: CURRENT_DATE(), DATE_DIFF(), FORMAT_DATE()
- Use SAFE_CAST() instead of CAST() to avoid errors
- Leverage IFNULL() or COALESCE() for NULL handling

**SQL Quality Guidelines:**
- Write explicit column names (avoid SELECT *)
- Use meaningful aliases with AS keyword
- Apply appropriate WHERE clauses for filtering
- Use proper GROUP BY with all non-aggregated columns
- Add ORDER BY for sorted results (essential for visualizations)
- Limit results appropriately (use LIMIT for large datasets)
- Handle NULLs explicitly to avoid incorrect aggregations

**Common Patterns:**
```sql
-- Aggregation with grouping
SELECT 
  column1,
  SUM(SAFE_CAST(column2 AS FLOAT64)) AS total_value,
  COUNT(*) AS count
FROM `{PROJECT_ID}.{DEFAULT_DATASET}.table_name`
WHERE column3 IS NOT NULL
GROUP BY column1
ORDER BY total_value DESC
LIMIT 100

-- Time-based analysis
SELECT 
  FORMAT_DATE('%Y-%m', date_column) AS month,
  AVG(numeric_column) AS avg_value
FROM `{PROJECT_ID}.{DEFAULT_DATASET}.table_name`
WHERE date_column >= DATE_SUB(CURRENT_DATE(), INTERVAL 6 MONTH)
GROUP BY month
ORDER BY month ASC

-- String matching
SELECT * 
FROM `{PROJECT_ID}.{DEFAULT_DATASET}.table_name`
WHERE LOWER(text_column) LIKE LOWER('%search_term%')
```

### 5. EXECUTE AND VALIDATE
- Run execute_query with the constructed SQL
- Check if results make sense (row count, value ranges)
- If error occurs, analyze the error message and fix:
  * Column name typos → recheck schema
  * Type mismatches → use SAFE_CAST()
  * NULL issues → add NULL handling
  * Syntax errors → verify BigQuery syntax

### 6. ADVANCED ANALYSIS WITH PYTHON (Optional)
For complex analysis beyond SQL capabilities, use execute_python with query results:

**When to Use Python:**
- Statistical analysis (correlation, regression, hypothesis testing)
- Custom data transformations or complex calculations
- Advanced visualizations (histograms, heatmaps, custom plots)
- Machine learning or predictive analytics
- Data cleaning or preprocessing

**Python Environment:**
- Query results available as 'df' (pandas DataFrame)
- Available libraries: pandas (pd), numpy (np), matplotlib.pyplot (plt), seaborn (sns)
- Use print() for text output, set 'result' variable for return values
- Use plt for custom visualizations (automatically captured)

**Example Usage:**
```python
# Statistical analysis
print(df.describe())
print(df.corr())

# Custom visualization
plt.figure(figsize=(10, 6))
plt.hist(df['column_name'], bins=30)
plt.title('Distribution of Column')
plt.xlabel('Value')
plt.ylabel('Frequency')

# Advanced calculation
result = df.groupby('category')['value'].agg(['mean', 'std', 'count'])
```

### 7. VISUALIZE INTELLIGENTLY
After getting query results, choose between suggest_chart (simple) or execute_python (advanced):

**suggest_chart - Use for standard visualizations:**
- **Line chart**: Time series data (x=date/time, y=metric), trends over time
- **Bar chart**: Categorical comparisons (x=category, y=value), rankings
- **Pie/Doughnut**: Part-to-whole with ≤10 categories, percentages, distributions
- **Scatter**: Correlation between two numeric variables (x=variable1, y=variable2)
- **None**: Single values, text-heavy data, or >50 data points without clear pattern

**execute_python - Use for custom visualizations:**
- Histograms, box plots, heatmaps, pair plots
- Multiple subplots or complex chart combinations
- Statistical plots (QQ plots, distribution fits)
- Custom styling or annotations

**Axis Selection (for suggest_chart):**
- x_axis: Use the dimension/category column (dates for trends, categories for comparisons)
- y_axis: Use the numeric metric column (always numeric values)
- Ensure columns match exactly with query result column names

### 8. PROVIDE INSIGHTS
- Explain what the data shows
- Highlight key findings (trends, outliers, patterns)
- Answer the user's original question clearly

## CRITICAL RULES
1. **Always check schemas before writing SQL** - Never assume column names
2. **Use exact column names from describe_table** - BigQuery is case-sensitive
3. **Test one table at a time** - Don't assume relationships exist
4. **Handle NULLs explicitly** - They affect aggregations
5. **Use appropriate data types** - SAFE_CAST when unsure
6. **Fully qualify table names** - Use `project.dataset.table` or `dataset.table`
7. **Always provide Japanese explanations** - Users expect Japanese responses

Think step-by-step and show your reasoning process."""}
        ]
        
        messages.extend(conversation_history)
        messages.append({"role": "user", "content": user_question})
        
        # Provider-based implementation
        if provider == 'gemini':
            # ===== GEMINI API IMPLEMENTATION (STREAMING) =====
            genai.configure(api_key=gemini_api_key)
            gemini_tools = build_gemini_tools_schema()
            model = genai.GenerativeModel(
                model_name='gemini-2.5-pro',
                tools=gemini_tools,
                system_instruction=messages[0]["content"]
            )
            
            # Convert conversation history to Gemini format
            gemini_contents = []
            for msg in conversation_history:
                if msg["role"] == "user":
                    gemini_contents.append({"role": "user", "parts": [{"text": msg["content"]}]})
                elif msg["role"] == "assistant":
                    gemini_contents.append({"role": "model", "parts": [{"text": msg["content"]}]})
            
            # Add current question
            gemini_contents.append({"role": "user", "parts": [{"text": user_question}]})
            
            # Main iteration loop
            for iteration in range(10):
                msg_queue.put({"type": "thinking", "message": "次のアクションを考えています..."})
                
                try:
                    # Generate response
                    response = model.generate_content(gemini_contents)
                    
                    if not response.candidates or not response.candidates[0].content.parts:
                        return {
                            "answer": "Geminiからレスポンスがありませんでした",
                            "steps": steps,
                            "data": result_data,
                            "charts": result_charts
                        }
                    
                    response_parts = response.candidates[0].content.parts
                    
                    # Check for function calls
                    has_function_call = False
                    for part in response_parts:
                        if hasattr(part, 'function_call') and part.function_call:
                            has_function_call = True
                            function_call = part.function_call
                            func_name = function_call.name
                            func_args = dict(function_call.args)
                            
                            steps.append(f"🔧 {func_name}({json.dumps(func_args, ensure_ascii=False)})")
                            msg_queue.put({"type": "tool_call", "name": func_name, "args": func_args})
                            
                            # Execute tool
                            try:
                                if func_name == "list_tables":
                                    msg_queue.put({"type": "tool_start", "tool": "list_tables", "message": "テーブル一覧を取得中..."})
                                    result = await list_tables(session, func_args.get("project", project_id))
                                    func_result = {"tables": result}
                                    msg_queue.put({"type": "tool_done", "tool": "list_tables", "message": f"{len(result)}個のテーブルを発見しました"})
                                    
                                elif func_name == "describe_table":
                                    table_name = func_args.get("table") or func_args.get("table_name")
                                    if not table_name:
                                        raise ValueError(f"Missing required parameter: 'table'")
                                    msg_queue.put({"type": "tool_start", "tool": "describe_table", "message": f"テーブル '{table_name}' のスキーマを取得中..."})
                                    result = await describe_table(
                                        session,
                                        func_args.get("project", project_id),
                                        func_args.get("dataset", dataset_id),
                                        table_name
                                    )
                                    func_result = {"columns": result}
                                    msg_queue.put({"type": "tool_done", "tool": "describe_table", "message": f"{len(result)}個のカラムを発見しました"})
                                    
                                elif func_name == "execute_query":
                                    query = func_args.get("query") or func_args.get("sql")
                                    if not query:
                                        raise ValueError(f"Missing required parameter: 'query'")
                                    msg_queue.put({"type": "tool_start", "tool": "execute_query", "message": "クエリを実行中...", "query": query})
                                    result = await execute_query(
                                        session,
                                        func_args.get("project", project_id),
                                        func_args.get("location", LOCATION),
                                        query
                                    )
                                    result_data = result
                                    msg = f"📊 {len(result)}行のデータを取得しました"
                                    steps.append(msg)
                                    msg_queue.put({"type": "tool_done", "tool": "execute_query", "message": msg})
                                    
                                    if len(result) > 100:
                                        func_result = {
                                            "row_count": len(result),
                                            "columns": list(result[0].keys()) if result else [],
                                            "sample_rows": result[:5],
                                            "message": f"Large dataset with {len(result)} rows."
                                        }
                                    else:
                                        func_result = {"rows": result}
                                        
                                elif func_name == "suggest_chart":
                                    msg_queue.put({"type": "tool_start", "tool": "suggest_chart", "message": "グラフ設定を提案中..."})
                                    chart_config = {
                                        "chart_type": func_args.get("chart_type", "bar"),
                                        "x_axis": func_args.get("x_axis"),
                                        "y_axis": func_args.get("y_axis"),
                                        "title": func_args.get("title", "")
                                    }
                                    result_charts.append(chart_config)
                                    func_result = {"chart": chart_config}
                                    chart_type_ja = {
                                        "bar": "棒グラフ",
                                        "line": "折れ線グラフ",
                                        "pie": "円グラフ",
                                        "doughnut": "ドーナツグラフ",
                                        "scatter": "散布図",
                                        "none": "グラフなし"
                                    }.get(chart_config["chart_type"], chart_config["chart_type"])
                                    msg = f"📈 {chart_type_ja}を提案しました"
                                    steps.append(msg)
                                    msg_queue.put({"type": "tool_done", "tool": "suggest_chart", "message": msg})
                                    
                                elif func_name == "execute_python":
                                    msg_queue.put({"type": "tool_start", "tool": "execute_python", "message": "Pythonコードを実行中..."})
                                    code = func_args.get("code")
                                    df_dict = {}
                                    if result_data:
                                        df_dict['df'] = pd.DataFrame(result_data)
                                    python_result = execute_python_code(code, df_dict, timeout=30)
                                    func_result = {
                                        "output": python_result.get('output', ''),
                                        "result": python_result.get('result'),
                                        "plots_count": len(python_result.get('plots', []))
                                    }
                                    msg_queue.put({"type": "tool_done", "tool": "execute_python", "message": "Python実行完了"})
                                    steps.append(f"🐍 Python実行完了")
                                    
                                else:
                                    func_result = {"error": "Unknown function"}
                                    
                            except Exception as e:
                                error_msg = str(e)
                                error_details = {
                                    "error": error_msg,
                                    "function": func_name,
                                    "arguments": func_args
                                }
                                steps.append(f"❌ {func_name} エラー: {error_msg}")
                                msg_queue.put({"type": "error", "message": f"エラー: {error_msg}"})
                                func_result = error_details
                            
                            # Add function call and response to conversation
                            gemini_contents.append({
                                "role": "model",
                                "parts": [{"function_call": {"name": func_name, "args": func_args}}]
                            })
                            gemini_contents.append({
                                "role": "user",
                                "parts": [{"function_response": {"name": func_name, "response": func_result}}]
                            })
                    
                    if not has_function_call:
                        # Final answer - stream it character by character
                        answer_text = ""
                        for part in response_parts:
                            if hasattr(part, 'text'):
                                for char in part.text:
                                    msg_queue.put({"type": "assistant_text", "text": char})
                                answer_text += part.text
                        return {
                            "answer": answer_text,
                            "steps": steps,
                            "data": result_data,
                            "charts": result_charts
                        }
                        
                except Exception as e:
                    error_msg = f"Gemini API error: {str(e)}"
                    print(f"ERROR: {error_msg}")
                    msg_queue.put({"type": "error", "message": f"エラーが発生しました: {error_msg}"})
                    return {
                        "answer": f"エラーが発生しました: {error_msg}",
                        "steps": steps,
                        "data": result_data,
                        "charts": result_charts
                    }
            
            return {
                "answer": "最大反復回数に達しました",
                "steps": steps,
                "data": result_data,
                "charts": result_charts
            }
        
        else:  # provider == 'openai'
            # ===== OPENAI API IMPLEMENTATION (STREAMING) =====
            client = OpenAI(api_key=api_key)
            tools = build_openai_tools_schema()
            
            for iteration in range(10):
                msg_queue.put({"type": "thinking", "message": "次のアクションを考えています..."})
                
                # gpt-5の場合は専用パラメーターを使用
                api_params = {
                    "model": OPENAI_MODEL,
                    "messages": messages,
                    "tools": tools,
                    "tool_choice": "auto",
                    "stream": True
                }
                
                # gpt-5の場合は reasoning_effort を追加
                if OPENAI_MODEL.startswith("gpt-5"):
                    api_params["reasoning_effort"] = "high"  # 深い推論を実行
                
                # GPT-5でストリーミングエラーが発生する場合は非ストリーミングにフォールバック
                try:
                    response = client.chat.completions.create(**api_params)
                except Exception as e:
                    if "stream" in str(e).lower() and OPENAI_MODEL.startswith("gpt-5"):
                        # ストリーミングなしで再試行
                        msg_queue.put({"type": "info", "message": "GPT-5: 非ストリーミングモードで実行中..."})
                        api_params["stream"] = False
                        response_obj = client.chat.completions.create(**api_params)
                        
                        # 非ストリーミングレスポンスをストリーミング形式に変換
                        assistant_message = response_obj.choices[0].message
                        if assistant_message.content:
                            for char in assistant_message.content:
                                msg_queue.put({"type": "assistant_text", "text": char})
                        
                        # メッセージを構築してツール呼び出しを処理
                        messages.append(assistant_message.model_dump())
                        
                        if not assistant_message.tool_calls:
                            return {
                                "answer": assistant_message.content,
                                "steps": steps,
                                "data": result_data,
                                "charts": result_charts  # 配列で返す
                            }
                        
                        # ツール呼び出し処理（既存のロジックと同じ）
                        for tool_call in assistant_message.tool_calls:
                            func_name = tool_call.function.name
                            func_args = json.loads(tool_call.function.arguments)
                            
                            steps.append(f"🔧 {func_name}({json.dumps(func_args, ensure_ascii=False)})")
                            msg_queue.put({"type": "tool_call", "name": func_name, "args": func_args})
                            
                            try:
                                if func_name == "list_tables":
                                    msg_queue.put({"type": "tool_start", "tool": "list_tables", "message": "テーブル一覧を取得中..."})
                                    result = await list_tables(session, func_args.get("project", PROJECT_ID))
                                    func_result = json.dumps({"tables": result}, ensure_ascii=False)
                                    msg_queue.put({"type": "tool_done", "tool": "list_tables", "message": f"{len(result)}個のテーブルを発見しました"})
                                    
                                elif func_name == "describe_table":
                                    # GPT-5とGPT-4で異なるパラメータ名を使う可能性があるため両方をチェック
                                    table_name = func_args.get("table") or func_args.get("table_name")
                                    if not table_name:
                                        raise ValueError(f"Missing required parameter: 'table'. Received: {func_args}")
                                    msg_queue.put({"type": "tool_start", "tool": "describe_table", "message": f"テーブル '{table_name}' のスキーマを取得中..."})
                                    result = await describe_table(
                                        session,
                                        func_args.get("project", PROJECT_ID),
                                        func_args.get("dataset", DEFAULT_DATASET),
                                        table_name
                                    )
                                    func_result = json.dumps({"columns": result}, ensure_ascii=False)
                                    msg_queue.put({"type": "tool_done", "tool": "describe_table", "message": f"{len(result)}個のカラムを発見しました"})
                                    
                                elif func_name == "execute_query":
                                    query = func_args.get("query") or func_args.get("sql")
                                    if not query:
                                        raise ValueError(f"Missing required parameter: 'query'. Received: {func_args}")
                                    msg_queue.put({"type": "tool_start", "tool": "execute_query", "message": f"クエリを実行中...", "query": query})
                                    result = await execute_query(
                                        session,
                                        func_args.get("project", PROJECT_ID),
                                        func_args.get("location", LOCATION),
                                        query
                                    )
                                    result_data = result
                                    msg = f"📊 {len(result)}行のデータを取得しました"
                                    steps.append(msg)
                                    msg_queue.put({"type": "tool_done", "tool": "execute_query", "message": msg})
                                    
                                    # If data is large, return summary to avoid token limit
                                    if len(result) > 100:
                                        summary = {
                                            "row_count": len(result),
                                            "columns": list(result[0].keys()) if result else [],
                                            "sample_rows": result[:5],
                                            "message": f"Large dataset with {len(result)} rows. Data stored for analysis."
                                        }
                                        func_result = json.dumps(summary, ensure_ascii=False, cls=DateTimeEncoder)
                                    else:
                                        func_result = json.dumps({"result": result}, ensure_ascii=False, cls=DateTimeEncoder)
                                    
                                elif func_name == "suggest_chart":
                                    msg_queue.put({"type": "tool_start", "tool": "suggest_chart", "message": "グラフ設定を提案中..."})
                                    chart_config = {
                                        "chart_type": func_args.get("chart_type", "bar"),
                                        "x_axis": func_args.get("x_axis"),
                                        "y_axis": func_args.get("y_axis"),
                                        "title": func_args.get("title", "")
                                    }
                                    result_charts.append(chart_config)  # 配列に追加
                                    func_result = json.dumps({"chart": chart_config}, ensure_ascii=False)
                                    chart_type_ja = {
                                        "bar": "棒グラフ",
                                        "line": "折れ線グラフ",
                                        "pie": "円グラフ",
                                        "doughnut": "ドーナツグラフ",
                                        "scatter": "散布図",
                                        "none": "グラフなし"
                                    }.get(chart_config["chart_type"], chart_config["chart_type"])
                                    msg = f"📈 {chart_type_ja}を提案しました"
                                    steps.append(msg)
                                    msg_queue.put({"type": "tool_done", "tool": "suggest_chart", "message": msg})
                                    
                                else:
                                    func_result = json.dumps({"error": "Unknown function"})
                                    
                            except Exception as e:
                                error_msg = str(e)
                                error_details = {
                                    "error": error_msg,
                                    "function": func_name,
                                    "arguments": func_args
                                }
                                
                                # エラーの種類に応じた修正提案を追加
                                if func_name == "execute_query":
                                    if "column" in error_msg.lower() or "field" in error_msg.lower():
                                        error_details["hint"] = "Column name error detected. Use describe_table to get exact column names (case-sensitive)."
                                    elif "table" in error_msg.lower():
                                        error_details["hint"] = "Table reference error. Use backticks: `project.dataset.table` or `dataset.table`"
                                    elif "type" in error_msg.lower() or "cast" in error_msg.lower():
                                        error_details["hint"] = "Data type error. Use SAFE_CAST(column AS FLOAT64) for type conversion."
                                    elif "syntax" in error_msg.lower():
                                        error_details["hint"] = "SQL syntax error. Check BigQuery syntax: use backticks for tables, proper GROUP BY, etc."
                                    else:
                                        error_details["hint"] = "Query execution failed. Review the error message and check: 1) Column names (case-sensitive), 2) Table references (use backticks), 3) Data types, 4) BigQuery syntax."
                                
                                steps.append(f"❌ {func_name} エラー: {error_msg}")
                                msg_queue.put({"type": "error", "message": f"{func_name} でエラーが発生しました: {error_msg}"})
                                func_result = json.dumps(error_details, ensure_ascii=False)
                            
                            messages.append({
                                "role": "tool",
                                "tool_call_id": tool_call.id,
                                "content": func_result
                            })
                        continue  # 次のイテレーションへ
                    else:
                        raise  # その他のエラーは再スロー
            
            # ストリーミングレスポンスを収集
            full_content = ""
            reasoning_content = ""
            tool_calls_data = []
            current_tool_call = None
            
            for chunk in response:
                delta = chunk.choices[0].delta
                
                # 推論過程のストリーミング（GPT-5など）
                if hasattr(delta, 'reasoning_content') and delta.reasoning_content:
                    reasoning_content += delta.reasoning_content
                    msg_queue.put({"type": "reasoning", "text": delta.reasoning_content})
                
                # コンテンツのストリーミング
                if delta.content:
                    full_content += delta.content
                    msg_queue.put({"type": "assistant_text", "text": delta.content})
                
                # ツール呼び出しの処理
                if delta.tool_calls:
                    for tc_chunk in delta.tool_calls:
                        if tc_chunk.index is not None:
                            while len(tool_calls_data) <= tc_chunk.index:
                                tool_calls_data.append({"id": "", "type": "function", "function": {"name": "", "arguments": ""}})
                            current_tool_call = tool_calls_data[tc_chunk.index]
                            
                            if tc_chunk.id:
                                current_tool_call["id"] = tc_chunk.id
                            if tc_chunk.function:
                                if tc_chunk.function.name:
                                    current_tool_call["function"]["name"] = tc_chunk.function.name
                                if tc_chunk.function.arguments:
                                    current_tool_call["function"]["arguments"] += tc_chunk.function.arguments
            
            # メッセージを構築
            assistant_message_dict = {"role": "assistant", "content": full_content or None}
            
            if tool_calls_data:
                # ツール呼び出しをOpenAI形式に変換
                from openai.types.chat import ChatCompletionMessageToolCall
                from openai.types.chat.chat_completion_message_tool_call import Function
                
                tool_calls_objs = []
                for tc in tool_calls_data:
                    tool_calls_objs.append(
                        ChatCompletionMessageToolCall(
                            id=tc["id"],
                            type="function",
                            function=Function(
                                name=tc["function"]["name"],
                                arguments=tc["function"]["arguments"]
                            )
                        )
                    )
                assistant_message_dict["tool_calls"] = [tc.model_dump() for tc in tool_calls_objs]
                
                messages.append(assistant_message_dict)
                
                # ツール実行
                for tool_call in tool_calls_objs:
                    func_name = tool_call.function.name
                    func_args = json.loads(tool_call.function.arguments)
                    
                    step_msg = f"🔧 {func_name}({json.dumps(func_args, ensure_ascii=False)})"
                    steps.append(step_msg)
                    
                    try:
                        if func_name == "list_tables":
                            msg_queue.put({"type": "tool_start", "tool": "list_tables", "message": "テーブル一覧を取得中..."})
                            result = await list_tables(session, func_args.get("project", PROJECT_ID))
                            func_result = json.dumps({"tables": result}, ensure_ascii=False)
                            msg_queue.put({"type": "tool_done", "tool": "list_tables", "message": f"{len(result)}個のテーブルを発見しました"})
                            
                        elif func_name == "describe_table":
                            # GPT-5とGPT-4で異なるパラメータ名を使う可能性があるため両方をチェック
                            table_name = func_args.get("table") or func_args.get("table_name")
                            if not table_name:
                                raise ValueError(f"Missing required parameter: 'table'. Received: {func_args}")
                            msg_queue.put({"type": "tool_start", "tool": "describe_table", "message": f"テーブル '{table_name}' のスキーマを取得中..."})
                            result = await describe_table(
                                session,
                                func_args.get("project", PROJECT_ID),
                                func_args.get("dataset", DEFAULT_DATASET),
                                table_name
                            )
                            func_result = json.dumps({"columns": result}, ensure_ascii=False)
                            msg_queue.put({"type": "tool_done", "tool": "describe_table", "message": f"{len(result)}個のカラムを発見しました"})
                            
                        elif func_name == "execute_query":
                            query = func_args.get("query") or func_args.get("sql")
                            if not query:
                                raise ValueError(f"Missing required parameter: 'query'. Received: {func_args}")
                            msg_queue.put({"type": "tool_start", "tool": "execute_query", "message": f"クエリを実行中...", "query": query})
                            result = await execute_query(
                                session,
                                func_args.get("project", PROJECT_ID),
                                func_args.get("location", LOCATION),
                                query
                            )
                            result_data = result
                            msg = f"📊 {len(result)}行のデータを取得しました"
                            steps.append(msg)
                            msg_queue.put({"type": "tool_done", "tool": "execute_query", "message": msg})
                            
                            # If data is large, return summary to avoid token limit
                            if len(result) > 100:
                                summary = {
                                    "row_count": len(result),
                                    "columns": list(result[0].keys()) if result else [],
                                    "sample_rows": result[:5],
                                    "message": f"Large dataset with {len(result)} rows. Data stored for analysis."
                                }
                                func_result = json.dumps(summary, ensure_ascii=False, cls=DateTimeEncoder)
                            else:
                                func_result = json.dumps({"rows": result}, ensure_ascii=False, cls=DateTimeEncoder)
                            
                        elif func_name == "suggest_chart":
                            msg_queue.put({"type": "tool_start", "tool": "suggest_chart", "message": "グラフ設定を提案中..."})
                            chart_config = {
                                "chart_type": func_args.get("chart_type", "bar"),
                                "x_axis": func_args.get("x_axis"),
                                "y_axis": func_args.get("y_axis"),
                                "title": func_args.get("title", "")
                            }
                            result_charts.append(chart_config)  # 配列に追加
                            func_result = json.dumps({"chart": chart_config}, ensure_ascii=False)
                            chart_type_ja = {
                                "bar": "棒グラフ",
                                "line": "折れ線グラフ",
                                "pie": "円グラフ",
                                "doughnut": "ドーナツグラフ",
                                "scatter": "散布図",
                                "none": "グラフなし"
                            }.get(chart_config["chart_type"], chart_config["chart_type"])
                            msg = f"📈 {chart_type_ja}を提案しました"
                            steps.append(msg)
                            msg_queue.put({"type": "tool_done", "tool": "suggest_chart", "message": msg})
                            
                        else:
                            func_result = json.dumps({"error": "Unknown function"})
                            
                    except Exception as e:
                        error_msg = str(e)
                        error_details = {
                            "error": error_msg,
                            "function": func_name,
                            "arguments": func_args
                        }
                        
                        # エラーの種類に応じた修正提案を追加
                        if func_name == "execute_query":
                            if "column" in error_msg.lower() or "field" in error_msg.lower():
                                error_details["hint"] = "Column name error detected. Use describe_table to get exact column names (case-sensitive)."
                            elif "table" in error_msg.lower():
                                error_details["hint"] = "Table reference error. Use backticks: `project.dataset.table` or `dataset.table`"
                            elif "type" in error_msg.lower() or "cast" in error_msg.lower():
                                error_details["hint"] = "Data type error. Use SAFE_CAST(column AS FLOAT64) for type conversion."
                            elif "syntax" in error_msg.lower():
                                error_details["hint"] = "SQL syntax error. Check BigQuery syntax: use backticks for tables, proper GROUP BY, etc."
                            else:
                                error_details["hint"] = "Query execution failed. Review the error message and check: 1) Column names (case-sensitive), 2) Table references (use backticks), 3) Data types, 4) BigQuery syntax."
                        
                        steps.append(f"❌ {func_name} エラー: {error_msg}")
                        msg_queue.put({"type": "error", "message": f"{func_name} でエラーが発生しました: {error_msg}"})
                        func_result = json.dumps(error_details, ensure_ascii=False)
                    
                    messages.append({
                        "role": "tool",
                        "tool_call_id": tool_call.id,
                        "content": func_result
                    })
            else:
                # ツール呼び出しなし = 最終回答
                messages.append(assistant_message_dict)
                return {
                    "answer": full_content,
                    "reasoning": reasoning_content,
                    "steps": steps,
                    "data": result_data,
                    "charts": result_charts  # 配列で返す
                }
        
        return {
            "answer": "Maximum iterations reached.",
            "reasoning": "",
            "steps": steps,
            "data": result_data,
            "charts": result_charts  # 配列で返す
        }

async def run_agent_with_steps(task_id: str, user_question: str, conversation_history: List[Dict[str, str]], 
                              api_key: str = None, project_id: str = None, dataset_id: str = None, 
                              service_account_json: str = None, project_db_id: int = None, user_id: int = None,
                              provider: str = 'openai', gemini_api_key: str = None, use_adc: bool = False,
                              max_scan_gb: float = None, enable_review: bool = False,
                              gemini_model: str = 'gemini-2.5-pro') -> Dict[str, Any]:
    """Main agent logic with progress tracking via task_id - supports OpenAI and Gemini"""
    import time
    
    start_time = time.time()
    
    def add_step(message: str):
        """Helper to add step to task - now uses database"""
        add_chat_task_step(task_id, message)
    
    def add_reasoning(text: str):
        """Helper to add reasoning content to task - now uses database"""
        append_chat_task_reasoning(task_id, text)
    
    def is_cancelled():
        """Check if task has been cancelled - now uses database"""
        return is_chat_task_cancelled(task_id)
    
    api_key = api_key or OPENAI_API_KEY
    project_id = project_id or PROJECT_ID
    dataset_id = dataset_id or DEFAULT_DATASET
    
    # Debug logging for troubleshooting
    print(f"DEBUG run_agent_with_steps: use_adc={use_adc}, project_id={project_id}, dataset_id={dataset_id}")
    
    # Validate project_id early
    if not project_id or project_id == "your-project-id":
        error_msg = "BigQuery プロジェクトIDが設定されていません。設定画面でプロジェクトIDを設定してください。"
        add_step(f"❌ エラー: {error_msg}")
        return {"error": error_msg, "steps": [f"❌ エラー: {error_msg}"]}
    
    # Only use service_account_json if not using ADC
    if not use_adc:
        service_account_json = service_account_json or GCP_SA_JSON
    else:
        service_account_json = None
        print(f"DEBUG: Using ADC (Application Default Credentials) for BigQuery authentication")
    
    # Load project memories if project_db_id is provided
    project_memories = []
    if project_db_id and user_id:
        try:
            conn = get_db_connection()
            cur = conn.cursor(cursor_factory=RealDictCursor)
            cur.execute('''
                SELECT memory_key, memory_value, category, table_name, updated_at
                FROM project_memories
                WHERE project_id = %s AND user_id = %s
                ORDER BY updated_at DESC
            ''', (project_db_id, user_id))
            project_memories = cur.fetchall()
            cur.close()
            conn.close()
        except Exception as e:
            print(f"Warning: Failed to load project memories: {e}")
    
    # Load numeric-consistency reference data for the review AI
    baseline_stats = load_baseline_stats(project_db_id, user_id)
    past_metrics = load_past_metrics(project_db_id, user_id)
    
    env = os.environ.copy()
    # Only set GOOGLE_APPLICATION_CREDENTIALS if not using ADC
    if service_account_json and not use_adc:
        env["GOOGLE_APPLICATION_CREDENTIALS"] = service_account_json
    
    print(f"DEBUG: Starting MCP server with project={project_id}, location={LOCATION}, use_adc={use_adc}")
    
    server_params = StdioServerParameters(
        command="mcp-server-bigquery",
        args=["--project", project_id, "--location", LOCATION],
        env=env
    )
    
    steps = []
    result_data = None
    result_charts = []
    latest_dataframe = None  # Store latest query result as DataFrame for Python execution
    python_results = []  # Store Python execution results
    executed_queries = []  # Track executed SQL for the review step
    review_round = 0  # Number of completed review rounds
    review_info = None  # Last review verdict info

    def run_review_step(answer_text):
        """Run the reviewer AI and record progress steps. Returns (approved, feedback_text)."""
        nonlocal review_round, review_info
        review_round += 1
        add_step(f"🔎 レビューAIが回答を検証中...（{review_round}回目・通常1分以内）")
        review = run_review_ai(provider, api_key, gemini_api_key,
                               user_question, executed_queries, result_data, answer_text,
                               doc_context_section=build_doc_context_section(project_memories),
                               baseline_section=build_baseline_section(baseline_stats),
                               past_metrics_section=build_past_metrics_section(past_metrics))
        review_info = review
        if review.get("review_failed"):
            add_step(f"⚠️ レビューをスキップしました: {review.get('reason', '')}")
            return True, None
        if review["verdict"] == "pass":
            add_step(f"✅ レビュー合格: {review.get('reason') or '問題は見つかりませんでした'}")
            return True, None
        issues_text = "\n".join(f"- {i}" for i in review["issues"]) or "- （詳細なし）"
        add_step(f"⚠️ レビューで問題が指摘されました: {review.get('reason', '')}")
        for issue in review["issues"][:5]:
            add_step(f"　🔸 {issue}")
        if review_round >= MAX_REVIEW_ROUNDS:
            add_step("⚠️ 最大レビュー回数に達したため、現在の回答を返します")
            return True, None
        add_step("🔄 レビュー指摘に対応中：再分析を実行しています（追加で数分かかる場合があります）")
        feedback_text = f"""レビュアーAIがあなたの分析に以下の問題を指摘しました。指摘に対処し、必要ならSQLを修正・再実行した上で、修正済みの最終回答を日本語で提示してください。

指摘事項:
{issues_text}

判定理由: {review.get('reason', '')}"""
        return False, feedback_text

    async with AsyncExitStack() as stack:
        stdio_transport = await stack.enter_async_context(stdio_client(server_params))
        stdio, write = stdio_transport
        session = await stack.enter_async_context(ClientSession(stdio, write))
        
        await session.initialize()
        
        client = OpenAI(api_key=api_key)
        
        # Build memory section if memories exist
        memory_section = build_memory_section(project_memories)
        
        # Build dataset info section
        if dataset_id:
            dataset_info = f"- Default Dataset: {dataset_id}"
        else:
            dataset_info = "- Default Dataset: (Not set - use list_tables to discover available datasets and tables)"
        
        messages = [
            {"role": "system", "content": f"""You are an expert BigQuery data analyst assistant with deep knowledge of SQL optimization and data analysis.

## ENVIRONMENT
- BigQuery Project: {project_id}
{dataset_info}
{memory_section}
## STEP-BY-STEP REASONING FRAMEWORK
Follow this systematic approach for every user query:

### 1. UNDERSTAND THE QUESTION
- Identify the key metrics, dimensions, and filters requested
- Determine the time period if applicable
- Clarify any ambiguous terms by analyzing context

### 2. DISCOVER RELEVANT TABLES
- Use list_tables to see all available tables
- Identify tables that likely contain the requested data based on naming patterns
- Consider multiple tables if joins might be needed

### 3. EXAMINE TABLE SCHEMAS
- Use describe_table for each relevant table
- Carefully note:
  * Column names and their exact spelling
  * Data types (STRING, INTEGER, FLOAT, TIMESTAMP, DATE, BOOL, etc.)
  * Which columns can be used for aggregation vs grouping
  * Potential join keys between tables

### 4. CONSTRUCT SQL QUERY WITH BEST PRACTICES

**BigQuery-Specific Syntax:**
- Always use backticks for table references: `project.dataset.table` or `dataset.table`
- Use DATETIME/TIMESTAMP functions: CURRENT_DATETIME(), DATETIME_DIFF(), FORMAT_DATETIME()
- For dates: CURRENT_DATE(), DATE_DIFF(), FORMAT_DATE()
- Use SAFE_CAST() instead of CAST() to avoid errors
- Leverage IFNULL() or COALESCE() for NULL handling

**SQL Quality Guidelines:**
- Write explicit column names (avoid SELECT *)
- Use meaningful aliases with AS keyword
- Apply appropriate WHERE clauses for filtering
- Use proper GROUP BY with all non-aggregated columns
- Add ORDER BY for sorted results (essential for visualizations)
- Limit results appropriately (use LIMIT for large datasets)
- Handle NULLs explicitly to avoid incorrect aggregations

**Common Patterns:**
```sql
-- Aggregation with grouping
SELECT 
  column1,
  SUM(SAFE_CAST(column2 AS FLOAT64)) AS total_value,
  COUNT(*) AS count
FROM `{PROJECT_ID}.{DEFAULT_DATASET}.table_name`
WHERE column3 IS NOT NULL
GROUP BY column1
ORDER BY total_value DESC
LIMIT 100

-- Time-based analysis
SELECT 
  FORMAT_DATE('%Y-%m', date_column) AS month,
  AVG(numeric_column) AS avg_value
FROM `{PROJECT_ID}.{DEFAULT_DATASET}.table_name`
WHERE date_column >= DATE_SUB(CURRENT_DATE(), INTERVAL 6 MONTH)
GROUP BY month
ORDER BY month ASC

-- String matching
SELECT * 
FROM `{PROJECT_ID}.{DEFAULT_DATASET}.table_name`
WHERE LOWER(text_column) LIKE LOWER('%search_term%')
```

### 5. EXECUTE AND VALIDATE
- Run execute_query with the constructed SQL
- Check if results make sense (row count, value ranges)
- If error occurs, analyze the error message and fix:
  * Column name typos → recheck schema
  * Type mismatches → use SAFE_CAST()
  * NULL issues → add NULL handling
  * Syntax errors → verify BigQuery syntax

### 6. ADVANCED ANALYSIS WITH PYTHON (Optional)
For complex analysis beyond SQL capabilities, use execute_python with query results:

**When to Use Python:**
- Statistical analysis (correlation, regression, hypothesis testing)
- Custom data transformations or complex calculations
- Advanced visualizations (histograms, heatmaps, custom plots)
- Machine learning or predictive analytics
- Data cleaning or preprocessing

**Python Environment:**
- Query results available as 'df' (pandas DataFrame)
- Available libraries: pandas (pd), numpy (np), matplotlib.pyplot (plt), seaborn (sns)
- Use print() for text output, set 'result' variable for return values
- Use plt for custom visualizations (automatically captured)

**Example Usage:**
```python
# Statistical analysis
print(df.describe())
print(df.corr())

# Custom visualization
plt.figure(figsize=(10, 6))
plt.hist(df['column_name'], bins=30)
plt.title('Distribution of Column')
plt.xlabel('Value')
plt.ylabel('Frequency')

# Advanced calculation
result = df.groupby('category')['value'].agg(['mean', 'std', 'count'])
```

### 7. VISUALIZE INTELLIGENTLY
After getting query results, choose between suggest_chart (simple) or execute_python (advanced):

**suggest_chart - Use for standard visualizations:**
- **Line chart**: Time series data (x=date/time, y=metric), trends over time
- **Bar chart**: Categorical comparisons (x=category, y=value), rankings
- **Pie/Doughnut**: Part-to-whole with ≤10 categories, percentages, distributions
- **Scatter**: Correlation between two numeric variables (x=variable1, y=variable2)
- **None**: Single values, text-heavy data, or >50 data points without clear pattern

**execute_python - Use for custom visualizations:**
- Histograms, box plots, heatmaps, pair plots
- Multiple subplots or complex chart combinations
- Statistical plots (QQ plots, distribution fits)
- Custom styling or annotations

**Axis Selection (for suggest_chart):**
- x_axis: Use the dimension/category column (dates for trends, categories for comparisons)
- y_axis: Use the numeric metric column (always numeric values)
- Ensure columns match exactly with query result column names

### 8. PROVIDE INSIGHTS
- Explain what the data shows
- Highlight key findings (trends, outliers, patterns)
- Answer the user's original question clearly

## CRITICAL RULES
1. **Always check schemas before writing SQL** - Never assume column names
2. **Use exact column names from describe_table** - BigQuery is case-sensitive
3. **Test one table at a time** - Don't assume relationships exist
4. **Handle NULLs explicitly** - They affect aggregations
5. **Use appropriate data types** - SAFE_CAST when unsure
6. **Fully qualify table names** - Use `project.dataset.table` or `dataset.table`
7. **Always provide Japanese explanations** - Users expect Japanese responses

Think step-by-step and show your reasoning process."""}
        ]
        
        messages.extend(conversation_history)
        messages.append({"role": "user", "content": user_question})
        
        # Provider-based implementation
        if provider == 'gemini':
            # ===== GEMINI API IMPLEMENTATION (WITH TASK TRACKING) =====
            client = genai_new.Client(api_key=gemini_api_key)
            gemini_tools = build_gemini_tools_schema()
            gemini_gen_config = genai_types.GenerateContentConfig(
                tools=[genai_types.Tool(function_declarations=gemini_tools)],
                system_instruction=messages[0]["content"]
            )
            gemini_model_name = gemini_model or 'gemini-2.5-pro'
            
            # Convert conversation history to Gemini format
            gemini_contents = []
            for msg in conversation_history:
                if msg["role"] == "user":
                    gemini_contents.append({"role": "user", "parts": [{"text": msg["content"]}]})
                elif msg["role"] == "assistant":
                    gemini_contents.append({"role": "model", "parts": [{"text": msg["content"]}]})
            
            # Add current question
            gemini_contents.append({"role": "user", "parts": [{"text": user_question}]})
            
            # Main iteration loop
            max_iterations = 20 if enable_review else 10
            for iteration in range(max_iterations):
                # Check for cancellation
                if is_cancelled():
                    return {
                        "answer": "処理がキャンセルされました",
                        "data": result_data,
                        "charts": result_charts,
                        "python_results": python_results,
                        "cancelled": True
                    }
                
                add_step("🤔 次のアクションを考えています...")
                
                try:
                    # Generate response (new google-genai SDK; preserves thought signatures for Gemini 3+)
                    response = client.models.generate_content(
                        model=gemini_model_name,
                        contents=gemini_contents,
                        config=gemini_gen_config
                    )
                    
                    if (not response.candidates or not response.candidates[0].content
                            or not response.candidates[0].content.parts):
                        add_step("❌ Geminiからレスポンスがありませんでした")
                        final_result = {
                            "answer": "Geminiからレスポンスがありませんでした",
                            "data": result_data,
                            "charts": result_charts,
                            "python_results": python_results
                        }
                        set_final_result(final_result)
                        return final_result
                    
                    response_parts = response.candidates[0].content.parts
                    
                    # Check for function calls
                    has_function_call = False
                    function_responses = []
                    for part in response_parts:
                        if getattr(part, 'function_call', None):
                            has_function_call = True
                            function_call = part.function_call
                            func_name = function_call.name
                            func_args = dict(function_call.args) if function_call.args else {}
                            
                            try:
                                if func_name == "list_tables":
                                    add_step("📋 テーブル一覧を取得中...")
                                    result = await list_tables(session, func_args.get("project", project_id))
                                    func_result = {"tables": result}
                                    add_step(f"✅ {len(result)}個のテーブルを発見しました")
                                    
                                elif func_name == "describe_table":
                                    table_name = func_args.get("table") or func_args.get("table_name")
                                    if not table_name:
                                        raise ValueError(f"Missing required parameter: 'table'")
                                    add_step(f"🔍 テーブル '{table_name}' のスキーマを取得中...")
                                    result = await describe_table(
                                        session,
                                        func_args.get("project", project_id),
                                        func_args.get("dataset", dataset_id),
                                        table_name
                                    )
                                    func_result = {"columns": result}
                                    add_step(f"✅ {len(result)}個のカラムを発見しました")
                                    
                                elif func_name == "execute_query":
                                    query = func_args.get("query") or func_args.get("sql")
                                    if not query:
                                        raise ValueError(f"Missing required parameter: 'query'")
                                    add_step("⚡ クエリを実行中...")
                                    result = await execute_query(
                                        session,
                                        func_args.get("project", project_id),
                                        func_args.get("location", LOCATION),
                                        query,
                                        credentials_path=service_account_json,
                                        max_scan_gb=max_scan_gb,
                                        add_step=add_step
                                    )
                                    result_data = result
                                    executed_queries.append(query)
                                    
                                    if result:
                                        latest_dataframe = pd.DataFrame(result)
                                        add_step(f"📊 {len(result)}行のデータを取得しました（Python処理可能）")
                                        
                                        if len(result) > 100:
                                            func_result = {
                                                "row_count": len(result),
                                                "columns": list(result[0].keys()),
                                                "sample_rows": result[:5],
                                                "message": f"Large dataset with {len(result)} rows. Full data available in Python as 'df'."
                                            }
                                        else:
                                            func_result = {"rows": result}
                                    else:
                                        latest_dataframe = None
                                        add_step("📊 0行のデータを取得しました")
                                        func_result = {"rows": []}
                                        
                                elif func_name == "execute_python":
                                    code = func_args.get("code")
                                    if not code:
                                        raise ValueError(f"Missing required parameter: 'code'")
                                    
                                    add_step("🐍 Pythonコードを実行中...")
                                    
                                    df_dict = {}
                                    if latest_dataframe is not None:
                                        df_dict['df'] = latest_dataframe
                                    
                                    python_result = execute_python_code(code, df_dict, timeout=30)
                                    python_results.append(python_result)
                                    
                                    func_result = {
                                        "output": python_result.get('output', ''),
                                        "result": python_result.get('result'),
                                        "plots_count": len(python_result.get('plots', []))
                                    }
                                    
                                    result_parts = []
                                    if python_result.get('output'):
                                        result_parts.append(f"出力あり")
                                    if python_result.get('result'):
                                        result_parts.append(f"結果あり")
                                    if python_result.get('plots'):
                                        result_parts.append(f"{len(python_result['plots'])}個のグラフ")
                                    add_step(f"✅ Python実行完了: {', '.join(result_parts) if result_parts else 'エラーなし'}")
                                    
                                elif func_name == "suggest_chart":
                                    add_step("📈 グラフ設定を提案中...")
                                    chart_config = {
                                        "chart_type": func_args.get("chart_type", "bar"),
                                        "x_axis": func_args.get("x_axis"),
                                        "y_axis": func_args.get("y_axis"),
                                        "title": func_args.get("title", "")
                                    }
                                    result_charts.append(chart_config)
                                    func_result = {"chart": chart_config}
                                    chart_type_ja = {
                                        "bar": "棒グラフ",
                                        "line": "折れ線グラフ",
                                        "pie": "円グラフ",
                                        "doughnut": "ドーナツグラフ",
                                        "scatter": "散布図",
                                        "none": "グラフなし"
                                    }.get(chart_config["chart_type"], chart_config["chart_type"])
                                    add_step(f"✅ {chart_type_ja}を提案しました")
                                    
                                else:
                                    func_result = {"error": "Unknown function"}
                                    add_step(f"❌ 未知の関数: {func_name}")
                                    
                            except Exception as e:
                                error_msg = str(e)
                                add_step(f"❌ {func_name} エラー: {error_msg}")
                                func_result = {
                                    "error": error_msg,
                                    "function": func_name,
                                    "arguments": func_args
                                }
                            
                            # Collect the response; sent back after all parts are processed
                            function_responses.append((func_name, func_result))
                    
                    if has_function_call:
                        # Append the model turn as-is (preserves thought signatures required by Gemini 3+),
                        # then a single user turn with all function responses
                        gemini_contents.append(response.candidates[0].content)
                        gemini_contents.append(genai_types.Content(
                            role="user",
                            parts=[
                                genai_types.Part.from_function_response(name=n, response=r)
                                for n, r in function_responses
                            ]
                        ))
                    
                    if not has_function_call:
                        # Final answer
                        answer_text = ""
                        for part in response_parts:
                            if getattr(part, 'text', None):
                                answer_text += part.text
                        
                        # Optional AI review of the analysis result
                        if enable_review and review_round < MAX_REVIEW_ROUNDS:
                            approved, feedback = run_review_step(answer_text)
                            if not approved:
                                gemini_contents.append({"role": "model", "parts": [{"text": answer_text}]})
                                gemini_contents.append({"role": "user", "parts": [{"text": feedback}]})
                                continue
                        
                        add_step("✅ 完了")
                        final_result = {
                            "answer": answer_text,
                            "data": result_data,
                            "charts": result_charts,
                            "python_results": python_results,
                            "review": review_info
                        }
                        return final_result
                        
                except Exception as e:
                    error_msg = f"Gemini API error: {str(e)}"
                    print(f"ERROR: {error_msg}")
                    add_step(f"❌ {error_msg}")
                    final_result = {
                        "answer": f"エラーが発生しました: {error_msg}",
                        "data": result_data,
                        "charts": result_charts,
                        "python_results": python_results
                    }
                    return final_result
            
            # Max iterations reached
            add_step("⚠️ 最大反復回数に達しました")
            final_result = {
                "answer": "最大反復回数に達しました",
                "data": result_data,
                "charts": result_charts,
                "python_results": python_results,
                "review": review_info
            }
            return final_result
        
        else:  # provider == 'openai'
            # ===== OPENAI API IMPLEMENTATION (WITH TASK TRACKING) =====
            client = OpenAI(api_key=api_key)
            tools = build_openai_tools_schema()
            
            max_iterations = 20 if enable_review else 10
            for iteration in range(max_iterations):
                # Check for cancellation
                if is_cancelled():
                    return {
                        "answer": "処理がキャンセルされました",
                        "data": result_data,
                        "charts": result_charts,
                        "python_results": python_results,
                        "cancelled": True
                    }
                
                add_step("🤔 次のアクションを考えています...")
                
                api_params = {
                    "model": OPENAI_MODEL,
                    "messages": messages,
                    "tools": tools,
                    "tool_choice": "auto"
                }
                
                if OPENAI_MODEL.startswith("gpt-5"):
                    api_params["reasoning_effort"] = "high"
                
                response = client.chat.completions.create(**api_params)
                assistant_message = response.choices[0].message
                
                # GPT-5の推論過程を取得してリアルタイム保存
                if hasattr(assistant_message, 'reasoning_content') and assistant_message.reasoning_content:
                    add_reasoning(assistant_message.reasoning_content)
                
                assistant_message_dict = {
                    "role": "assistant", 
                    "content": assistant_message.content or None
                }
                
                if assistant_message.tool_calls:
                    assistant_message_dict["tool_calls"] = [tc.model_dump() for tc in assistant_message.tool_calls]
                    messages.append(assistant_message_dict)
                    
                    for tool_call in assistant_message.tool_calls:
                        func_name = tool_call.function.name
                        func_args = json.loads(tool_call.function.arguments)
                        
                        try:
                            if func_name == "list_tables":
                                add_step("📋 テーブル一覧を取得中...")
                                result = await list_tables(session, func_args.get("project", PROJECT_ID))
                                func_result = json.dumps({"tables": result}, ensure_ascii=False)
                                add_step(f"✅ {len(result)}個のテーブルを発見しました")
                                
                            elif func_name == "describe_table":
                                table_name = func_args.get("table") or func_args.get("table_name")
                                if not table_name:
                                    raise ValueError(f"Missing required parameter: 'table'")
                                add_step(f"🔍 テーブル '{table_name}' のスキーマを取得中...")
                                result = await describe_table(
                                    session,
                                    func_args.get("project", PROJECT_ID),
                                    func_args.get("dataset", DEFAULT_DATASET),
                                    table_name
                                )
                                func_result = json.dumps({"columns": result}, ensure_ascii=False)
                                add_step(f"✅ {len(result)}個のカラムを発見しました")
                                
                            elif func_name == "execute_query":
                                query = func_args.get("query") or func_args.get("sql")
                                if not query:
                                    raise ValueError(f"Missing required parameter: 'query'")
                                add_step("⚡ クエリを実行中...")
                                result = await execute_query(
                                    session,
                                    func_args.get("project", PROJECT_ID),
                                    func_args.get("location", LOCATION),
                                    query,
                                    credentials_path=service_account_json,
                                    max_scan_gb=max_scan_gb,
                                    add_step=add_step
                                )
                                result_data = result
                                executed_queries.append(query)
                                
                                # Convert to DataFrame for Python execution
                                if result:
                                    latest_dataframe = pd.DataFrame(result)
                                    add_step(f"📊 {len(result)}行のデータを取得しました（Python処理可能）")
                                    
                                    # If data is large, return summary to avoid token limit
                                    if len(result) > 100:
                                        summary = {
                                            "row_count": len(result),
                                            "columns": list(result[0].keys()) if result else [],
                                            "sample_rows": result[:5],
                                            "message": f"Large dataset with {len(result)} rows. Full data available in Python as 'df'. Use execute_python for analysis."
                                        }
                                        func_result = json.dumps(summary, ensure_ascii=False, cls=DateTimeEncoder)
                                    else:
                                        func_result = json.dumps({"rows": result}, ensure_ascii=False, cls=DateTimeEncoder)
                                else:
                                    latest_dataframe = None
                                    add_step("📊 0行のデータを取得しました")
                                    func_result = json.dumps({"rows": []}, ensure_ascii=False)
                                
                            elif func_name == "execute_python":
                                code = func_args.get("code")
                                if not code:
                                    raise ValueError(f"Missing required parameter: 'code'")
                                
                                add_step("🐍 Pythonコードを実行中...")
                                
                                # Prepare DataFrame dict
                                df_dict = {}
                                if latest_dataframe is not None:
                                    df_dict['df'] = latest_dataframe
                                
                                # Execute Python code
                                python_result = execute_python_code(code, df_dict, timeout=30)
                                
                                # Store result
                                python_results.append(python_result)
                                
                                # Build result message
                                result_parts = []
                                if python_result.get('output'):
                                    result_parts.append(f"出力:\n{python_result['output']}")
                                if python_result.get('result'):
                                    result_parts.append(f"結果: {python_result['result']}")
                                if python_result.get('plots'):
                                    result_parts.append(f"{len(python_result['plots'])}個のグラフを生成しました")
                                
                                func_result = json.dumps({
                                    "output": python_result.get('output', ''),
                                    "result": python_result.get('result'),
                                    "plots_count": len(python_result.get('plots', []))
                                }, ensure_ascii=False)
                                
                                add_step(f"✅ Python実行完了: {', '.join(result_parts) if result_parts else 'エラーなし'}")
                                
                            elif func_name == "suggest_chart":
                                add_step("📈 グラフ設定を提案中...")
                                chart_config = {
                                    "chart_type": func_args.get("chart_type", "bar"),
                                    "x_axis": func_args.get("x_axis"),
                                    "y_axis": func_args.get("y_axis"),
                                    "title": func_args.get("title", "")
                                }
                                result_charts.append(chart_config)
                                func_result = json.dumps({"chart": chart_config}, ensure_ascii=False)
                                chart_type_ja = {
                                    "bar": "棒グラフ",
                                    "line": "折れ線グラフ",
                                    "pie": "円グラフ",
                                    "doughnut": "ドーナツグラフ",
                                    "scatter": "散布図",
                                    "none": "グラフなし"
                                }.get(chart_config["chart_type"], chart_config["chart_type"])
                                add_step(f"✅ {chart_type_ja}を提案しました")
                                
                            else:
                                func_result = json.dumps({"error": "Unknown function"})
                                
                        except Exception as e:
                            error_msg = str(e)
                            error_details = {
                                "error": error_msg,
                                "function": func_name,
                                "arguments": func_args
                            }
                            
                            if func_name == "execute_query":
                                if "column" in error_msg.lower() or "field" in error_msg.lower():
                                    error_details["hint"] = "Column name error detected. Use describe_table to get exact column names (case-sensitive)."
                                elif "table" in error_msg.lower():
                                    error_details["hint"] = "Table reference error. Use backticks: `project.dataset.table` or `dataset.table`"
                                elif "type" in error_msg.lower() or "cast" in error_msg.lower():
                                    error_details["hint"] = "Data type error. Use SAFE_CAST(column AS FLOAT64) for type conversion."
                                elif "syntax" in error_msg.lower():
                                    error_details["hint"] = "SQL syntax error. Check BigQuery syntax: use backticks for tables, proper GROUP BY, etc."
                                else:
                                    error_details["hint"] = "Query execution failed. Review the error message and check: 1) Column names (case-sensitive), 2) Table references (use backticks), 3) Data types, 4) BigQuery syntax."
                            
                            add_step(f"❌ {func_name} エラー: {error_msg}")
                            func_result = json.dumps(error_details, ensure_ascii=False)
                        
                        messages.append({
                            "role": "tool",
                            "tool_call_id": tool_call.id,
                            "content": func_result
                        })
                else:
                    messages.append(assistant_message_dict)
                    
                    # Optional AI review of the analysis result
                    if enable_review and review_round < MAX_REVIEW_ROUNDS:
                        approved, feedback = run_review_step(assistant_message.content or "")
                        if not approved:
                            messages.append({"role": "user", "content": feedback})
                            continue
                    
                    add_step("✅ 分析完了")
                    processing_time = time.time() - start_time
                    
                    # GPT-5の場合、推論過程を取得
                    reasoning_content = ""
                    if hasattr(assistant_message, 'reasoning_content') and assistant_message.reasoning_content:
                        reasoning_content = assistant_message.reasoning_content
                        add_reasoning(reasoning_content)  # リアルタイムで chat_tasks に保存
                    
                    return {
                        "answer": assistant_message.content,
                        "reasoning": reasoning_content,
                        "steps": steps,
                        "data": result_data,
                        "charts": result_charts,
                        "python_results": python_results,
                        "steps_count": len(steps),
                        "processing_time": processing_time,
                        "review": review_info
                    }
            
            add_step("⚠️ 最大反復回数に達しました")
            processing_time = time.time() - start_time
            return {
                "answer": "Maximum iterations reached.",
                "reasoning": "",
                "steps": steps,
                "data": result_data,
                "charts": result_charts,
                "python_results": python_results,
                "steps_count": len(steps),
                "processing_time": processing_time,
                "review": review_info
            }

@app.route('/health')
def health_check():
    """Health check endpoint for deployment"""
    return jsonify({"status": "ok"}), 200

@app.route('/')
@login_required
def index():
    """Render dashboard page"""
    # Check if user has any projects
    conn = get_db_connection()
    cur = conn.cursor(cursor_factory=RealDictCursor)
    cur.execute('SELECT COUNT(*) as count FROM projects WHERE user_id = %s', (current_user.id,))
    result = cur.fetchone()
    cur.close()
    conn.close()
    
    # If no projects, redirect to project management page
    if result['count'] == 0:
        flash('まずプロジェクトを作成してください。', 'info')
        return redirect(url_for('projects'))
    
    return render_template('dashboard.html')

@app.route('/login', methods=['GET', 'POST'])
def login():
    """Handle login"""
    if current_user.is_authenticated:
        return redirect(url_for('index'))
    
    if request.method == 'POST':
        username = request.form.get('username')
        password = request.form.get('password')
        
        conn = get_db_connection()
        cur = conn.cursor(cursor_factory=RealDictCursor)
        cur.execute('SELECT * FROM users WHERE username = %s', (username,))
        user_data = cur.fetchone()
        cur.close()
        conn.close()
        
        if user_data and bcrypt.checkpw(password.encode('utf-8'), user_data['password_hash'].encode('utf-8')):
            user = User(user_data['id'], user_data['username'], user_data['email'])
            login_user(user, remember=request.form.get('remember'))
            
            # Safe redirect - only allow same-origin redirects
            next_page = get_redirect_target()
            if not next_page or not is_safe_url(next_page):
                next_page = url_for('index')
            return redirect(next_page)
        else:
            flash('ユーザー名またはパスワードが正しくありません。', 'error')
    
    return render_template('login.html')

@app.route('/signup', methods=['GET', 'POST'])
def signup():
    """Handle user registration"""
    if current_user.is_authenticated:
        return redirect(url_for('index'))
    
    if request.method == 'POST':
        username = request.form.get('username')
        password = request.form.get('password')
        email = request.form.get('email', '').strip() or None  # 空文字列はNULLに変換
        
        # Validation
        if not username or not password:
            flash('ユーザー名とパスワードは必須です。', 'error')
            return render_template('signup.html')
        
        if len(password) < 6:
            flash('パスワードは6文字以上にしてください。', 'error')
            return render_template('signup.html')
        
        # Check if username exists
        conn = get_db_connection()
        cur = conn.cursor(cursor_factory=RealDictCursor)
        cur.execute('SELECT id FROM users WHERE username = %s', (username,))
        existing_user = cur.fetchone()
        
        if existing_user:
            flash('このユーザー名は既に使用されています。', 'error')
            cur.close()
            conn.close()
            return render_template('signup.html')
        
        # Check if email exists
        if email:
            cur.execute('SELECT id FROM users WHERE email = %s', (email,))
            existing_email = cur.fetchone()
            
            if existing_email:
                flash('このメールアドレスは既に使用されています。', 'error')
                cur.close()
                conn.close()
                return render_template('signup.html')
        
        # Hash password
        password_hash = bcrypt.hashpw(password.encode('utf-8'), bcrypt.gensalt()).decode('utf-8')
        
        # Insert new user
        try:
            cur.execute(
                'INSERT INTO users (username, password_hash, email) VALUES (%s, %s, %s) RETURNING id',
                (username, password_hash, email)
            )
            user_id = cur.fetchone()['id']
            conn.commit()
        except Exception as e:
            conn.rollback()
            if 'users_email_key' in str(e):
                flash('このメールアドレスは既に使用されています。', 'error')
            elif 'users_username_key' in str(e):
                flash('このユーザー名は既に使用されています。', 'error')
            else:
                flash('アカウント作成中にエラーが発生しました。', 'error')
            cur.close()
            conn.close()
            return render_template('signup.html')
        
        cur.close()
        conn.close()
        
        flash('アカウントが作成されました。ログインしてください。', 'success')
        return redirect(url_for('login'))
    
    return render_template('signup.html')

@app.route('/logout')
@login_required
def logout():
    """Handle logout"""
    logout_user()
    return redirect(url_for('login'))

@app.route('/account')
@login_required
def account_settings():
    """Render account settings page"""
    return render_template('account_settings.html')

@app.route('/agent-chat')
@login_required
def agent_chat_redirect():
    """Redirect to latest chat session or create new one"""
    try:
        # Check if user has any projects first
        conn = get_db_connection()
        cur = conn.cursor(cursor_factory=RealDictCursor)
        
        cur.execute('SELECT id FROM projects WHERE user_id = %s LIMIT 1', (current_user.id,))
        has_any_project = cur.fetchone() is not None
        
        if not has_any_project:
            cur.close()
            conn.close()
            flash('まずプロジェクトを作成してください。', 'info')
            return redirect(url_for('projects'))
        
        # Get active project
        cur.execute('''
            SELECT id FROM projects
            WHERE user_id = %s AND is_active = true
            LIMIT 1
        ''', (current_user.id,))
        project = cur.fetchone()
        
        if not project:
            # Has projects but no active one - show page with warning to select project
            cur.close()
            conn.close()
            return render_template('agent_chat.html', session_id=None)
        
        # Get latest session for this project
        cur.execute('''
            SELECT id FROM chat_sessions
            WHERE project_id = %s
            ORDER BY updated_at DESC
            LIMIT 1
        ''', (project['id'],))
        latest_session = cur.fetchone()
        
        if latest_session:
            # Redirect to latest session
            session_id = latest_session['id']
            cur.close()
            conn.close()
            return redirect(url_for('agent_chat', session_id=session_id))
        
        # No existing sessions - create new one
        cur.execute('''
            INSERT INTO chat_sessions (project_id, title)
            VALUES (%s, %s)
            RETURNING id
        ''', (project['id'], 'New Chat'))
        
        session = cur.fetchone()
        conn.commit()
        cur.close()
        conn.close()
        
        # Redirect to session URL
        return redirect(url_for('agent_chat', session_id=session['id']))
    except Exception as e:
        print(f"Error handling chat redirect: {e}")
        return render_template('agent_chat.html', session_id=None)

@app.route('/agent-chat/<int:session_id>')
@login_required
def agent_chat(session_id):
    """Render agent chat page for specific session"""
    try:
        # Verify session belongs to user
        conn = get_db_connection()
        cur = conn.cursor(cursor_factory=RealDictCursor)
        
        cur.execute('''
            SELECT s.id 
            FROM chat_sessions s
            JOIN projects p ON p.id = s.project_id
            WHERE s.id = %s AND p.user_id = %s
            LIMIT 1
        ''', (session_id, current_user.id))
        
        session = cur.fetchone()
        cur.close()
        conn.close()
        
        if not session:
            # Session not found or doesn't belong to user
            flash('セッションが見つかりませんでした。', 'error')
            return redirect(url_for('agent_chat_redirect'))
        
        return render_template('agent_chat.html', session_id=session_id)
    except Exception as e:
        print(f"Error loading session: {e}")
        return redirect(url_for('agent_chat_redirect'))

@app.route('/api/chat', methods=['POST'])
@login_required
def chat():
    """Start chat task and return task_id for polling"""
    try:
        data = request.json
        question = data.get('question', '')
        history = data.get('history', [])
        session_id = data.get('session_id')
        
        if not question:
            return jsonify({"error": "Question is required"}), 400
        
        # Get user_id from current_user before background thread
        user_id = current_user.id
        
        # Get active project configuration
        config = get_active_project_config(user_id)
        
        # Use provider from request if provided, otherwise use project config
        provider = data.get('provider', config.get('provider', 'openai'))
        
        # Validate API key based on provider
        if provider == 'openai' and not config['api_key']:
            return jsonify({"error": "OpenAI API key not configured"}), 500
        elif provider == 'gemini' and not config.get('gemini_api_key'):
            return jsonify({"error": "Gemini API key not configured"}), 500
        
        if not config['project_id']:
            return jsonify({"error": "BigQuery project not configured"}), 500
        
        # Get active project ID
        conn = get_db_connection()
        cur = conn.cursor(cursor_factory=RealDictCursor)
        cur.execute('SELECT id FROM projects WHERE user_id = %s AND is_active = true LIMIT 1', (user_id,))
        active_project = cur.fetchone()
        active_project_id = active_project['id'] if active_project else None
        
        # Create or verify session
        if session_id:
            cur.execute('''
                SELECT id FROM chat_sessions 
                WHERE id = %s AND project_id = %s
            ''', (session_id, active_project_id))
            if not cur.fetchone():
                session_id = None
        
        if not session_id and active_project_id:
            cur.execute('''
                INSERT INTO chat_sessions (project_id, title)
                VALUES (%s, %s)
                RETURNING id
            ''', (active_project_id, 'New Chat'))
            new_session = cur.fetchone()
            session_id = new_session['id']
            conn.commit()
        
        cur.close()
        conn.close()
        
        # Generate task ID
        import uuid
        task_id = str(uuid.uuid4())
        
        # Initialize task state in database (Cloud Run compatible)
        if not create_chat_task(task_id, user_id, active_project_id):
            return jsonify({"error": "Failed to create task"}), 500
        add_chat_task_step(task_id, '🔧 処理を開始しました...')
        
        # Run task in background thread
        def run_task():
            try:
                result = asyncio.run(run_agent_with_steps(
                    task_id, question, history,
                    api_key=config['api_key'],
                    project_id=config['project_id'],
                    dataset_id=config['dataset_id'],
                    service_account_json=config['service_account_json'],
                    project_db_id=active_project_id,
                    user_id=user_id,
                    provider=provider,
                    gemini_api_key=config.get('gemini_api_key'),
                    use_adc=config.get('use_adc', False),
                    max_scan_gb=config.get('max_scan_gb'),
                    enable_review=config.get('enable_review', False),
                    gemini_model=config.get('gemini_model', 'gemini-2.5-pro')
                ))
                
                # Save to database
                if active_project_id and session_id:
                    conn = get_db_connection()
                    cur = conn.cursor(cursor_factory=RealDictCursor)
                    cur.execute('''
                        INSERT INTO chat_history (project_id, session_id, user_message, ai_response, query_result, steps_count, processing_time, reasoning_process)
                        VALUES (%s, %s, %s, %s, %s, %s, %s, %s)
                    ''', (
                        active_project_id,
                        session_id,
                        question,
                        result.get('answer', ''),
                        json.dumps(result, ensure_ascii=False, cls=DateTimeEncoder),
                        result.get('steps_count', 0),
                        result.get('processing_time', 0),
                        result.get('reasoning', '')
                    ))
                    
                    # Update session timestamp and title
                    cur.execute('''
                        UPDATE chat_sessions 
                        SET updated_at = CURRENT_TIMESTAMP 
                        WHERE id = %s
                    ''', (session_id,))
                    
                    cur.execute('''
                        SELECT title, 
                               (SELECT COUNT(*) FROM chat_history WHERE session_id = %s) as msg_count
                        FROM chat_sessions 
                        WHERE id = %s
                    ''', (session_id, session_id))
                    session_info = cur.fetchone()
                    
                    if session_info and session_info['title'] == 'New Chat' and session_info['msg_count'] == 1:
                        auto_title = question[:50] + ('...' if len(question) > 50 else '')
                        cur.execute('''
                            UPDATE chat_sessions 
                            SET title = %s 
                            WHERE id = %s
                        ''', (auto_title, session_id))
                    
                    conn.commit()
                    cur.close()
                    conn.close()
                
                # Extract and store key metrics from the answer for future consistency checks
                if active_project_id and not result.get('error'):
                    try:
                        save_analysis_metrics(
                            provider, config['api_key'], config.get('gemini_api_key'),
                            config.get('gemini_model'), active_project_id, user_id,
                            session_id, question, result.get('answer', '')
                        )
                    except Exception as metric_err:
                        print(f"Warning: metric extraction failed: {metric_err}")
                
                # Update task status in database (Cloud Run compatible)
                update_chat_task(task_id, 
                    status='completed',
                    result=result,
                    session_id=session_id,
                    reasoning=result.get('reasoning', '')
                )
                    
            except Exception as e:
                import traceback
                error_traceback = traceback.format_exc()
                print(f"ERROR in run_task: {error_traceback}", flush=True)
                # Update task error status in database (Cloud Run compatible)
                add_chat_task_step(task_id, f"❌ エラー: {str(e)}")
                update_chat_task(task_id,
                    status='error',
                    error=str(e),
                    traceback=error_traceback
                )
        
        thread = threading.Thread(target=run_task)
        thread.daemon = True
        thread.start()
        
        return jsonify({
            "task_id": task_id,
            "session_id": session_id
        })
    
    except Exception as e:
        import traceback
        return jsonify({
            "error": str(e),
            "traceback": traceback.format_exc()
        }), 500

@app.route('/api/chat/status/<task_id>', methods=['GET'])
@login_required
def chat_status(task_id):
    """Get chat task status and progress - now uses database for Cloud Run compatibility"""
    task = get_chat_task(task_id)
    if not task:
        return jsonify({"error": "Task not found"}), 404
    
    return jsonify({
        "status": task['status'],
        "steps": task['steps'] if task['steps'] else [],
        "result": task['result'],
        "error": task['error'],
        "session_id": task.get('session_id'),
        "reasoning": task.get('reasoning', ''),
        "cancelled": task.get('cancelled', False)
    })

@app.route('/api/chat/cancel/<task_id>', methods=['POST'])
@login_required
def cancel_chat_task_endpoint(task_id):
    """Cancel a running chat task - now uses database for Cloud Run compatibility"""
    task = get_chat_task(task_id)
    if not task:
        return jsonify({"error": "Task not found"}), 404
    
    if task['status'] != 'running':
        return jsonify({"error": "Task is not running"}), 400
    
    add_chat_task_step(task_id, '⛔ 処理がキャンセルされました')
    update_chat_task(task_id, cancelled=True, status='cancelled')
    
    return jsonify({"success": True, "message": "Task cancelled"})

@app.route('/api/config', methods=['GET'])
def get_config():
    """Get configuration status"""
    return jsonify({
        "project_id": PROJECT_ID,
        "dataset": DEFAULT_DATASET,
        "openai_configured": bool(OPENAI_API_KEY),
        "gcp_configured": bool(GCP_SA_JSON)
    })

@app.route('/projects')
@login_required
def projects():
    """Render projects page"""
    return render_template('projects.html')

@app.route('/settings')
@login_required
def settings():
    """Render settings page"""
    # Check if user has any projects
    conn = get_db_connection()
    cur = conn.cursor(cursor_factory=RealDictCursor)
    cur.execute('SELECT id FROM projects WHERE user_id = %s LIMIT 1', (current_user.id,))
    has_project = cur.fetchone() is not None
    cur.close()
    conn.close()
    
    if not has_project:
        flash('まずプロジェクトを作成してください。', 'info')
        return redirect(url_for('projects'))
    
    return render_template('settings.html')

@app.route('/api/settings', methods=['GET'])
@login_required
def get_settings():
    """Get active project settings"""
    try:
        conn = get_db_connection()
        cur = conn.cursor(cursor_factory=RealDictCursor)
        cur.execute('''
            SELECT id, name, description, bigquery_project_id, bigquery_dataset_id,
                   openai_api_key IS NOT NULL as has_api_key,
                   gemini_api_key IS NOT NULL as has_gemini_key,
                   service_account_json,
                   use_env_json,
                   (service_account_json IS NOT NULL OR use_env_json = true) as has_service_account,
                   ai_provider,
                   max_scan_gb,
                   enable_review,
                   gemini_model
            FROM projects
            WHERE user_id = %s AND is_active = true
            ORDER BY updated_at DESC
            LIMIT 1
        ''', (current_user.id,))
        project = cur.fetchone()
        print(f"DEBUG get_settings: active project={project}")
        cur.close()
        conn.close()
        
        if not project:
            return jsonify({
                "success": False,
                "error": "アクティブなプロジェクトが選択されていません。プロジェクトを作成・選択してください。",
                "no_project": True
            }), 404
        
        # Get project_id: prefer explicit setting, then extract from service account JSON
        bigquery_project_id = project['bigquery_project_id']
        sa_extracted_project_id = None
        if project['service_account_json']:
            sa_extracted_project_id = extract_project_id_from_sa_json(project['service_account_json'])
        
        # Use explicit setting if available, otherwise use extracted from SA JSON
        effective_project_id = bigquery_project_id or sa_extracted_project_id or ''
        
        return jsonify({
            "success": True,
            "project_id": effective_project_id,
            "project_id_source": "explicit" if bigquery_project_id else ("service_account" if sa_extracted_project_id else "none"),
            "sa_extracted_project_id": sa_extracted_project_id,
            "dataset": project['bigquery_dataset_id'] or '',
            "has_api_key": project['has_api_key'],
            "has_gemini_key": project['has_gemini_key'],
            "has_service_account": project['has_service_account'],
            "project_name": project['name'],
            "project_description": project['description'] or '',
            "ai_provider": project['ai_provider'] or 'openai',
            "openai_model": OPENAI_MODEL,
            "gemini_model": project.get('gemini_model') or "gemini-2.5-pro",
            "location": LOCATION,
            "max_scan_gb": float(project['max_scan_gb']) if project.get('max_scan_gb') else DEFAULT_MAX_SCAN_GB,
            "enable_review": bool(project.get('enable_review'))
        })
    except Exception as e:
        import traceback
        return jsonify({
            "error": str(e),
            "traceback": traceback.format_exc()
        }), 500

@app.route('/api/settings', methods=['POST'])
@login_required
def save_settings():
    """Save settings to active project or create first project"""
    try:
        conn = get_db_connection()
        cur = conn.cursor(cursor_factory=RealDictCursor)
        
        # Check if user has any projects
        cur.execute('''
            SELECT id FROM projects
            WHERE user_id = %s AND is_active = true
            ORDER BY updated_at DESC
            LIMIT 1
        ''', (current_user.id,))
        project = cur.fetchone()
        print(f"DEBUG save_settings: active project={project}")
        
        # Track if this is a new project creation
        is_new_project = False
        
        # If no active project exists, create a new one (first-time setup)
        if not project:
            project_name = request.form.get('project_name', '新しいプロジェクト')
            project_description = request.form.get('project_description', '')
            ai_provider = request.form.get('ai_provider', 'openai')
            openai_key = request.form.get('openai_key', '')
            gemini_key = request.form.get('gemini_key', '')
            bigquery_project_id = request.form.get('project_id', '')
            bigquery_dataset = request.form.get('default_dataset', '')
            
            # Validate scan limit (GB) if provided
            max_scan_gb_value = None
            if request.form.get('max_scan_gb'):
                try:
                    max_scan_gb_value = float(request.form.get('max_scan_gb'))
                    if max_scan_gb_value <= 0 or max_scan_gb_value > 100000:
                        raise ValueError("out of range")
                except ValueError:
                    cur.close()
                    conn.close()
                    return jsonify({
                        "success": False,
                        "error": "スキャン上限（GB）は 0 より大きい数値で入力してください"
                    }), 400
            
            enable_review_value = request.form.get('enable_review', 'false') == 'true'
            gemini_model_value = (request.form.get('gemini_model') or '').strip() or 'gemini-2.5-pro'
            if not gemini_model_value.startswith('gemini'):
                cur.close()
                conn.close()
                return jsonify({
                    "success": False,
                    "error": "無効なGeminiモデル名です"
                }), 400
            
            # Create new project with initial settings
            cur.execute('''
                INSERT INTO projects (
                    user_id, name, description, ai_provider, openai_api_key, gemini_api_key,
                    bigquery_project_id, bigquery_dataset_id, max_scan_gb, enable_review, gemini_model,
                    is_active, created_at, updated_at
                )
                VALUES (%s, %s, %s, %s, %s, %s, %s, %s, COALESCE(%s, 10), %s, %s, true, CURRENT_TIMESTAMP, CURRENT_TIMESTAMP)
                RETURNING id, name
            ''', (current_user.id, project_name, project_description, ai_provider,
                  openai_key or None, gemini_key or None, bigquery_project_id or None, bigquery_dataset or None,
                  max_scan_gb_value, enable_review_value, gemini_model_value))
            new_project = cur.fetchone()
            project_id = new_project['id']
            is_new_project = True
            
            # Handle GCP JSON authentication
            json_source = request.form.get('json_source', 'upload')
            if json_source == 'adc' or json_source == 'env':
                # Use Application Default Credentials (ADC)
                cur.execute('''
                    UPDATE projects 
                    SET use_env_json = true, service_account_json = NULL, original_json_filename = NULL, updated_at = CURRENT_TIMESTAMP
                    WHERE id = %s
                ''', (project_id,))
            elif 'gcp_json' in request.files:
                json_file = request.files['gcp_json']
                if json_file.filename:
                    original_filename = json_file.filename
                    json_path = os.path.join(os.getcwd(), f'gcp_credentials_project_{project_id}.json')
                    json_file.save(json_path)
                    cur.execute('''
                        UPDATE projects 
                        SET service_account_json = %s, original_json_filename = %s, use_env_json = false, updated_at = CURRENT_TIMESTAMP
                        WHERE id = %s
                    ''', (json_path, original_filename, project_id))
            
            conn.commit()
            cur.close()
            conn.close()
            
            return jsonify({
                "success": True,
                "message": "プロジェクトを作成しました",
                "project_name": new_project['name'],
                "project_id": project_id,
                "is_new": True
            })
        
        # Existing project - build update fields
        project_id = project['id']
        update_fields = []
        params = []
        
        # Handle project name
        if request.form.get('project_name'):
            update_fields.append('name = %s')
            params.append(request.form.get('project_name'))
        
        # Handle project description
        if 'project_description' in request.form:
            update_fields.append('description = %s')
            params.append(request.form.get('project_description'))
        
        # Handle AI provider
        if request.form.get('ai_provider'):
            update_fields.append('ai_provider = %s')
            params.append(request.form.get('ai_provider'))
        
        # Handle Gemini model selection
        if request.form.get('gemini_model'):
            gemini_model_value = request.form.get('gemini_model').strip()
            if not gemini_model_value.startswith('gemini'):
                cur.close()
                conn.close()
                return jsonify({
                    "success": False,
                    "error": "無効なGeminiモデル名です"
                }), 400
            update_fields.append('gemini_model = %s')
            params.append(gemini_model_value)
        
        # Handle OpenAI API key
        if request.form.get('openai_key'):
            update_fields.append('openai_api_key = %s')
            params.append(request.form.get('openai_key'))
        
        # Handle Gemini API key
        if request.form.get('gemini_key'):
            update_fields.append('gemini_api_key = %s')
            params.append(request.form.get('gemini_key'))
        
        # Handle BigQuery project ID
        if request.form.get('project_id'):
            update_fields.append('bigquery_project_id = %s')
            params.append(request.form.get('project_id'))
        
        # Handle BigQuery dataset
        if request.form.get('default_dataset'):
            update_fields.append('bigquery_dataset_id = %s')
            params.append(request.form.get('default_dataset'))
        
        # Handle BigQuery scan limit (GB)
        if request.form.get('max_scan_gb'):
            try:
                max_scan_gb_value = float(request.form.get('max_scan_gb'))
                if max_scan_gb_value <= 0 or max_scan_gb_value > 100000:
                    raise ValueError("out of range")
                update_fields.append('max_scan_gb = %s')
                params.append(max_scan_gb_value)
            except ValueError:
                cur.close()
                conn.close()
                return jsonify({
                    "success": False,
                    "error": "スキャン上限（GB）は 0 より大きい数値で入力してください"
                }), 400
        
        # Handle AI review toggle
        if 'enable_review' in request.form:
            update_fields.append('enable_review = %s')
            params.append(request.form.get('enable_review') == 'true')
        
        # Handle GCP JSON authentication
        json_source = request.form.get('json_source', 'upload')
        if json_source == 'adc' or json_source == 'env':
            # Use Application Default Credentials (ADC)
            update_fields.append('use_env_json = %s')
            params.append(True)
            update_fields.append('service_account_json = %s')
            params.append(None)
            update_fields.append('original_json_filename = %s')
            params.append(None)
        elif json_source == 'upload':
            # User selected to use uploaded JSON file (not ADC)
            # Always set use_env_json to False when upload is selected
            update_fields.append('use_env_json = %s')
            params.append(False)
            
            # If a new file was uploaded, save it
            if 'gcp_json' in request.files:
                json_file = request.files['gcp_json']
                if json_file.filename:
                    # Save the JSON file with project-specific name
                    original_filename = json_file.filename
                    json_path = os.path.join(os.getcwd(), f'gcp_credentials_project_{project_id}.json')
                    json_file.save(json_path)
                    update_fields.append('service_account_json = %s')
                    params.append(json_path)
                    update_fields.append('original_json_filename = %s')
                    params.append(original_filename)
        
        if not update_fields:
            cur.close()
            conn.close()
            return jsonify({
                "success": False,
                "error": "更新する設定がありません"
            }), 400
        
        # Update existing project
        update_fields.append('updated_at = CURRENT_TIMESTAMP')
        params.extend([project_id, current_user.id])
        
        query = f'''
            UPDATE projects
            SET {', '.join(update_fields)}
            WHERE id = %s AND user_id = %s
        '''
        
        cur.execute(query, params)
        conn.commit()
        cur.close()
        conn.close()
        
        return jsonify({
            "success": True,
            "message": "設定を保存しました"
        })
    
    except Exception as e:
        import traceback
        return jsonify({
            "error": str(e),
            "traceback": traceback.format_exc()
        }), 500

@app.route('/api/get-api-key')
@login_required
def get_api_key():
    """Get API key for active project (OpenAI or Gemini)"""
    try:
        provider = request.args.get('provider', 'openai')
        
        conn = get_db_connection()
        cur = conn.cursor(cursor_factory=RealDictCursor)
        
        if provider == 'gemini':
            cur.execute('''
                SELECT gemini_api_key FROM projects
                WHERE user_id = %s AND is_active = true
                LIMIT 1
            ''', (current_user.id,))
            project = cur.fetchone()
            cur.close()
            conn.close()
            
            if not project:
                return jsonify({"api_key": ""})
            
            return jsonify({
                "api_key": project['gemini_api_key'] or ""
            })
        else:  # default to openai
            cur.execute('''
                SELECT openai_api_key FROM projects
                WHERE user_id = %s AND is_active = true
                LIMIT 1
            ''', (current_user.id,))
            project = cur.fetchone()
            cur.close()
            conn.close()
            
            if not project:
                return jsonify({"api_key": ""})
            
            return jsonify({
                "api_key": project['openai_api_key'] or ""
            })
    except Exception as e:
        return jsonify({"api_key": ""}), 500

@app.route('/api/gemini-models', methods=['GET'])
@login_required
def list_gemini_models():
    """List available Gemini models (generateContent-capable) using the active project's API key"""
    try:
        conn = get_db_connection()
        cur = conn.cursor(cursor_factory=RealDictCursor)
        cur.execute('''
            SELECT gemini_api_key FROM projects
            WHERE user_id = %s AND is_active = true
            LIMIT 1
        ''', (current_user.id,))
        project = cur.fetchone()
        cur.close()
        conn.close()
        
        if not project or not project['gemini_api_key']:
            return jsonify({
                "success": False,
                "error": "Gemini APIキーが設定されていません。先にAPIキーを保存してください。"
            }), 400
        
        client = genai_new.Client(api_key=project['gemini_api_key'])
        models = []
        for m in client.models.list():
            name = (m.name or '').replace('models/', '')
            actions = getattr(m, 'supported_actions', None) or []
            if not name.startswith('gemini'):
                continue
            if actions and 'generateContent' not in actions:
                continue
            models.append({
                "name": name,
                "display_name": getattr(m, 'display_name', '') or name,
                "description": (getattr(m, 'description', '') or '')[:200],
                "input_token_limit": getattr(m, 'input_token_limit', None),
                "output_token_limit": getattr(m, 'output_token_limit', None)
            })
        
        # Newest-first ordering (name descending puts higher versions first)
        models.sort(key=lambda x: x['name'], reverse=True)
        
        return jsonify({"success": True, "models": models})
    except Exception as e:
        return jsonify({
            "success": False,
            "error": f"モデル一覧の取得に失敗しました: {str(e)}"
        }), 500

@app.route('/api/check-env-json')
@login_required
def check_env_json():
    """Check if Application Default Credentials (ADC) is available"""
    try:
        # Check if running in Cloud Run environment
        # Cloud Run sets K_SERVICE environment variable
        is_cloud_run = os.environ.get('K_SERVICE') is not None
        
        if is_cloud_run:
            # In Cloud Run, ADC is always available via metadata server
            # The service account attached to Cloud Run will be used
            return jsonify({
                "available": True,
                "environment": "cloud_run"
            })
        
        # For non-Cloud Run environments, try to get default credentials
        import google.auth
        from google.auth.exceptions import DefaultCredentialsError
        
        try:
            credentials, project = google.auth.default(
                scopes=["https://www.googleapis.com/auth/bigquery.readonly"]
            )
            available = True
        except DefaultCredentialsError:
            available = False
        
        return jsonify({
            "available": available,
            "environment": "local"
        })
    except Exception as e:
        print(f"Error checking ADC availability: {e}")
        return jsonify({"available": False}), 500

@app.route('/api/get-json-file-info')
@login_required
def get_json_file_info():
    """Get information about uploaded JSON file for active project"""
    try:
        conn = get_db_connection()
        cur = conn.cursor(cursor_factory=RealDictCursor)
        cur.execute('''
            SELECT id, service_account_json, use_env_json, original_json_filename FROM projects
            WHERE user_id = %s AND is_active = true
            ORDER BY updated_at DESC
            LIMIT 1
        ''', (current_user.id,))
        project = cur.fetchone()
        print(f"DEBUG get-json-file-info: project={project}")
        cur.close()
        conn.close()
        
        if not project:
            return jsonify({"exists": False, "use_env": False})
        
        # Check if using ADC (Application Default Credentials)
        if project.get('use_env_json'):
            return jsonify({
                "exists": True,
                "use_env": True,
                "filename": "Cloud Run サービスアカウント（ADC）"
            })
        
        json_path = project['service_account_json']
        original_filename = project.get('original_json_filename')
        print(f"DEBUG get-json-file-info: json_path={json_path}, exists={os.path.exists(json_path) if json_path else False}")
        
        if json_path and os.path.exists(json_path):
            # Use original filename if available, otherwise fall back to the stored filename
            filename = original_filename if original_filename else os.path.basename(json_path)
            print(f"DEBUG get-json-file-info: File found - {filename}")
            return jsonify({
                "exists": True,
                "use_env": False,
                "filename": filename,
                "path": json_path
            })
        
        print(f"DEBUG get-json-file-info: No file - returning exists=False")
        return jsonify({"exists": False, "use_env": False})
    except Exception as e:
        return jsonify({"exists": False, "use_env": False}), 500

@app.route('/api/delete-json-file', methods=['POST'])
@login_required
def delete_json_file():
    """Delete uploaded JSON file for active project"""
    try:
        # Get active project
        conn = get_db_connection()
        cur = conn.cursor(cursor_factory=RealDictCursor)
        cur.execute('''
            SELECT id, service_account_json FROM projects
            WHERE user_id = %s AND is_active = true
            LIMIT 1
        ''', (current_user.id,))
        project = cur.fetchone()
        
        if not project:
            cur.close()
            conn.close()
            return jsonify({
                "success": False,
                "error": "アクティブなプロジェクトが選択されていません"
            }), 404
        
        # Delete physical file if it exists
        json_path = project['service_account_json']
        if json_path and os.path.exists(json_path):
            os.remove(json_path)
        
        # Update database to remove JSON path
        cur.execute('''
            UPDATE projects
            SET service_account_json = NULL, original_json_filename = NULL, updated_at = CURRENT_TIMESTAMP
            WHERE id = %s AND user_id = %s
        ''', (project['id'], current_user.id))
        
        conn.commit()
        cur.close()
        conn.close()
        
        return jsonify({
            "success": True,
            "message": "JSON file deleted successfully"
        })
    
    except Exception as e:
        return jsonify({
            "error": str(e)
        }), 500

@app.route('/api/test-connection', methods=['POST'])
@login_required
def test_connection():
    """Test AI provider (OpenAI/Gemini) and BigQuery connection for active project"""
    try:
        # Get active project configuration
        conn = get_db_connection()
        cur = conn.cursor(cursor_factory=RealDictCursor)
        cur.execute('''
            SELECT openai_api_key, gemini_api_key, ai_provider, bigquery_project_id, bigquery_dataset_id, service_account_json
            FROM projects
            WHERE user_id = %s AND is_active = true
            LIMIT 1
        ''', (current_user.id,))
        project = cur.fetchone()
        cur.close()
        conn.close()
        
        if not project:
            return jsonify({
                "success": False,
                "error": "アクティブなプロジェクトが選択されていません"
            }), 404
        
        openai_api_key = project['openai_api_key']
        gemini_api_key = project['gemini_api_key']
        ai_provider = project['ai_provider'] or 'openai'
        project_id = project['bigquery_project_id']
        gcp_json = project['service_account_json']
        
        # Check if at least one AI provider is configured
        if not openai_api_key and not gemini_api_key:
            return jsonify({
                "error": "OpenAI API キーまたは Gemini API キーのどちらかを設定してください"
            }), 400
        
        ai_test_result = None
        ai_provider_tested = None
        
        # Test the selected AI provider
        if ai_provider == 'openai':
            if not openai_api_key:
                return jsonify({
                    "error": "OpenAI が選択されていますが、OpenAI API キーが設定されていません。キーを入力するか、Gemini に切り替えてください。"
                }), 400
            
            try:
                client = OpenAI(api_key=openai_api_key)
                client.models.list()
                ai_test_result = "success"
                ai_provider_tested = "OpenAI"
            except Exception as e:
                error_msg = str(e)
                if "invalid_api_key" in error_msg.lower() or "incorrect api key" in error_msg.lower() or "401" in error_msg:
                    return jsonify({
                        "error": "OpenAI API キーが無効です。正しいキーを入力してください。"
                    }), 400
                elif "insufficient_quota" in error_msg.lower() or "rate_limit" in error_msg.lower():
                    return jsonify({
                        "error": "OpenAI API の利用制限に達しています。プランを確認するか、しばらく待ってから再試行してください。"
                    }), 400
                else:
                    return jsonify({
                        "error": f"OpenAI 接続エラー: {error_msg}"
                    }), 400
        
        elif ai_provider == 'gemini':
            if not gemini_api_key:
                return jsonify({
                    "error": "Gemini が選択されていますが、Gemini API キーが設定されていません。キーを入力するか、OpenAI に切り替えてください。"
                }), 400
            
            try:
                _client = genai_new.Client(api_key=gemini_api_key)
                # Simple test to verify the API key works
                response = _client.models.generate_content(model='gemini-2.0-flash', contents="Hello", config=genai_types.GenerateContentConfig(max_output_tokens=5))
                ai_test_result = "success"
                ai_provider_tested = "Gemini"
            except Exception as e:
                error_msg = str(e)
                if "api_key" in error_msg.lower() or "invalid" in error_msg.lower() or "401" in error_msg or "403" in error_msg:
                    return jsonify({
                        "error": "Gemini API キーが無効です。正しいキーを入力してください。"
                    }), 400
                elif "quota" in error_msg.lower() or "rate" in error_msg.lower():
                    return jsonify({
                        "error": "Gemini API の利用制限に達しています。プランを確認するか、しばらく待ってから再試行してください。"
                    }), 400
                else:
                    return jsonify({
                        "error": f"Gemini 接続エラー: {error_msg}"
                    }), 400
        
        # Test BigQuery
        if not project_id or not gcp_json:
            return jsonify({
                "error": "BigQuery が完全に設定されていません（プロジェクトIDとサービスアカウントが必要です）"
            }), 400
        
        if not os.path.exists(gcp_json):
            return jsonify({
                "error": f"GCP 認証ファイルが見つかりません"
            }), 400
        
        return jsonify({
            "success": True,
            "message": f"すべての接続が成功しました！{ai_provider_tested} と BigQuery が正しく設定されています。"
        })
    
    except Exception as e:
        import traceback
        return jsonify({
            "error": str(e),
            "traceback": traceback.format_exc()
        }), 500

@app.route('/api/test-service-account', methods=['POST'])
@login_required
def test_service_account():
    """Test service account connection and get details (ADC or JSON file)"""
    try:
        from google.cloud import bigquery
        from google.oauth2 import service_account
        import google.auth
        
        # Get active project configuration
        conn = get_db_connection()
        cur = conn.cursor(cursor_factory=RealDictCursor)
        cur.execute('''
            SELECT bigquery_project_id, service_account_json, use_env_json
            FROM projects
            WHERE user_id = %s AND is_active = true
            LIMIT 1
        ''', (current_user.id,))
        project = cur.fetchone()
        cur.close()
        conn.close()
        
        if not project:
            return jsonify({
                "success": False,
                "error": "アクティブなプロジェクトが選択されていません"
            }), 404
        
        use_adc = project.get('use_env_json', False)
        service_account_json = project.get('service_account_json')
        bq_project_id = project.get('bigquery_project_id')
        
        service_account_email = None
        auth_method = None
        project_id_from_credentials = None
        bq_permissions = []
        datasets_accessible = []
        
        try:
            if use_adc:
                # Test ADC (Application Default Credentials)
                auth_method = "ADC (Application Default Credentials)"
                credentials, detected_project = google.auth.default(
                    scopes=['https://www.googleapis.com/auth/bigquery']
                )
                project_id_from_credentials = detected_project
                
                # Get service account email from ADC
                service_account_email = None
                
                # Method 1: Try to get from credentials object
                if hasattr(credentials, 'service_account_email') and credentials.service_account_email:
                    service_account_email = credentials.service_account_email
                elif hasattr(credentials, '_service_account_email') and credentials._service_account_email:
                    service_account_email = credentials._service_account_email
                
                # Method 2: Try to refresh and get from credentials
                if not service_account_email or service_account_email == 'default':
                    try:
                        from google.auth.transport import requests as google_requests
                        credentials.refresh(google_requests.Request())
                        if hasattr(credentials, 'service_account_email') and credentials.service_account_email:
                            service_account_email = credentials.service_account_email
                    except:
                        pass
                
                # Method 3: Query GCP Metadata Server (works in Cloud Run/GCE)
                if not service_account_email or service_account_email == 'default':
                    try:
                        import urllib.request
                        metadata_url = "http://metadata.google.internal/computeMetadata/v1/instance/service-accounts/default/email"
                        req = urllib.request.Request(metadata_url, headers={"Metadata-Flavor": "Google"})
                        with urllib.request.urlopen(req, timeout=2) as response:
                            service_account_email = response.read().decode('utf-8').strip()
                    except Exception as meta_error:
                        print(f"Metadata server query failed: {meta_error}")
                
                # Method 4: Use IAM API to get service account info
                if not service_account_email or service_account_email == 'default':
                    try:
                        from google.auth.transport import requests as google_requests
                        from google.oauth2 import id_token
                        # Try to get identity token info
                        credentials.refresh(google_requests.Request())
                        if hasattr(credentials, 'token') and credentials.token:
                            import urllib.request
                            import json as json_module
                            token_info_url = f"https://oauth2.googleapis.com/tokeninfo?access_token={credentials.token}"
                            req = urllib.request.Request(token_info_url)
                            with urllib.request.urlopen(req, timeout=5) as response:
                                token_info = json_module.loads(response.read().decode('utf-8'))
                                if 'email' in token_info:
                                    service_account_email = token_info['email']
                    except Exception as token_error:
                        print(f"Token info query failed: {token_error}")
                
                if not service_account_email:
                    service_account_email = "(サービスアカウント名を取得できません)"
                
                # Create BigQuery client with ADC
                client = bigquery.Client(project=bq_project_id or detected_project)
                
            else:
                # Test JSON file credentials
                auth_method = "サービスアカウントJSON"
                
                if not service_account_json:
                    return jsonify({
                        "success": False,
                        "error": "サービスアカウントJSONファイルが設定されていません"
                    }), 400
                
                if not os.path.exists(service_account_json):
                    return jsonify({
                        "success": False,
                        "error": "サービスアカウントJSONファイルが見つかりません"
                    }), 400
                
                # Load credentials from JSON file
                credentials = service_account.Credentials.from_service_account_file(
                    service_account_json,
                    scopes=['https://www.googleapis.com/auth/bigquery']
                )
                service_account_email = credentials.service_account_email
                project_id_from_credentials = credentials.project_id
                
                # Create BigQuery client with JSON credentials
                client = bigquery.Client(
                    project=bq_project_id or project_id_from_credentials,
                    credentials=credentials
                )
            
            # Test BigQuery access by listing datasets
            try:
                datasets = list(client.list_datasets(max_results=5))
                for ds in datasets:
                    datasets_accessible.append(ds.dataset_id)
                bq_permissions.append("bigquery.datasets.get")
                bq_permissions.append("bigquery.datasets.list")
            except Exception as ds_error:
                error_str = str(ds_error)
                if "403" in error_str or "Permission" in error_str:
                    return jsonify({
                        "success": False,
                        "error": f"BigQueryへのアクセス権限がありません。サービスアカウントに「BigQuery データ閲覧者」ロールを付与してください。",
                        "auth_method": auth_method,
                        "service_account": service_account_email,
                        "details": error_str
                    }), 400
                else:
                    raise ds_error
            
            # Test query execution permission
            try:
                test_query = "SELECT 1 as test"
                query_job = client.query(test_query)
                list(query_job.result())
                bq_permissions.append("bigquery.jobs.create")
            except Exception as query_error:
                error_str = str(query_error)
                if "403" in error_str or "Permission" in error_str:
                    bq_permissions.append("bigquery.jobs.create (制限あり)")
                else:
                    pass  # Non-permission error, ignore
            
            return jsonify({
                "success": True,
                "auth_method": auth_method,
                "service_account": service_account_email,
                "project_id_from_credentials": project_id_from_credentials,
                "configured_project_id": bq_project_id,
                "permissions": bq_permissions,
                "datasets_accessible": datasets_accessible[:5],  # Limit to 5
                "message": f"サービスアカウント接続テスト成功！"
            })
            
        except Exception as auth_error:
            error_str = str(auth_error)
            return jsonify({
                "success": False,
                "error": f"認証エラー: {error_str}",
                "auth_method": auth_method,
                "service_account": service_account_email
            }), 400
    
    except Exception as e:
        import traceback
        return jsonify({
            "success": False,
            "error": str(e),
            "traceback": traceback.format_exc()
        }), 500

@app.route('/api/validate-api-key', methods=['POST'])
@login_required
def validate_api_key():
    """Validate an API key before saving"""
    try:
        data = request.json
        provider = data.get('provider', '')
        api_key = data.get('api_key', '')
        
        if not api_key:
            return jsonify({
                "valid": False,
                "error": "API キーが入力されていません"
            }), 400
        
        if provider == 'openai':
            try:
                client = OpenAI(api_key=api_key)
                client.models.list()
                return jsonify({
                    "valid": True,
                    "message": "OpenAI API キーは有効です"
                })
            except Exception as e:
                error_msg = str(e)
                if "invalid_api_key" in error_msg.lower() or "incorrect api key" in error_msg.lower() or "401" in error_msg:
                    return jsonify({
                        "valid": False,
                        "error": "OpenAI API キーが無効です。正しいキーを入力してください。"
                    }), 400
                elif "insufficient_quota" in error_msg.lower() or "rate_limit" in error_msg.lower():
                    return jsonify({
                        "valid": False,
                        "error": "OpenAI API の利用制限に達しています。"
                    }), 400
                else:
                    return jsonify({
                        "valid": False,
                        "error": f"OpenAI 接続エラー: {error_msg}"
                    }), 400
        
        elif provider == 'gemini':
            try:
                _client = genai_new.Client(api_key=api_key)
                response = _client.models.generate_content(model='gemini-2.0-flash', contents="Hello", config=genai_types.GenerateContentConfig(max_output_tokens=5))
                return jsonify({
                    "valid": True,
                    "message": "Gemini API キーは有効です"
                })
            except Exception as e:
                error_msg = str(e)
                if "api_key" in error_msg.lower() or "invalid" in error_msg.lower() or "401" in error_msg or "403" in error_msg:
                    return jsonify({
                        "valid": False,
                        "error": "Gemini API キーが無効です。正しいキーを入力してください。"
                    }), 400
                elif "quota" in error_msg.lower() or "rate" in error_msg.lower():
                    return jsonify({
                        "valid": False,
                        "error": "Gemini API の利用制限に達しています。"
                    }), 400
                else:
                    return jsonify({
                        "valid": False,
                        "error": f"Gemini 接続エラー: {error_msg}"
                    }), 400
        else:
            return jsonify({
                "valid": False,
                "error": "不明なプロバイダーです"
            }), 400
    
    except Exception as e:
        import traceback
        return jsonify({
            "valid": False,
            "error": str(e)
        }), 500

# ============================================
# Project Management API
# ============================================

@app.route('/api/projects', methods=['GET'])
@login_required
def get_projects():
    """Get all projects for current user"""
    try:
        conn = get_db_connection()
        cur = conn.cursor(cursor_factory=RealDictCursor)
        cur.execute('''
            SELECT id, name, description, bigquery_project_id, bigquery_dataset_id,
                   is_active, created_at, updated_at
            FROM projects
            WHERE user_id = %s
            ORDER BY created_at DESC
        ''', (current_user.id,))
        projects = cur.fetchall()
        cur.close()
        conn.close()
        
        return jsonify({
            "success": True,
            "projects": [dict(p) for p in projects]
        })
    except Exception as e:
        import traceback
        return jsonify({
            "error": str(e),
            "traceback": traceback.format_exc()
        }), 500

@app.route('/api/projects', methods=['POST'])
@login_required
def create_project():
    """Create a new project"""
    try:
        data = request.json
        name = data.get('name', '').strip()
        description = data.get('description', '').strip()
        
        if not name:
            return jsonify({"error": "Project name is required"}), 400
        
        conn = get_db_connection()
        cur = conn.cursor(cursor_factory=RealDictCursor)
        
        # Create project
        cur.execute('''
            INSERT INTO projects (user_id, name, description)
            VALUES (%s, %s, %s)
            RETURNING id, name, description, is_active, created_at, updated_at
        ''', (current_user.id, name, description))
        
        project = cur.fetchone()
        conn.commit()
        cur.close()
        conn.close()
        
        return jsonify({
            "success": True,
            "project": dict(project),
            "message": "Project created successfully"
        })
    except Exception as e:
        import traceback
        return jsonify({
            "error": str(e),
            "traceback": traceback.format_exc()
        }), 500

@app.route('/api/projects/<int:project_id>', methods=['GET'])
@login_required
def get_project(project_id):
    """Get project details"""
    try:
        conn = get_db_connection()
        cur = conn.cursor(cursor_factory=RealDictCursor)
        cur.execute('''
            SELECT id, name, description, bigquery_project_id, bigquery_dataset_id,
                   is_active, created_at, updated_at,
                   openai_api_key IS NOT NULL as has_api_key,
                   use_env_json,
                   (service_account_json IS NOT NULL OR use_env_json = true) as has_service_account
            FROM projects
            WHERE id = %s AND user_id = %s
        ''', (project_id, current_user.id))
        
        project = cur.fetchone()
        cur.close()
        conn.close()
        
        if not project:
            return jsonify({"error": "Project not found"}), 404
        
        return jsonify({
            "success": True,
            "project": dict(project)
        })
    except Exception as e:
        import traceback
        return jsonify({
            "error": str(e),
            "traceback": traceback.format_exc()
        }), 500

@app.route('/api/projects/<int:project_id>', methods=['PUT'])
@login_required
def update_project(project_id):
    """Update project"""
    try:
        data = request.json
        
        conn = get_db_connection()
        cur = conn.cursor(cursor_factory=RealDictCursor)
        
        # Check if project exists and belongs to user
        cur.execute('SELECT id FROM projects WHERE id = %s AND user_id = %s', 
                   (project_id, current_user.id))
        if not cur.fetchone():
            cur.close()
            conn.close()
            return jsonify({"error": "Project not found"}), 404
        
        # Build update query dynamically
        update_fields = []
        params = []
        
        if 'name' in data:
            update_fields.append('name = %s')
            params.append(data['name'])
        if 'description' in data:
            update_fields.append('description = %s')
            params.append(data['description'])
        if 'openai_api_key' in data:
            update_fields.append('openai_api_key = %s')
            params.append(data['openai_api_key'])
        if 'bigquery_project_id' in data:
            update_fields.append('bigquery_project_id = %s')
            params.append(data['bigquery_project_id'])
        if 'bigquery_dataset_id' in data:
            update_fields.append('bigquery_dataset_id = %s')
            params.append(data['bigquery_dataset_id'])
        if 'service_account_json' in data:
            update_fields.append('service_account_json = %s')
            params.append(data['service_account_json'])
        
        if not update_fields:
            return jsonify({"error": "No fields to update"}), 400
        
        update_fields.append('updated_at = CURRENT_TIMESTAMP')
        params.extend([project_id, current_user.id])
        
        query = f'''
            UPDATE projects
            SET {', '.join(update_fields)}
            WHERE id = %s AND user_id = %s
            RETURNING id, name, description, bigquery_project_id, bigquery_dataset_id,
                      is_active, created_at, updated_at
        '''
        
        cur.execute(query, params)
        project = cur.fetchone()
        conn.commit()
        cur.close()
        conn.close()
        
        return jsonify({
            "success": True,
            "project": dict(project),
            "message": "Project updated successfully"
        })
    except Exception as e:
        import traceback
        return jsonify({
            "error": str(e),
            "traceback": traceback.format_exc()
        }), 500

@app.route('/api/projects/<int:project_id>', methods=['DELETE'])
@login_required
def delete_project(project_id):
    """Delete project"""
    try:
        conn = get_db_connection()
        cur = conn.cursor()
        
        # Check if project exists and belongs to user
        cur.execute('SELECT id FROM projects WHERE id = %s AND user_id = %s', 
                   (project_id, current_user.id))
        if not cur.fetchone():
            cur.close()
            conn.close()
            return jsonify({"error": "Project not found"}), 404
        
        # Delete project (cascade will delete chat history)
        cur.execute('DELETE FROM projects WHERE id = %s AND user_id = %s', 
                   (project_id, current_user.id))
        conn.commit()
        cur.close()
        conn.close()
        
        return jsonify({
            "success": True,
            "message": "Project deleted successfully"
        })
    except Exception as e:
        import traceback
        return jsonify({
            "error": str(e),
            "traceback": traceback.format_exc()
        }), 500

@app.route('/api/projects/<int:project_id>/activate', methods=['POST'])
@login_required
def activate_project(project_id):
    """Set project as active"""
    try:
        conn = get_db_connection()
        cur = conn.cursor(cursor_factory=RealDictCursor)
        
        # Check if project exists and belongs to user
        cur.execute('SELECT id FROM projects WHERE id = %s AND user_id = %s', 
                   (project_id, current_user.id))
        if not cur.fetchone():
            cur.close()
            conn.close()
            return jsonify({"error": "Project not found"}), 404
        
        # Deactivate all projects for this user
        cur.execute('UPDATE projects SET is_active = false WHERE user_id = %s', 
                   (current_user.id,))
        
        # Activate selected project
        cur.execute('''
            UPDATE projects SET is_active = true 
            WHERE id = %s AND user_id = %s
            RETURNING id, name, description, bigquery_project_id, bigquery_dataset_id,
                      is_active, created_at, updated_at
        ''', (project_id, current_user.id))
        
        project = cur.fetchone()
        conn.commit()
        cur.close()
        conn.close()
        
        return jsonify({
            "success": True,
            "project": dict(project),
            "message": f"Project '{project['name']}' is now active"
        })
    except Exception as e:
        import traceback
        return jsonify({
            "error": str(e),
            "traceback": traceback.format_exc()
        }), 500

@app.route('/api/projects/active', methods=['GET'])
@login_required
def get_active_project():
    """Get currently active project"""
    try:
        conn = get_db_connection()
        cur = conn.cursor(cursor_factory=RealDictCursor)
        cur.execute('''
            SELECT id, name, description, bigquery_project_id, bigquery_dataset_id,
                   openai_api_key, service_account_json, is_active, created_at, updated_at
            FROM projects
            WHERE user_id = %s AND is_active = true
            LIMIT 1
        ''', (current_user.id,))
        
        project = cur.fetchone()
        cur.close()
        conn.close()
        
        if not project:
            return jsonify({
                "success": True,
                "project": None
            })
        
        return jsonify({
            "success": True,
            "project": dict(project)
        })
    except Exception as e:
        import traceback
        return jsonify({
            "error": str(e),
            "traceback": traceback.format_exc()
        }), 500

# ============================================
# Chat Session Management API
# ============================================

@app.route('/api/chat-sessions', methods=['POST'])
@login_required
def create_chat_session():
    """Create a new chat session for active project"""
    try:
        # Get active project
        config = get_active_project_config(current_user.id)
        conn = get_db_connection()
        cur = conn.cursor(cursor_factory=RealDictCursor)
        
        cur.execute('''
            SELECT id FROM projects
            WHERE user_id = %s AND is_active = true
            LIMIT 1
        ''', (current_user.id,))
        project = cur.fetchone()
        
        if not project:
            cur.close()
            conn.close()
            return jsonify({"error": "No active project"}), 404
        
        # Create new session
        cur.execute('''
            INSERT INTO chat_sessions (project_id, title)
            VALUES (%s, %s)
            RETURNING id, project_id, title, created_at, updated_at
        ''', (project['id'], 'New Chat'))
        
        session = cur.fetchone()
        conn.commit()
        cur.close()
        conn.close()
        
        return jsonify({
            "success": True,
            "session": dict(session)
        })
    except Exception as e:
        import traceback
        return jsonify({
            "error": str(e),
            "traceback": traceback.format_exc()
        }), 500

@app.route('/api/chat-sessions', methods=['GET'])
@login_required
def get_chat_sessions():
    """Get all chat sessions for active project"""
    try:
        # Get active project
        conn = get_db_connection()
        cur = conn.cursor(cursor_factory=RealDictCursor)
        
        cur.execute('''
            SELECT id FROM projects
            WHERE user_id = %s AND is_active = true
            LIMIT 1
        ''', (current_user.id,))
        project = cur.fetchone()
        
        if not project:
            cur.close()
            conn.close()
            return jsonify({
                "success": True,
                "sessions": []
            })
        
        # Get sessions with message count
        cur.execute('''
            SELECT 
                s.id, 
                s.project_id, 
                s.title, 
                s.created_at, 
                s.updated_at,
                COUNT(h.id) as message_count
            FROM chat_sessions s
            LEFT JOIN chat_history h ON h.session_id = s.id
            WHERE s.project_id = %s
            GROUP BY s.id, s.project_id, s.title, s.created_at, s.updated_at
            ORDER BY s.updated_at DESC
        ''', (project['id'],))
        
        sessions = cur.fetchall()
        cur.close()
        conn.close()
        
        return jsonify({
            "success": True,
            "sessions": [dict(s) for s in sessions]
        })
    except Exception as e:
        import traceback
        return jsonify({
            "error": str(e),
            "traceback": traceback.format_exc()
        }), 500

@app.route('/api/chat-sessions/<int:session_id>/messages', methods=['GET'])
@login_required
def get_session_messages(session_id):
    """Get all messages for a specific chat session"""
    try:
        conn = get_db_connection()
        cur = conn.cursor(cursor_factory=RealDictCursor)
        
        # Verify session belongs to user's active project
        cur.execute('''
            SELECT s.id 
            FROM chat_sessions s
            JOIN projects p ON p.id = s.project_id
            WHERE s.id = %s AND p.user_id = %s
            LIMIT 1
        ''', (session_id, current_user.id))
        
        if not cur.fetchone():
            cur.close()
            conn.close()
            return jsonify({"error": "Session not found"}), 404
        
        # Get all messages for this session
        cur.execute('''
            SELECT id, user_message, ai_response, query_result, created_at, steps_count, processing_time, reasoning_process
            FROM chat_history
            WHERE session_id = %s
            ORDER BY created_at ASC
        ''', (session_id,))
        
        messages = cur.fetchall()
        cur.close()
        conn.close()
        
        return jsonify({
            "success": True,
            "messages": [dict(m) for m in messages]
        })
    except Exception as e:
        import traceback
        return jsonify({
            "error": str(e),
            "traceback": traceback.format_exc()
        }), 500

@app.route('/api/input-history', methods=['GET'])
@login_required
def get_input_history():
    """Get recent distinct user inputs across all sessions of the active project"""
    try:
        conn = get_db_connection()
        cur = conn.cursor(cursor_factory=RealDictCursor)
        cur.execute('''
            SELECT DISTINCT ON (h.user_message) h.user_message, h.created_at
            FROM chat_history h
            JOIN chat_sessions s ON s.id = h.session_id
            JOIN projects p ON p.id = s.project_id
            WHERE p.user_id = %s AND p.is_active = TRUE
              AND h.user_message IS NOT NULL AND h.user_message <> ''
            ORDER BY h.user_message, h.created_at DESC
        ''', (current_user.id,))
        rows = cur.fetchall()
        cur.close()
        conn.close()
        rows.sort(key=lambda r: r['created_at'], reverse=True)
        return jsonify({
            "success": True,
            "history": [
                {"message": r['user_message'], "created_at": r['created_at'].isoformat() if r['created_at'] else None}
                for r in rows[:50]
            ]
        })
    except Exception as e:
        return jsonify({"success": False, "error": str(e)}), 500


@app.route('/api/chat-sessions/<int:session_id>/title', methods=['PUT'])
@login_required
def update_session_title(session_id):
    """Update chat session title"""
    try:
        data = request.json
        new_title = data.get('title', '').strip()
        
        if not new_title:
            return jsonify({"error": "Title is required"}), 400
        
        conn = get_db_connection()
        cur = conn.cursor(cursor_factory=RealDictCursor)
        
        # Verify session belongs to user's active project
        cur.execute('''
            SELECT s.id 
            FROM chat_sessions s
            JOIN projects p ON p.id = s.project_id
            WHERE s.id = %s AND p.user_id = %s
            LIMIT 1
        ''', (session_id, current_user.id))
        
        if not cur.fetchone():
            cur.close()
            conn.close()
            return jsonify({"error": "Session not found"}), 404
        
        # Update title
        cur.execute('''
            UPDATE chat_sessions
            SET title = %s, updated_at = CURRENT_TIMESTAMP
            WHERE id = %s
            RETURNING id, project_id, title, created_at, updated_at
        ''', (new_title, session_id))
        
        session = cur.fetchone()
        conn.commit()
        cur.close()
        conn.close()
        
        return jsonify({
            "success": True,
            "session": dict(session)
        })
    except Exception as e:
        import traceback
        return jsonify({
            "error": str(e),
            "traceback": traceback.format_exc()
        }), 500

@app.route('/api/chat-sessions/<int:session_id>', methods=['DELETE'])
@login_required
def delete_session(session_id):
    """Delete a chat session"""
    try:
        conn = get_db_connection()
        cur = conn.cursor(cursor_factory=RealDictCursor)
        
        # Verify session belongs to user
        cur.execute('''
            SELECT s.id 
            FROM chat_sessions s
            JOIN projects p ON p.id = s.project_id
            WHERE s.id = %s AND p.user_id = %s
            LIMIT 1
        ''', (session_id, current_user.id))
        
        if not cur.fetchone():
            cur.close()
            conn.close()
            return jsonify({"error": "Session not found"}), 404
        
        # Delete session (cascade will delete related chat_history)
        cur.execute('DELETE FROM chat_sessions WHERE id = %s', (session_id,))
        
        conn.commit()
        cur.close()
        conn.close()
        
        return jsonify({
            "success": True,
            "message": "Session deleted successfully"
        })
    except Exception as e:
        import traceback
        return jsonify({
            "error": str(e),
            "traceback": traceback.format_exc()
        }), 500

# ============================================
# Account Settings API
# ============================================

@app.route('/api/account/password', methods=['PUT'])
@login_required
def change_password():
    """Change user password"""
    try:
        data = request.json
        current_password = data.get('current_password', '')
        new_password = data.get('new_password', '')
        
        if not current_password or not new_password:
            return jsonify({"error": "Current password and new password are required"}), 400
        
        if len(new_password) < 6:
            return jsonify({"error": "New password must be at least 6 characters"}), 400
        
        conn = get_db_connection()
        cur = conn.cursor(cursor_factory=RealDictCursor)
        
        # Get current password hash
        cur.execute('SELECT password_hash FROM users WHERE id = %s', (current_user.id,))
        user = cur.fetchone()
        
        if not user:
            cur.close()
            conn.close()
            return jsonify({"error": "User not found"}), 404
        
        # Verify current password
        if not bcrypt.checkpw(current_password.encode('utf-8'), user['password_hash'].encode('utf-8')):
            cur.close()
            conn.close()
            return jsonify({"error": "Current password is incorrect"}), 401
        
        # Hash new password
        new_password_hash = bcrypt.hashpw(new_password.encode('utf-8'), bcrypt.gensalt()).decode('utf-8')
        
        # Update password
        cur.execute('''
            UPDATE users
            SET password_hash = %s
            WHERE id = %s
        ''', (new_password_hash, current_user.id))
        
        conn.commit()
        cur.close()
        conn.close()
        
        return jsonify({
            "success": True,
            "message": "Password changed successfully"
        })
    except Exception as e:
        import traceback
        return jsonify({
            "error": str(e),
            "traceback": traceback.format_exc()
        }), 500

@app.route('/api/account/email', methods=['PUT'])
@login_required
def update_email():
    """Update user email"""
    try:
        data = request.json
        email = data.get('email', '').strip()
        
        conn = get_db_connection()
        cur = conn.cursor(cursor_factory=RealDictCursor)
        
        # Update email
        cur.execute('''
            UPDATE users
            SET email = %s
            WHERE id = %s
        ''', (email if email else None, current_user.id))
        
        conn.commit()
        cur.close()
        conn.close()
        
        return jsonify({
            "success": True,
            "message": "Email updated successfully"
        })
    except Exception as e:
        import traceback
        return jsonify({
            "error": str(e),
            "traceback": traceback.format_exc()
        }), 500

@app.route('/api/account', methods=['DELETE'])
@login_required
def delete_account():
    """Delete user account"""
    try:
        data = request.json
        password = data.get('password', '')
        
        if not password:
            return jsonify({"error": "Password is required"}), 400
        
        conn = get_db_connection()
        cur = conn.cursor(cursor_factory=RealDictCursor)
        
        # Get current password hash
        cur.execute('SELECT password_hash FROM users WHERE id = %s', (current_user.id,))
        user = cur.fetchone()
        
        if not user:
            cur.close()
            conn.close()
            return jsonify({"error": "User not found"}), 404
        
        # Verify password
        if not bcrypt.checkpw(password.encode('utf-8'), user['password_hash'].encode('utf-8')):
            cur.close()
            conn.close()
            return jsonify({"error": "Password is incorrect"}), 401
        
        # Delete user (cascades to projects, chat_sessions, chat_history)
        cur.execute('DELETE FROM users WHERE id = %s', (current_user.id,))
        
        conn.commit()
        cur.close()
        conn.close()
        
        # Logout user
        logout_user()
        
        return jsonify({
            "success": True,
            "message": "Account deleted successfully"
        })
    except Exception as e:
        import traceback
        return jsonify({
            "error": str(e),
            "traceback": traceback.format_exc()
        }), 500

# ============================================
# Memory Management API
# ============================================
@app.route('/api/projects/<int:project_id>/memories', methods=['GET'])
@login_required
def get_project_memories(project_id):
    """Get all memories for a project"""
    try:
        conn = get_db_connection()
        cur = conn.cursor(cursor_factory=RealDictCursor)
        
        # Verify project ownership
        cur.execute('SELECT id FROM projects WHERE id = %s AND user_id = %s', (project_id, current_user.id))
        if not cur.fetchone():
            cur.close()
            conn.close()
            return jsonify({"error": "Project not found or access denied"}), 404
        
        # Get all memories for this project
        cur.execute('''
            SELECT id, memory_key, memory_value, category, table_name, created_at, updated_at
            FROM project_memories
            WHERE project_id = %s AND user_id = %s
            ORDER BY updated_at DESC
        ''', (project_id, current_user.id))
        
        memories = cur.fetchall()
        cur.close()
        conn.close()
        
        return jsonify({
            "success": True,
            "memories": memories
        })
    except Exception as e:
        import traceback
        return jsonify({
            "error": str(e),
            "traceback": traceback.format_exc()
        }), 500

@app.route('/api/projects/<int:project_id>/memories', methods=['POST'])
@login_required
def create_memory(project_id):
    """Create a new memory for a project"""
    try:
        data = request.json
        memory_key = data.get('memory_key', '').strip()
        memory_value = data.get('memory_value', '').strip()
        category = (data.get('category') or 'general').strip()
        table_name = (data.get('table_name') or '').strip() or None
        
        if category not in ('general', 'table_info', 'business_rule', 'document_context'):
            return jsonify({"error": "Invalid category. Must be one of: general, table_info, business_rule, document_context"}), 400
        if category == 'table_info' and not table_name:
            return jsonify({"error": "テーブル情報にはテーブル名を指定してください"}), 400
        
        if not memory_key or not memory_value:
            return jsonify({"error": "Memory key and value are required"}), 400
        
        conn = get_db_connection()
        cur = conn.cursor(cursor_factory=RealDictCursor)
        
        # Verify project ownership
        cur.execute('SELECT id FROM projects WHERE id = %s AND user_id = %s', (project_id, current_user.id))
        if not cur.fetchone():
            cur.close()
            conn.close()
            return jsonify({"error": "Project not found or access denied"}), 404
        
        # Insert memory
        cur.execute('''
            INSERT INTO project_memories (project_id, user_id, memory_key, memory_value, category, table_name, updated_at)
            VALUES (%s, %s, %s, %s, %s, %s, CURRENT_TIMESTAMP)
            RETURNING id, memory_key, memory_value, category, table_name, created_at, updated_at
        ''', (project_id, current_user.id, memory_key, memory_value, category, table_name))
        
        memory = cur.fetchone()
        conn.commit()
        cur.close()
        conn.close()
        
        return jsonify({
            "success": True,
            "message": "Memory created successfully",
            "memory": memory
        })
    except Exception as e:
        import traceback
        return jsonify({
            "error": str(e),
            "traceback": traceback.format_exc()
        }), 500

@app.route('/api/memories/<int:memory_id>', methods=['PUT'])
@login_required
def update_memory(memory_id):
    """Update an existing memory"""
    try:
        data = request.json
        memory_key = data.get('memory_key', '').strip()
        memory_value = data.get('memory_value', '').strip()
        category = (data.get('category') or 'general').strip()
        table_name = (data.get('table_name') or '').strip() or None
        
        if category not in ('general', 'table_info', 'business_rule', 'document_context'):
            return jsonify({"error": "Invalid category. Must be one of: general, table_info, business_rule, document_context"}), 400
        if category == 'table_info' and not table_name:
            return jsonify({"error": "テーブル情報にはテーブル名を指定してください"}), 400
        
        if not memory_key or not memory_value:
            return jsonify({"error": "Memory key and value are required"}), 400
        
        conn = get_db_connection()
        cur = conn.cursor(cursor_factory=RealDictCursor)
        
        # Verify memory ownership
        cur.execute('SELECT id FROM project_memories WHERE id = %s AND user_id = %s', (memory_id, current_user.id))
        if not cur.fetchone():
            cur.close()
            conn.close()
            return jsonify({"error": "Memory not found or access denied"}), 404
        
        # Update memory
        cur.execute('''
            UPDATE project_memories
            SET memory_key = %s, memory_value = %s, category = %s, table_name = %s, updated_at = CURRENT_TIMESTAMP
            WHERE id = %s
            RETURNING id, memory_key, memory_value, category, table_name, created_at, updated_at
        ''', (memory_key, memory_value, category, table_name, memory_id))
        
        memory = cur.fetchone()
        conn.commit()
        cur.close()
        conn.close()
        
        return jsonify({
            "success": True,
            "message": "Memory updated successfully",
            "memory": memory
        })
    except Exception as e:
        import traceback
        return jsonify({
            "error": str(e),
            "traceback": traceback.format_exc()
        }), 500

@app.route('/api/memories/<int:memory_id>', methods=['DELETE'])
@login_required
def delete_memory(memory_id):
    """Delete a memory"""
    try:
        conn = get_db_connection()
        cur = conn.cursor(cursor_factory=RealDictCursor)
        
        # Verify memory ownership
        cur.execute('SELECT id FROM project_memories WHERE id = %s AND user_id = %s', (memory_id, current_user.id))
        if not cur.fetchone():
            cur.close()
            conn.close()
            return jsonify({"error": "Memory not found or access denied"}), 404
        
        # Delete memory
        cur.execute('DELETE FROM project_memories WHERE id = %s', (memory_id,))
        conn.commit()
        cur.close()
        conn.close()
        
        return jsonify({
            "success": True,
            "message": "Memory deleted successfully"
        })
    except Exception as e:
        import traceback
        return jsonify({
            "error": str(e),
            "traceback": traceback.format_exc()
        }), 500

# ============================================================
# Project documents (proposal / minutes) for context-aware insights
# ============================================================

MAX_DOCUMENT_SIZE = 20 * 1024 * 1024  # 20MB
MAX_EXTRACT_CHARS = 30000  # cap text sent to the AI for extraction

def extract_text_from_document(file_storage, filename):
    """Extract plain text from an uploaded PPTX/PDF/TXT file. Returns (text, file_type)."""
    lower = filename.lower()
    if lower.endswith('.pptx'):
        from pptx import Presentation
        from io import BytesIO
        prs = Presentation(BytesIO(file_storage.read()))
        texts = []
        for i, slide in enumerate(prs.slides, 1):
            slide_texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    t = shape.text_frame.text.strip()
                    if t:
                        slide_texts.append(t)
                if getattr(shape, 'has_table', False) and shape.has_table:
                    for row in shape.table.rows:
                        cells = [c.text.strip() for c in row.cells if c.text.strip()]
                        if cells:
                            slide_texts.append(" | ".join(cells))
            if getattr(slide, 'has_notes_slide', False) and slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    slide_texts.append(f"（ノート）{notes}")
            if slide_texts:
                texts.append(f"--- スライド {i} ---\n" + "\n".join(slide_texts))
        return "\n\n".join(texts), 'pptx'
    elif lower.endswith('.pdf'):
        from pypdf import PdfReader
        from io import BytesIO
        reader = PdfReader(BytesIO(file_storage.read()))
        texts = []
        for i, page in enumerate(reader.pages, 1):
            t = (page.extract_text() or "").strip()
            if t:
                texts.append(f"--- ページ {i} ---\n{t}")
        return "\n\n".join(texts), 'pdf'
    elif lower.endswith('.txt') or lower.endswith('.md'):
        raw = file_storage.read()
        try:
            text = raw.decode('utf-8')
        except UnicodeDecodeError:
            text = raw.decode('cp932', errors='replace')
        return text.strip(), 'text'
    else:
        raise ValueError("対応していないファイル形式です。PPTX / PDF / TXT / MD をアップロードしてください。")


def extract_document_context_with_ai(provider, api_key, gemini_api_key, gemini_model, filename, document_text):
    """Use the project's AI provider to extract objectives / analysis policy / KPI definitions.

    Returns a list of {"memory_key": ..., "memory_value": ...} entries.
    Raises on failure (caller returns error to user; no silent fallback).
    """
    doc_text = document_text[:MAX_EXTRACT_CHARS]
    extraction_prompt = f"""以下は、データ分析プロジェクトに関する文書（提案書または議事録）「{filename}」から抽出したテキストです。

この文書から、今後のデータ分析の前提として常に参照すべき情報を抽出してください:
1. **分析の目的** — このプロジェクト/案件で何を明らかにしたいのか、ビジネス上のゴール
2. **分析方針** — どのような切り口・アプローチ・優先順位で分析を進めるか
3. **KPI定義** — 重要指標の名称と定義（計算方法、対象期間、フィルタ条件など）
4. **その他の重要な前提** — 分析時に必ず考慮すべき制約・決定事項（あれば）

文書に該当する情報がない項目は省略してください。内容は文書に書かれていることに忠実に、簡潔にまとめてください（推測で補わないこと）。

以下のJSON形式のみで回答してください:
{{"items": [{{"title": "分析の目的", "content": "..."}}, {{"title": "分析方針", "content": "..."}}, ...]}}

## 文書テキスト
{doc_text}"""

    if provider == 'gemini':
        client = genai_new.Client(api_key=gemini_api_key)
        response = client.models.generate_content(
            model=gemini_model or 'gemini-2.5-flash',
            contents=extraction_prompt,
            config=genai_types.GenerateContentConfig(response_mime_type="application/json")
        )
        text = response.text
    else:
        client = OpenAI(api_key=api_key)
        response = client.chat.completions.create(
            model=OPENAI_MODEL,
            messages=[
                {"role": "system", "content": "You are a precise document analyst. Respond only with valid JSON."},
                {"role": "user", "content": extraction_prompt}
            ],
            response_format={"type": "json_object"}
        )
        text = response.choices[0].message.content

    parsed = json.loads(text)
    items = parsed.get("items") or []
    results = []
    for item in items:
        title = str(item.get("title", "")).strip()
        content = str(item.get("content", "")).strip()
        if title and content:
            results.append({
                "memory_key": f"【{filename}】{title}",
                "memory_value": content
            })
    return results


MAX_METRICS_PER_ANSWER = 10


def extract_analysis_metrics_with_ai(provider, api_key, gemini_api_key, gemini_model, question, answer):
    """Extract key numeric metrics (name + value) from an analysis answer via AI.

    Returns a list of {"name": str, "value": float, "unit": str}. Raises on failure
    (callers treat metric extraction as best-effort and log the error).
    """
    extraction_prompt = f"""以下はデータ分析AIの回答です。回答に含まれる主要な数値指標を抽出してください。

## ユーザーの質問
{question[:500]}

## 分析回答
{answer[:6000]}

ルール:
- 分析結果の核心となる指標のみ抽出（合計値、平均値、件数、割合など）。最大{MAX_METRICS_PER_ANSWER}件
- 指標名は内容が特定できる具体的な日本語（例:「2024年3月の総売上」「メール開封率」）
- value は数値のみ（カンマ・単位を除去した数値）。unit は単位（円、件、% など。なければ空文字）
- 例示・仮定の数値や、SQLのLIMIT値などの指標でない数値は含めない
- 指標がなければ空配列

以下のJSON形式のみで回答してください:
{{"metrics": [{{"name": "...", "value": 123.4, "unit": "..."}}]}}"""

    if provider == 'gemini':
        client = genai_new.Client(api_key=gemini_api_key)
        response = client.models.generate_content(
            model=gemini_model or 'gemini-2.5-flash',
            contents=extraction_prompt,
            config=genai_types.GenerateContentConfig(response_mime_type="application/json")
        )
        text = response.text
    else:
        client = OpenAI(api_key=api_key)
        response = client.chat.completions.create(
            model=OPENAI_MODEL,
            messages=[
                {"role": "system", "content": "You are a precise metrics extractor. Respond only with valid JSON."},
                {"role": "user", "content": extraction_prompt}
            ],
            response_format={"type": "json_object"}
        )
        text = response.choices[0].message.content

    parsed = json.loads(text)
    metrics = []
    for item in (parsed.get("metrics") or [])[:MAX_METRICS_PER_ANSWER]:
        name = str(item.get("name", "")).strip()[:255]
        unit = str(item.get("unit", "") or "").strip()[:50]
        try:
            value = float(item.get("value"))
        except (TypeError, ValueError):
            continue
        if name and value == value and value not in (float('inf'), float('-inf')):
            metrics.append({"name": name, "value": value, "unit": unit})
    return metrics


def save_analysis_metrics(provider, api_key, gemini_api_key, gemini_model,
                          project_db_id, user_id, session_id, question, answer):
    """Extract metrics from an analysis answer and persist them. Best-effort; raises on failure."""
    if not answer or not answer.strip():
        return
    metrics = extract_analysis_metrics_with_ai(
        provider, api_key, gemini_api_key, gemini_model, question, answer)
    if not metrics:
        return
    conn = get_db_connection()
    try:
        cur = conn.cursor()
        for m in metrics:
            cur.execute('''
                INSERT INTO analysis_metrics (project_id, user_id, session_id, question, metric_name, metric_value, metric_unit)
                VALUES (%s, %s, %s, %s, %s, %s, %s)
            ''', (project_db_id, user_id, session_id, (question or '')[:2000], m['name'], m['value'], m['unit'] or None))
        conn.commit()
        cur.close()
        print(f"DEBUG: Saved {len(metrics)} analysis metrics for project {project_db_id}")
    finally:
        conn.close()


def _run_document_extraction(cur, project, project_id, document_id, filename, text):
    """Run AI extraction and insert document_context memories. Returns list of created memories."""
    provider = project.get('ai_provider') or 'openai'
    extracted = extract_document_context_with_ai(
        provider,
        project.get('openai_api_key') or OPENAI_API_KEY,
        project.get('gemini_api_key'),
        project.get('gemini_model'),
        filename, text
    )
    memories = []
    for entry in extracted:
        cur.execute('''
            INSERT INTO project_memories (project_id, user_id, memory_key, memory_value, category, document_id, updated_at)
            VALUES (%s, %s, %s, %s, 'document_context', %s, CURRENT_TIMESTAMP)
            RETURNING id, memory_key, memory_value, category, table_name, document_id, created_at, updated_at
        ''', (project_id, current_user.id, entry['memory_key'], entry['memory_value'], document_id))
        memories.append(cur.fetchone())
    return memories


@app.route('/api/projects/<int:project_id>/documents', methods=['GET'])
@login_required
def list_documents(project_id):
    """List uploaded documents for a project (with their extracted memory counts)"""
    try:
        conn = get_db_connection()
        cur = conn.cursor(cursor_factory=RealDictCursor)
        cur.execute('SELECT id FROM projects WHERE id = %s AND user_id = %s', (project_id, current_user.id))
        if not cur.fetchone():
            cur.close()
            conn.close()
            return jsonify({"error": "Project not found or access denied"}), 404
        cur.execute('''
            SELECT d.id, d.filename, d.file_type, d.created_at,
                   COUNT(m.id) AS memory_count
            FROM project_documents d
            LEFT JOIN project_memories m ON m.document_id = d.id
            WHERE d.project_id = %s AND d.user_id = %s
            GROUP BY d.id
            ORDER BY d.created_at DESC
        ''', (project_id, current_user.id))
        documents = cur.fetchall()
        cur.close()
        conn.close()
        return jsonify({"success": True, "documents": documents})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route('/api/projects/<int:project_id>/documents', methods=['POST'])
@login_required
def upload_document(project_id):
    """Upload a proposal/minutes document, extract text, and derive analysis context via AI"""
    try:
        conn = get_db_connection()
        cur = conn.cursor(cursor_factory=RealDictCursor)
        cur.execute('''
            SELECT id, ai_provider, openai_api_key, gemini_api_key, gemini_model
            FROM projects WHERE id = %s AND user_id = %s
        ''', (project_id, current_user.id))
        project = cur.fetchone()
        if not project:
            cur.close()
            conn.close()
            return jsonify({"error": "Project not found or access denied"}), 404

        if 'file' not in request.files or not request.files['file'].filename:
            cur.close()
            conn.close()
            return jsonify({"error": "ファイルが選択されていません"}), 400
        file = request.files['file']
        filename = file.filename

        file.seek(0, os.SEEK_END)
        size = file.tell()
        file.seek(0)
        if size > MAX_DOCUMENT_SIZE:
            cur.close()
            conn.close()
            return jsonify({"error": "ファイルサイズが上限（20MB）を超えています"}), 400

        try:
            text, file_type = extract_text_from_document(file, filename)
        except ValueError as ve:
            cur.close()
            conn.close()
            return jsonify({"error": str(ve)}), 400

        if not text.strip():
            cur.close()
            conn.close()
            return jsonify({"error": "文書からテキストを抽出できませんでした（画像のみのファイルの可能性があります）"}), 400

        cur.execute('''
            INSERT INTO project_documents (project_id, user_id, filename, file_type, extracted_text)
            VALUES (%s, %s, %s, %s, %s)
            RETURNING id, filename, file_type, created_at
        ''', (project_id, current_user.id, filename, file_type, text))
        document = cur.fetchone()

        try:
            memories = _run_document_extraction(cur, project, project_id, document['id'], filename, text)
        except Exception as ai_err:
            conn.rollback()
            cur.close()
            conn.close()
            return jsonify({"error": f"AIによる目的・方針の抽出に失敗しました: {ai_err}"}), 502

        conn.commit()
        cur.close()
        conn.close()
        return jsonify({
            "success": True,
            "document": document,
            "memories": memories,
            "message": f"文書をアップロードし、{len(memories)}件のコンテキストを抽出しました"
        })
    except Exception as e:
        import traceback
        return jsonify({"error": str(e), "traceback": traceback.format_exc()}), 500


@app.route('/api/documents/<int:document_id>/reextract', methods=['POST'])
@login_required
def reextract_document(document_id):
    """Re-run AI extraction for an uploaded document (replaces its previous memories)"""
    try:
        conn = get_db_connection()
        cur = conn.cursor(cursor_factory=RealDictCursor)
        cur.execute('''
            SELECT d.id, d.project_id, d.filename, d.extracted_text,
                   p.ai_provider, p.openai_api_key, p.gemini_api_key, p.gemini_model
            FROM project_documents d
            JOIN projects p ON p.id = d.project_id
            WHERE d.id = %s AND d.user_id = %s
        ''', (document_id, current_user.id))
        doc = cur.fetchone()
        if not doc:
            cur.close()
            conn.close()
            return jsonify({"error": "Document not found or access denied"}), 404
        if not (doc.get('extracted_text') or '').strip():
            cur.close()
            conn.close()
            return jsonify({"error": "この文書には抽出可能なテキストがありません"}), 400

        cur.execute('DELETE FROM project_memories WHERE document_id = %s', (document_id,))
        try:
            memories = _run_document_extraction(cur, doc, doc['project_id'], document_id, doc['filename'], doc['extracted_text'])
        except Exception as ai_err:
            conn.rollback()
            cur.close()
            conn.close()
            return jsonify({"error": f"AIによる再抽出に失敗しました: {ai_err}"}), 502

        conn.commit()
        cur.close()
        conn.close()
        return jsonify({
            "success": True,
            "memories": memories,
            "message": f"{len(memories)}件のコンテキストを再抽出しました"
        })
    except Exception as e:
        import traceback
        return jsonify({"error": str(e), "traceback": traceback.format_exc()}), 500


@app.route('/api/documents/<int:document_id>', methods=['DELETE'])
@login_required
def delete_document(document_id):
    """Delete a document and its extracted memories (via ON DELETE CASCADE)"""
    try:
        conn = get_db_connection()
        cur = conn.cursor(cursor_factory=RealDictCursor)
        cur.execute('SELECT id FROM project_documents WHERE id = %s AND user_id = %s', (document_id, current_user.id))
        if not cur.fetchone():
            cur.close()
            conn.close()
            return jsonify({"error": "Document not found or access denied"}), 404
        cur.execute('DELETE FROM project_documents WHERE id = %s', (document_id,))
        conn.commit()
        cur.close()
        conn.close()
        return jsonify({"success": True, "message": "文書と抽出済みコンテキストを削除しました"})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


def _make_bq_client_for_project(project_id, project_row):
    """Create a BigQuery client for a project row (same credential resolution as schema-tree).

    Returns (client, bq_project_id, bq_dataset_id) or raises ValueError with a user-facing message.
    """
    from google.oauth2 import service_account
    import google.auth

    bq_project_id = project_row['bigquery_project_id']
    bq_dataset_id = project_row.get('bigquery_dataset_id')
    if not bq_project_id:
        raise ValueError("BigQueryプロジェクトIDが設定されていません")

    if project_row.get('use_env_json'):
        credentials, _ = google.auth.default(
            scopes=["https://www.googleapis.com/auth/bigquery.readonly"]
        )
        return bigquery.Client(credentials=credentials, project=bq_project_id), bq_project_id, bq_dataset_id

    credentials_file = f"gcp_credentials_project_{project_id}.json"
    if os.path.exists(credentials_file):
        credentials = service_account.Credentials.from_service_account_file(credentials_file)
        return bigquery.Client(credentials=credentials, project=bq_project_id), bq_project_id, bq_dataset_id
    sa_json = project_row.get('service_account_json')
    if sa_json:
        # The column may hold either raw JSON or a path to a key file
        if os.path.exists(sa_json):
            credentials = service_account.Credentials.from_service_account_file(sa_json)
            return bigquery.Client(credentials=credentials, project=bq_project_id), bq_project_id, bq_dataset_id
        try:
            creds_data = json.loads(sa_json)
        except json.JSONDecodeError:
            raise ValueError("サービスアカウントJSONの形式が不正です")
        credentials = service_account.Credentials.from_service_account_info(creds_data)
        return bigquery.Client(credentials=credentials, project=bq_project_id), bq_project_id, bq_dataset_id
    return bigquery.Client(project=bq_project_id), bq_project_id, bq_dataset_id


MAX_BASELINE_TABLES = 20
MAX_BASELINE_NUMERIC_COLS = 8
MAX_BASELINE_DATE_COLS = 3
BQ_NUMERIC_TYPES = {"INTEGER", "INT64", "FLOAT", "FLOAT64", "NUMERIC", "BIGNUMERIC"}
BQ_DATE_TYPES = {"DATE", "DATETIME", "TIMESTAMP"}


def _bq_ident(name: str) -> str:
    """Sanitize a BigQuery identifier for safe backtick quoting (defense-in-depth)."""
    return str(name).replace('`', '').replace('\\', '')


def _compute_table_baseline(bq_client, bq_project_id, dataset_id, table_id, max_bytes_billed):
    """Compute row count, numeric column summaries, and date ranges for one table."""
    bq_project_id, dataset_id, table_id = _bq_ident(bq_project_id), _bq_ident(dataset_id), _bq_ident(table_id)
    table = bq_client.get_table(f"{bq_project_id}.{dataset_id}.{table_id}")
    numeric_cols = [_bq_ident(f.name) for f in table.schema
                    if f.field_type in BQ_NUMERIC_TYPES and (f.mode or "NULLABLE") != "REPEATED"][:MAX_BASELINE_NUMERIC_COLS]
    date_cols = [_bq_ident(f.name) for f in table.schema
                 if f.field_type in BQ_DATE_TYPES and (f.mode or "NULLABLE") != "REPEATED"][:MAX_BASELINE_DATE_COLS]

    select_parts = ["COUNT(*) AS __row_count"]
    for col in numeric_cols:
        select_parts.append(f"SUM(SAFE_CAST(`{col}` AS FLOAT64)) AS `sum_{col}`")
        select_parts.append(f"MIN(SAFE_CAST(`{col}` AS FLOAT64)) AS `min_{col}`")
        select_parts.append(f"MAX(SAFE_CAST(`{col}` AS FLOAT64)) AS `max_{col}`")
    for col in date_cols:
        select_parts.append(f"MIN(`{col}`) AS `dmin_{col}`")
        select_parts.append(f"MAX(`{col}`) AS `dmax_{col}`")

    query = f"SELECT {', '.join(select_parts)} FROM `{bq_project_id}.{dataset_id}.{table_id}`"
    job_config = bigquery.QueryJobConfig(maximum_bytes_billed=max_bytes_billed)
    row = list(bq_client.query(query, job_config=job_config).result())[0]

    numeric_stats = []
    for col in numeric_cols:
        numeric_stats.append({
            "column": col,
            "sum": _serialize_bq_value(row[f"sum_{col}"]),
            "min": _serialize_bq_value(row[f"min_{col}"]),
            "max": _serialize_bq_value(row[f"max_{col}"]),
        })
    date_ranges = []
    for col in date_cols:
        date_ranges.append({
            "column": col,
            "min": _serialize_bq_value(row[f"dmin_{col}"]),
            "max": _serialize_bq_value(row[f"dmax_{col}"]),
        })
    return row["__row_count"], numeric_stats, date_ranges


@app.route('/api/projects/<int:project_id>/baseline-stats', methods=['GET'])
@login_required
def get_baseline_stats(project_id):
    """List saved table baseline stats for a project"""
    try:
        conn = get_db_connection()
        cur = conn.cursor(cursor_factory=RealDictCursor)
        cur.execute('SELECT id FROM projects WHERE id = %s AND user_id = %s', (project_id, current_user.id))
        if not cur.fetchone():
            cur.close()
            conn.close()
            return jsonify({"error": "Project not found or access denied"}), 404
        cur.execute('''
            SELECT id, dataset_id, table_name, row_count, numeric_stats, date_ranges, error, computed_at
            FROM table_baseline_stats
            WHERE project_id = %s AND user_id = %s
            ORDER BY dataset_id, table_name
        ''', (project_id, current_user.id))
        stats = cur.fetchall()
        cur.close()
        conn.close()
        return jsonify({"success": True, "stats": stats})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route('/api/projects/<int:project_id>/baseline-stats/refresh', methods=['POST'])
@login_required
def refresh_baseline_stats(project_id):
    """Compute and save baseline aggregates (row counts, numeric summaries, date ranges)
    for the tables in the project's configured dataset. Respects max_scan_gb per query."""
    try:
        conn = get_db_connection()
        cur = conn.cursor(cursor_factory=RealDictCursor)
        cur.execute('''
            SELECT id, bigquery_project_id, bigquery_dataset_id, service_account_json, use_env_json, max_scan_gb
            FROM projects WHERE id = %s AND user_id = %s
        ''', (project_id, current_user.id))
        project = cur.fetchone()
        if not project:
            cur.close()
            conn.close()
            return jsonify({"error": "Project not found or access denied"}), 404

        try:
            bq_client, bq_project_id, bq_dataset_id = _make_bq_client_for_project(project_id, project)
        except ValueError as ve:
            cur.close()
            conn.close()
            return jsonify({"error": str(ve)}), 400

        if not bq_dataset_id:
            cur.close()
            conn.close()
            return jsonify({"error": "データセットが設定されていません。プロジェクト設定でデータセットIDを設定してください。"}), 400

        max_scan_gb = float(project.get('max_scan_gb') or 10)
        max_bytes_billed = int(max_scan_gb * (1024 ** 3))

        tables = [t.table_id for t in bq_client.list_tables(bq_dataset_id)][:MAX_BASELINE_TABLES]
        if not tables:
            cur.close()
            conn.close()
            return jsonify({"error": f"データセット {bq_dataset_id} にテーブルが見つかりませんでした"}), 404

        # Drop stale stats from other datasets (e.g. after the project dataset was changed)
        cur.execute('''
            DELETE FROM table_baseline_stats
            WHERE project_id = %s AND dataset_id <> %s
        ''', (project_id, bq_dataset_id))

        results = []
        for table_id in tables:
            row_count = None
            numeric_stats = []
            date_ranges = []
            error_msg = None
            try:
                row_count, numeric_stats, date_ranges = _compute_table_baseline(
                    bq_client, bq_project_id, bq_dataset_id, table_id, max_bytes_billed)
            except Exception as tbl_err:
                error_msg = str(tbl_err)[:1000]
                print(f"WARNING: baseline stats failed for {bq_dataset_id}.{table_id}: {tbl_err}")
            cur.execute('''
                INSERT INTO table_baseline_stats
                    (project_id, user_id, dataset_id, table_name, row_count, numeric_stats, date_ranges, error, computed_at)
                VALUES (%s, %s, %s, %s, %s, %s, %s, %s, CURRENT_TIMESTAMP)
                ON CONFLICT (project_id, dataset_id, table_name)
                DO UPDATE SET user_id = EXCLUDED.user_id,
                              row_count = EXCLUDED.row_count,
                              numeric_stats = EXCLUDED.numeric_stats,
                              date_ranges = EXCLUDED.date_ranges,
                              error = EXCLUDED.error,
                              computed_at = CURRENT_TIMESTAMP
            ''', (project_id, current_user.id, bq_dataset_id, table_id, row_count,
                  json.dumps(numeric_stats, ensure_ascii=False),
                  json.dumps(date_ranges, ensure_ascii=False),
                  error_msg))
            results.append({
                "table_name": table_id,
                "row_count": row_count,
                "numeric_columns": len(numeric_stats),
                "date_columns": len(date_ranges),
                "error": error_msg
            })
        conn.commit()
        cur.close()
        conn.close()

        succeeded = sum(1 for r in results if not r["error"])
        failed = len(results) - succeeded
        return jsonify({
            "success": True,
            "dataset_id": bq_dataset_id,
            "tables_processed": len(results),
            "succeeded": succeeded,
            "failed": failed,
            "results": results
        })
    except Exception as e:
        import traceback
        print(f"ERROR refresh_baseline_stats: {traceback.format_exc()}")
        return jsonify({"error": str(e)}), 500


@app.route('/api/projects/<int:project_id>/analysis-metrics', methods=['GET'])
@login_required
def list_analysis_metrics(project_id):
    """List recently saved analysis metrics for a project"""
    try:
        conn = get_db_connection()
        cur = conn.cursor(cursor_factory=RealDictCursor)
        cur.execute('SELECT id FROM projects WHERE id = %s AND user_id = %s', (project_id, current_user.id))
        if not cur.fetchone():
            cur.close()
            conn.close()
            return jsonify({"error": "Project not found or access denied"}), 404
        cur.execute('''
            SELECT id, question, metric_name, metric_value, metric_unit, created_at
            FROM analysis_metrics
            WHERE project_id = %s AND user_id = %s
            ORDER BY created_at DESC
            LIMIT 100
        ''', (project_id, current_user.id))
        metrics = cur.fetchall()
        cur.close()
        conn.close()
        return jsonify({"success": True, "metrics": metrics})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route('/api/projects/<int:project_id>/analysis-metrics/<int:metric_id>', methods=['DELETE'])
@login_required
def delete_analysis_metric(project_id, metric_id):
    """Delete a saved analysis metric"""
    try:
        conn = get_db_connection()
        cur = conn.cursor(cursor_factory=RealDictCursor)
        cur.execute('''
            DELETE FROM analysis_metrics
            WHERE id = %s AND project_id = %s AND user_id = %s
            RETURNING id
        ''', (metric_id, project_id, current_user.id))
        deleted = cur.fetchone()
        conn.commit()
        cur.close()
        conn.close()
        if not deleted:
            return jsonify({"error": "Metric not found or access denied"}), 404
        return jsonify({"success": True})
    except Exception as e:
        return jsonify({"error": str(e)}), 500


@app.route('/api/projects/<int:project_id>/schema-tree', methods=['GET'])
@login_required
def get_schema_tree(project_id):
    """Get BigQuery schema tree for a project using BigQuery client library or ADC"""
    try:
        conn = get_db_connection()
        cur = conn.cursor(cursor_factory=RealDictCursor)
        
        cur.execute('''
            SELECT bigquery_project_id, bigquery_dataset_id, service_account_json, use_env_json
            FROM projects
            WHERE id = %s AND user_id = %s
        ''', (project_id, current_user.id))
        
        project = cur.fetchone()
        cur.close()
        conn.close()
        
        if not project:
            return jsonify({"error": "Project not found"}), 404
        
        if not project['bigquery_project_id']:
            return jsonify({"error": "BigQuery project not configured"}), 400
        
        bq_project_id = project['bigquery_project_id']
        bq_dataset_id = project['bigquery_dataset_id']
        service_account_json = project['service_account_json']
        use_adc = project.get('use_env_json', False)
        
        temp_file = None
        try:
            from google.cloud import bigquery
            from google.oauth2 import service_account
            import google.auth
            
            bq_client = None
            
            if use_adc:
                # Use Application Default Credentials (ADC) for Cloud Run
                credentials, adc_project = google.auth.default(
                    scopes=["https://www.googleapis.com/auth/bigquery.readonly"]
                )
                bq_client = bigquery.Client(credentials=credentials, project=bq_project_id)
            else:
                credentials_file = f"gcp_credentials_project_{project_id}.json"
                if os.path.exists(credentials_file):
                    credentials = service_account.Credentials.from_service_account_file(credentials_file)
                    bq_client = bigquery.Client(credentials=credentials, project=bq_project_id)
                elif service_account_json:
                    try:
                        creds_data = json.loads(service_account_json)
                        credentials = service_account.Credentials.from_service_account_info(creds_data)
                        bq_client = bigquery.Client(credentials=credentials, project=bq_project_id)
                    except json.JSONDecodeError:
                        return jsonify({"error": "Invalid service account JSON format"}), 400
                else:
                    bq_client = bigquery.Client(project=bq_project_id)
            
            dataset_tree = {}
            
            if bq_dataset_id:
                datasets_to_process = [bq_dataset_id]
            else:
                datasets = list(bq_client.list_datasets())
                datasets_to_process = [ds.dataset_id for ds in datasets[:10]]
            
            for dataset_id in datasets_to_process:
                try:
                    tables = list(bq_client.list_tables(dataset_id))
                    dataset_tree[dataset_id] = []
                    
                    for table_ref in tables[:50]:
                        try:
                            table = bq_client.get_table(f"{bq_project_id}.{dataset_id}.{table_ref.table_id}")
                            columns = []
                            for field in table.schema:
                                columns.append({
                                    "name": field.name,
                                    "type": field.field_type,
                                    "mode": field.mode or "NULLABLE",
                                    "description": field.description or ""
                                })
                            
                            dataset_tree[dataset_id].append({
                                "name": table_ref.table_id,
                                "columns": columns,
                                "num_rows": table.num_rows,
                                "size_bytes": table.num_bytes
                            })
                        except Exception as table_err:
                            dataset_tree[dataset_id].append({
                                "name": table_ref.table_id,
                                "error": str(table_err)
                            })
                except Exception as dataset_err:
                    dataset_tree[dataset_id] = [{
                        "error": str(dataset_err)
                    }]
            
            result = []
            for dataset_name, tables in dataset_tree.items():
                result.append({
                    "name": dataset_name,
                    "tables": tables
                })
            
            response_data = {
                "success": True,
                "schema_tree": {"datasets": result}
            }
            print(f"Schema tree response: {json.dumps(response_data, default=str)[:500]}")
            return jsonify(response_data)
        finally:
            if temp_file and os.path.exists(temp_file):
                os.unlink(temp_file)
        
    except Exception as e:
        import traceback
        error_msg = str(e)
        
        # Provide user-friendly error messages for common permission issues
        if "403" in error_msg or "Forbidden" in error_msg or "Access Denied" in error_msg:
            return jsonify({
                "error": "BigQuery へのアクセス権限がありません。サービスアカウントに以下の権限が付与されているか確認してください:\n• BigQuery データ閲覧者 (roles/bigquery.dataViewer)\n• BigQuery ジョブユーザー (roles/bigquery.jobUser)",
                "error_type": "permission_denied"
            }), 403
        elif "404" in error_msg or "Not found" in error_msg:
            return jsonify({
                "error": "指定されたプロジェクトまたはデータセットが見つかりません。プロジェクトIDが正しいか確認してください。",
                "error_type": "not_found"
            }), 404
        elif "Could not automatically determine credentials" in error_msg or "DefaultCredentialsError" in error_msg:
            return jsonify({
                "error": "認証情報が見つかりません。Cloud Run の場合はサービスアカウントがアタッチされているか確認してください。",
                "error_type": "credentials_not_found"
            }), 401
        
        return jsonify({
            "error": f"BigQuery エラー: {error_msg}",
            "traceback": traceback.format_exc()
        }), 500

@app.route('/api/list-datasets', methods=['GET'])
@login_required
def list_datasets():
    """List available BigQuery datasets using service account credentials or ADC"""
    try:
        conn = get_db_connection()
        cur = conn.cursor(cursor_factory=RealDictCursor)
        
        cur.execute('''
            SELECT bigquery_project_id, service_account_json, use_env_json
            FROM projects
            WHERE user_id = %s AND is_active = true
            LIMIT 1
        ''', (current_user.id,))
        
        project = cur.fetchone()
        cur.close()
        conn.close()
        
        if not project:
            return jsonify({"error": "No active project found", "datasets": []}), 404
        
        use_adc = project.get('use_env_json', False)
        service_account_json_path = project['service_account_json']
        
        # Get project ID from explicit setting or extract from service account
        bq_project_id = project['bigquery_project_id']
        if not bq_project_id and service_account_json_path and not use_adc:
            bq_project_id = extract_project_id_from_sa_json(service_account_json_path)
        
        if not bq_project_id:
            return jsonify({"error": "BigQuery project ID not configured", "datasets": []}), 400
        
        if not use_adc and not service_account_json_path:
            return jsonify({"error": "Service account not configured", "datasets": []}), 400
        
        # Use Google Cloud BigQuery client directly
        from google.cloud import bigquery
        from google.oauth2 import service_account
        import google.auth
        
        temp_file = None
        try:
            if use_adc:
                # Use Application Default Credentials (ADC) for Cloud Run
                credentials, adc_project = google.auth.default(
                    scopes=["https://www.googleapis.com/auth/bigquery.readonly"]
                )
            elif os.path.exists(service_account_json_path):
                credentials = service_account.Credentials.from_service_account_file(
                    service_account_json_path,
                    scopes=["https://www.googleapis.com/auth/bigquery.readonly"]
                )
            else:
                # Assume it's JSON content - write to temp file
                temp_file = tempfile.NamedTemporaryFile(mode='w', suffix='.json', delete=False)
                temp_file.write(service_account_json_path)
                temp_file.close()
                credentials = service_account.Credentials.from_service_account_file(
                    temp_file.name,
                    scopes=["https://www.googleapis.com/auth/bigquery.readonly"]
                )
            
            client = bigquery.Client(project=bq_project_id, credentials=credentials)
            
            datasets = []
            for dataset in client.list_datasets():
                datasets.append({
                    "id": dataset.dataset_id,
                    "full_id": f"{bq_project_id}.{dataset.dataset_id}",
                    "friendly_name": dataset.friendly_name or dataset.dataset_id
                })
            
            return jsonify({
                "success": True,
                "project_id": bq_project_id,
                "datasets": datasets
            })
        finally:
            if temp_file and hasattr(temp_file, 'name') and os.path.exists(temp_file.name):
                os.unlink(temp_file.name)
        
    except Exception as e:
        import traceback
        error_msg = str(e)
        
        # Provide user-friendly error messages for common permission issues
        if "403" in error_msg or "Forbidden" in error_msg or "Access Denied" in error_msg:
            return jsonify({
                "success": False,
                "error": "BigQuery へのアクセス権限がありません。サービスアカウントに以下の権限が付与されているか確認してください:\n• BigQuery データ閲覧者 (roles/bigquery.dataViewer)\n• BigQuery ジョブユーザー (roles/bigquery.jobUser)",
                "error_type": "permission_denied",
                "datasets": []
            }), 403
        elif "404" in error_msg or "Not found" in error_msg:
            return jsonify({
                "success": False,
                "error": "指定されたプロジェクトが見つかりません。プロジェクトIDが正しいか確認してください。",
                "error_type": "not_found",
                "datasets": []
            }), 404
        elif "Could not automatically determine credentials" in error_msg or "DefaultCredentialsError" in error_msg:
            return jsonify({
                "success": False,
                "error": "認証情報が見つかりません。Cloud Run の場合はサービスアカウントがアタッチされているか確認してください。",
                "error_type": "credentials_not_found",
                "datasets": []
            }), 401
        
        return jsonify({
            "success": False,
            "error": f"BigQuery エラー: {error_msg}",
            "datasets": []
        }), 500


if __name__ == '__main__':
    # Local development server
    # In production (Cloud Run), gunicorn is used instead
    port = int(os.getenv('PORT', 5000))
    app.run(host='0.0.0.0', port=port, debug=True)
